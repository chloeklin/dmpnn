#!/usr/bin/env python3
"""Build the weekly progress deck: re-measurement + the monomer-B split."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn
from pathlib import Path
from PIL import Image
import copy

OUT = "week_review_monomerB_split.pptx"

INK   = RGBColor(0x2B, 0x12, 0x19)
BERRY = RGBColor(0x6D, 0x2E, 0x46)
ROSE  = RGBColor(0xA2, 0x67, 0x69)
TEAL  = RGBColor(0x1C, 0x72, 0x93)
CREAM = RGBColor(0xF3, 0xED, 0xE7)
MUTED = RGBColor(0x6B, 0x6B, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE  = RGBColor(0xD8, 0xCF, 0xC8)

HEAD = "Cambria"
BODY = "Calibri"

W, H = 13.333, 7.5
M = 0.7

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    if dark:
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(H))
        bg.fill.solid()
        bg.fill.fore_color.rgb = INK
        bg.line.fill.background()
        bg.shadow.inherit = False
    return s


def txt(s, x, y, w, h, text, size=15, bold=False, color=INK, font=BODY,
        align=PP_ALIGN.LEFT, italic=False, space_after=6, line=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = text if isinstance(text, list) else [text]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if line:
            p.line_spacing = line
        segs = ln if isinstance(ln, list) else [ln]
        for seg in segs:
            content, opts = (seg if isinstance(seg, tuple) else (seg, {}))
            r = p.add_run()
            r.text = content
            f = r.font
            f.size = Pt(opts.get("size", size))
            f.bold = opts.get("bold", bold)
            f.italic = opts.get("italic", italic)
            f.color.rgb = opts.get("color", color)
            f.name = opts.get("font", font)
    return tb


def bullets(s, x, y, w, h, items, size=14, color=INK, gap=9, marker=ROSE):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.15
        segs = it if isinstance(it, list) else [it]
        r0 = p.add_run()
        r0.text = "— "
        r0.font.size = Pt(size)
        r0.font.color.rgb = marker
        r0.font.bold = True
        r0.font.name = BODY
        for seg in segs:
            content, opts = (seg if isinstance(seg, tuple) else (seg, {}))
            r = p.add_run()
            r.text = content
            r.font.size = Pt(opts.get("size", size))
            r.font.bold = opts.get("bold", False)
            r.font.italic = opts.get("italic", False)
            r.font.color.rgb = opts.get("color", color)
            r.font.name = BODY
    return tb


def card(s, x, y, w, h, fill=CREAM, radius=True, line=None):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    if radius:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    return shp


def circle(s, x, y, d, label, fill=BERRY, tcolor=WHITE, size=15):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = tcolor
    r.font.name = BODY
    return shp


def heading(s, text, sub=None, dark=False, kicker=None, sub_w=11.9):
    y = 0.45
    if kicker:
        txt(s, M, y, 11.9, 0.28, kicker.upper(), size=11, bold=True,
            color=ROSE if not dark else ROSE, font=BODY)
        y += 0.36
    txt(s, M, y, 11.9, 0.75, text, size=33, bold=True,
        color=WHITE if dark else INK, font=HEAD)
    y += 0.86
    if sub:
        # estimate wrapped lines: ~9.9 chars per inch of width at 15pt Calibri
        nlines = max(1, -(-len(sub) // int(sub_w * 9.9)))
        h = 0.29 * nlines
        txt(s, M, y, sub_w, h, sub, size=15,
            color=CREAM if dark else MUTED, font=BODY, line=1.2)
        y += h + 0.34
    return y


def plain_table(s, x, y, w, rows, col_w, row_h=0.34, head_h=0.38,
                size=12, head_size=12, head_fill=BERRY, head_color=WHITE,
                zebra=True, bold_first_col=False, left_cols=None):
    nrows, ncols = len(rows), len(rows[0])
    total_h = head_h + row_h * (nrows - 1)
    gf = s.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(total_h))
    tbl = gf.table
    # plain grid style
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        for child in list(tblPr):
            if child.tag == qn('a:tableStyleId'):
                tblPr.remove(child)
        el = tblPr.makeelement(qn('a:tableStyleId'), {})
        el.text = "{5940675A-B579-460E-94D1-54222C63F5DA}"
        tblPr.append(el)
    tbl.first_row = True
    tbl.horz_banding = False
    scale = w / sum(col_w)
    for i, cw in enumerate(col_w):
        tbl.columns[i].width = Emu(int(Inches(cw * scale)))
    tbl.rows[0].height = Inches(head_h)
    for r in range(1, nrows):
        tbl.rows[r].height = Inches(row_h)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = head_fill
            elif zebra and r % 2 == 0:
                cell.fill.fore_color.rgb = CREAM
            else:
                cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            lefties = {0} if left_cols is None else set(left_cols)
            p.alignment = PP_ALIGN.LEFT if c in lefties else PP_ALIGN.CENTER
            content, opts = (val if isinstance(val, tuple) else (val, {}))
            run = p.add_run()
            run.text = content
            f = run.font
            f.size = Pt(head_size if r == 0 else opts.get("size", size))
            f.name = BODY
            f.bold = opts.get("bold", r == 0 or (bold_first_col and c == 0))
            f.color.rgb = opts.get("color", head_color if r == 0 else INK)
    return gf


def style_chart(chart, colors, value_fmt='0.00', show_values=True,
                legend=True, cat_size=11, val_size=10):
    chart.has_title = False
    chart.font.name = BODY
    chart.font.size = Pt(cat_size)
    if legend:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(11)
        chart.legend.font.color.rgb = MUTED
    else:
        chart.has_legend = False
    for i, ser in enumerate(chart.series):
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = colors[i % len(colors)]
        ser.format.line.fill.background()
    plot = chart.plots[0]
    plot.gap_width = 60
    if show_values:
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.number_format = value_fmt
        dl.number_format_is_linked = False
        dl.font.size = Pt(val_size)
        dl.font.bold = True
        dl.font.color.rgb = INK
        dl.font.name = BODY
        try:
            dl.position = XL_LABEL_POSITION.OUTSIDE_END
        except Exception:
            pass
    ca = chart.category_axis
    ca.has_major_gridlines = False
    ca.tick_labels.font.size = Pt(cat_size)
    ca.tick_labels.font.color.rgb = INK
    ca.tick_labels.font.name = BODY
    ca.format.line.color.rgb = RULE
    va = chart.value_axis
    va.has_major_gridlines = True
    va.major_gridlines.format.line.color.rgb = RULE
    va.major_gridlines.format.line.width = Pt(0.5)
    va.tick_labels.font.size = Pt(10)
    va.tick_labels.font.color.rgb = MUTED
    va.tick_labels.font.name = BODY
    va.format.line.fill.background()
    return chart


def footer(s, text, dark=False):
    txt(s, M, H - 0.55, 11.9, 0.3, text, size=10,
        color=CREAM if dark else MUTED, italic=True, font=BODY)



def _figdir():
    """figures/ lives next to this script; fall back to the supervisor-update folder."""
    here = Path(__file__).resolve().parent
    for cand in (here / "figures",
                 here.parent / "29-07-2026 supervisor_update" / "figures",
                 Path.cwd() / "figures"):
        if cand.is_dir():
            return cand
    return here / "figures"


FIGDIR = _figdir()


def image(s, name, x, y, w, h):
    """Place figures/<name>.png scaled to fit the (x, y, w, h) box, centred."""
    path = FIGDIR / f"{name}.png"
    if not path.is_file():
        txt(s, x, y + h / 2, w, 0.4, f"[missing figure: {name}.png]", size=13,
            color=MUTED, font=BODY, align=PP_ALIGN.CENTER)
        return None
    iw, ih = Image.open(path).size
    scale = min(w / iw, h / ih)
    dw, dh = iw * scale, ih * scale
    return s.shapes.add_picture(str(path), Inches(x + (w - dw) / 2),
                                Inches(y + (h - dh) / 2), Inches(dw), Inches(dh))


def figure_slide(title, kicker, figure, takeaway, sub=None):
    s = slide()
    y = heading(s, title, kicker=kicker, sub=sub)
    image(s, figure, M, y, 11.9, H - y - 1.42)
    card(s, M, H - 1.32, 11.9, 0.96, CREAM)
    txt(s, M + 0.25, H - 1.16, 11.4, 0.8, takeaway, size=12.5, color=INK, font=BODY, line=1.2)
    return s

# ============================================================ 1. TITLE
s = slide(dark=True)
txt(s, M, 1.45, 11.9, 0.35, "WEEKLY REVIEW  ·  22 – 29 JULY 2026", size=12,
    bold=True, color=ROSE, font=BODY)
txt(s, M, 2.0, 11.4, 1.9,
    ["Re-measuring the diagnostic,",
     "and a second chemical axis"],
    size=42, bold=True, color=WHITE, font=HEAD, space_after=2, line=1.05)
txt(s, M, 4.05, 10.6, 0.5,
    "Why last week's results are being withdrawn, and the monomer-B split built in their place",
    size=16, color=CREAM, font=BODY)

stats = [
    ("+0.048 eV", "cost of a model-selection\nbug, present since commit 1"),
    ("0.45 → 0.98", "group-mean R² across three\nidentical runs of one fold"),
    ("9 → 682", "monomers on the new\nevaluation axis"),
]
for i, (big, small) in enumerate(stats):
    x = M + i * 4.05
    txt(s, x, 5.15, 3.8, 0.6, big, size=27, bold=True, color=ROSE, font=HEAD)
    txt(s, x, 5.78, 3.7, 0.9, small.split("\n"), size=12, color=CREAM,
        font=BODY, space_after=0, line=1.25)
txt(s, M, H - 0.75, 6, 0.3, "Chloe Lin", size=13, color=CREAM, font=BODY)
s.notes_slide.notes_text_frame.text = (
    "Headline: a reproduction check failed, which uncovered a bug in every runner and a "
    "variance problem. Results withdrawn and re-measuring. The week's constructive output "
    "is the monomer-B split plus two methodological findings."
)

# ============================================================ 2. WHAT CHANGED
s = slide()
y = heading(s, "Three things happened", kicker="summary")
items = [
    ("1", "A reproduction check failed",
     "Rerunning one stored result did not reproduce it. Tracing that uncovered a "
     "model-selection bug in all five training runners, present in each one's first commit."),
    ("2", "Training is far noisier than assumed",
     "Three identical runs of EA fold 1 gave group-mean R² of 0.450, 0.790 and 0.978. "
     "Most differences we have reported are smaller than that spread."),
    ("3", "Built the monomer-B split",
     "The A axis has only 9 monomers. The B axis has 682. Along the way: the B-monomer "
     "space turns out to be dominated by two scaffold families."),
]
for i, (n, head, body) in enumerate(items):
    yy = y + i * 1.42
    card(s, M, yy, 11.9, 1.22, CREAM)
    circle(s, M + 0.3, yy + 0.36, 0.5, n, BERRY, WHITE, 16)
    txt(s, M + 1.0, yy + 0.2, 10.6, 0.35, head, size=17, bold=True, color=BERRY, font=HEAD)
    txt(s, M + 1.0, yy + 0.62, 10.6, 0.5, body, size=13, color=INK, font=BODY, line=1.2)
txt(s, M, y + 4.42, 11.9, 0.4,
    [[("Net position:  ", {"bold": True, "color": BERRY}),
     ("the 22 July model results cannot be defended and are being re-measured. "
      "The methodological contribution is larger than it was a week ago.", {})]],
    size=14, color=INK, font=BODY)
s.notes_slide.notes_text_frame.text = "Set expectations early: this is a correction week, not a results week."

# ============================================================ 3. THE BUG
s = slide()
y = heading(s, "The bug: we tested the wrong model", kicker="finding 1")
txt(s, M, y, 7.0, 0.95,
    "Training tracked validation loss and saved the best checkpoint — then prediction "
    "ignored it and used the final model left in memory after early-stopping patience expired.",
    size=15, color=INK, font=BODY, line=1.25)
bullets(s, M, y + 1.15, 7.0, 3.1, [
    [("Present in all five runners", {"bold": True}),
     (" — hpg_hier, wdmpnn, and the stage2d path for ChemArch / GlobalArch / frac.", {})],
    [("From each runner's first commit", {"bold": True}),
     (" — confirmed by git blame, so no stored prediction predates it.", {})],
    [("The model used was, by construction, worse", {"bold": True}),
     (" than the one training had already selected — and worse by a random amount.", {})],
    [("Explains a puzzle:", {"bold": True}),
     (" validation loss did not predict test performance, because the model being "
      "tested was not the one the validation loss described.", {})],
], size=13)

card(s, 8.15, y, 4.45, 3.85, INK)
txt(s, 8.45, y + 0.32, 3.9, 0.35, "WHAT IT ACTUALLY COST", size=11, bold=True, color=ROSE, font=BODY)
txt(s, 8.45, y + 0.72, 3.9, 0.8, "erratic", size=34, bold=True, color=WHITE, font=HEAD)
txt(s, 8.45, y + 1.52, 3.9, 1.1,
    ["Measured on all 266 regenerated runs:", "mean −0.003 to −0.022 eV by model,",
     "median −0.003 eV, worst case +0.144."],
    size=11.5, color=CREAM, font=BODY, space_after=0, line=1.3)
txt(s, 8.45, y + 2.75, 3.9, 0.95,
    ["Not a systematic penalty — a random one.", "The +0.048 eV from six pilot runs was",
     "not representative."],
    size=11.5, bold=True, color=ROSE, font=BODY, space_after=0, line=1.3)
txt(s, M, y + 4.15, 11.9, 0.5,
    "Fixed in all five runners. Best-checkpoint prediction is now the default, and both "
    "predictions are stored so the cost can be measured per model family.",
    size=13, color=INK, font=BODY, line=1.2)
footer(s, "Predicting from the selected checkpoint is still correct practice — but the bug's damage was random, not one-directional.")

# ============================================================ 4. VARIANCE
s = slide()
y = heading(s, "The same run, three times", kicker="finding 2",
            sub="HPG-hier, EA, monomer-heldout. Same seed, same split, same code, same GPU model. "
                "Only the process differs.")
cd = CategoryChartData()
cd.categories = ["EA fold 0", "EA fold 1"]
cd.add_series("run 1", (0.962, 0.790))
cd.add_series("run 2", (0.982, 0.450))
cd.add_series("run 3", (0.986, 0.978))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(M), Inches(y + 0.05),
                        Inches(7.2), Inches(3.85), cd)
ch = style_chart(gf.chart, [BERRY, ROSE, TEAL], value_fmt='0.000')
ch.value_axis.maximum_scale = 1.0
ch.value_axis.minimum_scale = 0.0

card(s, 8.35, y + 0.05, 4.25, 1.55, CREAM)
txt(s, 8.6, y + 0.28, 3.8, 0.3, "GROUP-MEAN R² SD", size=11, bold=True, color=BERRY, font=BODY)
txt(s, 8.6, y + 0.66, 1.7, 0.6, "0.018", size=26, bold=True, color=INK, font=HEAD)
txt(s, 10.5, y + 0.66, 1.9, 0.6, "0.091", size=26, bold=True, color=BERRY, font=HEAD)
txt(s, 8.6, y + 1.18, 1.7, 0.28, "fold 0", size=12, color=MUTED, font=BODY)
txt(s, 10.5, y + 1.18, 1.9, 0.28, "fold 1", size=12, color=MUTED, font=BODY)

bullets(s, 8.35, y + 1.8, 4.25, 2.6, [
    "Fold 1's spread covers the entire quality range we have been reporting.",
    "Noise is fold-dependent — fold 1's SD is five times fold 0's.",
    [("The \"pathological folds\"", {"bold": True}),
     (" (EA 1, EA 6, IP 5, IP 2) are now suspected to be simply the high-variance folds.", {})],
    [("Fold 1 was recorded at 0.575", {"bold": True, "color": BERRY}),
     (" — one draw from this distribution. The unmodified model reaches 0.978 on it about "
      "one run in three, so per-fold conclusions from single runs are withdrawn.", {})],
], size=12.5, gap=7)
footer(s, "Six V100 runs. Wall times 2,237–13,464 s — runs terminate at very different points.")

# ============================================================ 5. THE FIX
s = slide()
y = heading(s, "The fix, and what we rejected", kicker="protocol")
keep = [
    ("Validation design", "unchanged", "one held-out monomer, as designed"),
    ("Model selection", "best checkpoint", "the bug fix"),
    ("Replicates", "3 seeds, averaged", "42 / 43 / 44, mean prediction reported"),
]
txt(s, M, y, 5.7, 0.3, "FROZEN PROTOCOL", size=11, bold=True, color=TEAL, font=BODY)
for i, (a, b, c) in enumerate(keep):
    yy = y + 0.42 + i * 0.92
    card(s, M, yy, 5.7, 0.8, CREAM)
    txt(s, M + 0.22, yy + 0.11, 5.3, 0.28, a.upper(), size=10.5, bold=True, color=MUTED, font=BODY)
    txt(s, M + 0.22, yy + 0.36, 3.1, 0.3, b, size=15, bold=True, color=TEAL, font=HEAD)
    txt(s, M + 2.6, yy + 0.4, 2.9, 0.28, c, size=11.5, color=INK, font=BODY)

txt(s, 7.1, y, 5.5, 0.3, "TESTED AND REJECTED", size=11, bold=True, color=BERRY, font=BODY)
rej = [
    ("Row-level validation", "Stabilised fold 1 but destabilised fold 0 (SD 0.018 → 0.034), "
                             "doubled wall time, and would mean selecting on familiar "
                             "chemistry while testing on unfamiliar."),
    ("Minimum-epoch floor", "Aimed at a mechanism the data did not support — runs that trained "
                            "longer were not better."),
]
for i, (a, b) in enumerate(rej):
    yy = y + 0.42 + i * 1.38
    card(s, 7.1, yy, 5.5, 1.24, WHITE, line=RULE)
    txt(s, 7.32, yy + 0.14, 5.05, 0.3, a, size=14, bold=True, color=BERRY, font=HEAD)
    txt(s, 7.32, yy + 0.48, 5.05, 0.72, b, size=11.5, color=INK, font=BODY, line=1.2)

card(s, M, y + 3.3, 11.9, 1.05, INK)
txt(s, M + 0.3, y + 3.55, 11.3, 0.7,
    [[("Averaging predictions over repeats improved every metric in every cell tested — ", {}),
     ("6/6 on MAE, 24/24 overall", {"bold": True, "color": ROSE}),
     (".  We no longer depend on model selection behaving; we average over it.", {})]],
    size=14, color=WHITE, font=BODY, line=1.25)
footer(s, "Every model gets identical treatment — stated wherever results appear.")

# ============================================================ 6. DIVIDER
s = slide(dark=True)
txt(s, M, 2.7, 11.9, 0.4, "PART TWO", size=13, bold=True, color=ROSE, font=BODY)
txt(s, M, 3.2, 11.6, 1.4, "Two chemical axes, two different tests",
    size=40, bold=True, color=WHITE, font=HEAD, line=1.05)
txt(s, M, 4.75, 10.4, 0.5,
    "Why the monomer-A split cannot measure what we claimed, and what the monomer-B split can",
    size=16, color=CREAM, font=BODY)

# ============================================================ 7. A SPLIT
s = slide()
y = heading(s, "How the monomer-A split works", kicker="the existing split",
            sub="The dataset is an exact factorial: 9 A monomers × 682 B monomers × 7 "
                "composition/architecture cells = 42,966 rows. Every A/B pair appears exactly 7 times.")
txt(s, M, y, 11.9, 0.3, "ONE FOLD OF THE A-HELDOUT SPLIT — 9 DONOR MONOMERS", size=11,
    bold=True, color=MUTED, font=BODY)
bw, bh, by = 1.18, 0.85, y + 0.42
labels = ["train"] * 7 + ["val", "TEST"]
for i in range(9):
    x = M + i * (bw + 0.13)
    fill = CREAM if i < 7 else (ROSE if i == 7 else BERRY)
    tc = INK if i < 7 else WHITE
    card(s, x, by, bw, bh, fill)
    txt(s, x, by + 0.16, bw, 0.28, f"A{i+1}", size=13, bold=True, color=tc,
        font=BODY, align=PP_ALIGN.CENTER)
    txt(s, x, by + 0.47, bw, 0.24, labels[i], size=10, bold=(i >= 7), color=tc,
        font=BODY, align=PP_ALIGN.CENTER)
txt(s, M, by + 1.02, 11.9, 0.3,
    "4,774 rows per monomer  ·  33,418 train / 4,774 validation / 4,774 test",
    size=12, color=MUTED, font=BODY, italic=True)

rows = [
    ["consequence", "why it follows"],
    [("Models train on 7 donors", {"bold": True}),
     "The split removes the test monomer and a second one for validation. Our \"unseen chemistry\" results are 7-example extrapolations."],
    [("Within a fold, A is constant", {"bold": True}),
     "All variation across groups comes from B, composition and architecture — all seen thousands of times. An error on the held-out A enters as a near-constant offset."],
    [("Novelty is one number per fold", {"bold": True}),
     "A single held-out monomer means no within-fold similarity variation, so no performance-versus-novelty analysis is possible."],
]
plain_table(s, M, by + 1.5, 11.9, rows, [3.3, 8.6], row_h=0.62, head_h=0.36, size=12.5,
            left_cols=[0, 1])
footer(s, "A identity carries 42% of EA variance and 50% of IP variance — across only 9 monomers.")

# ============================================================ 8. THE NULL FLOOR
s = slide()
y = heading(s, "The limitation, quantified", kicker="evaluation floor",
            sub="A null predictor that ignores the held-out monomer entirely: predict each test row "
                "from the training mean of its (partner monomer, composition, architecture) cell.")
rows = [
    ["split", "target", "null floor (median group-mean R²)", "headroom to 1.0"],
    ["A-heldout", "EA", ("0.676", {"bold": True, "color": BERRY}), ("0.324", {"bold": True, "color": BERRY})],
    ["A-heldout", "IP", "−0.034", ("1.034", {"bold": True, "color": TEAL})],
]
plain_table(s, M, y, 6.8, rows, [1.8, 1.0, 2.8, 1.5], row_h=0.44, head_h=0.56, size=13)

card(s, 7.95, y, 4.65, 2.05, INK)
txt(s, 8.22, y + 0.25, 4.3, 0.3, "ON EA FOLD 2", size=11, bold=True, color=ROSE, font=BODY)
txt(s, 8.22, y + 0.62, 2.0, 0.6, "0.961", size=28, bold=True, color=WHITE, font=HEAD)
txt(s, 10.55, y + 0.62, 2.0, 0.6, "0.922", size=28, bold=True, color=ROSE, font=HEAD)
txt(s, 8.22, y + 1.18, 2.2, 0.28, "the null", size=12, color=CREAM, font=BODY)
txt(s, 10.55, y + 1.18, 2.2, 0.28, "HPG-hier", size=12, color=CREAM, font=BODY)
txt(s, 8.22, y + 1.54, 4.3, 0.3, "A predictor with no knowledge of the",
    size=11.5, color=CREAM, font=BODY)
txt(s, 8.22, y + 1.74, 4.3, 0.3, "held-out monomer beats our model.",
    size=11.5, color=CREAM, font=BODY)

txt(s, M, y + 1.75, 6.8, 1.5,
    [[("Why EA and IP diverge.  ", {"bold": True, "color": BERRY}),
     ("The null's error is essentially the held-out donor's own offset. Whether that "
      "matters depends on its size relative to variation across groups within the fold, "
      "which comes from B.", {})]],
    size=13, color=INK, font=BODY, line=1.25)
bullets(s, M, y + 3.0, 11.9, 1.5, [
    [("EA", {"bold": True, "color": BERRY}),
     (" — B spread is large (2.93 eV), donor offsets are small, so the null reconstructs most "
      "of the signal. The metric is near-degenerate.", {})],
    [("IP", {"bold": True, "color": TEAL}),
     (" — A carries more variance (0.500) and the B spread is smaller (1.76 eV), so the missing "
      "offset dominates and the null collapses. The metric is sound.", {})],
], size=13)
footer(s, "This also undermines the earlier claim that wDMPNN wins EA chemistry, 0.965 vs 0.925.")

# ============================================================ 9. VARIANCE SHARES
s = slide()
y = heading(s, "Which axis carries the signal", kicker="why a second split",
            sub="Both monomer roles carry comparable variance. What differs by two orders of "
                "magnitude is how many monomers carry it.")
cd = CategoryChartData()
cd.categories = ["EA", "IP"]
cd.add_series("monomer A identity  (9 monomers)", (0.418, 0.500))
cd.add_series("monomer B identity  (682 monomers)", (0.457, 0.329))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(M), Inches(y),
                        Inches(6.5), Inches(3.6), cd)
ch = style_chart(gf.chart, [BERRY, TEAL], value_fmt='0.000')
ch.value_axis.maximum_scale = 0.6
ch.value_axis.minimum_scale = 0.0

rows = [
    ["", "monomer A", "monomer B"],
    ["distinct monomers", ("9", {"bold": True, "color": BERRY}), ("682", {"bold": True, "color": TEAL})],
    ["rows per monomer", "4,774", "63"],
    ["EA variance share", "0.418", "0.457"],
    ["IP variance share", "0.500", "0.329"],
    ["spread of means, EA", "1.36 eV", "2.93 eV"],
    ["spread of means, IP", "1.18 eV", "1.76 eV"],
    ["Murcko scaffolds", "9, closely related", "112"],
]
plain_table(s, 7.55, y, 5.05, rows, [2.2, 1.4, 1.45], row_h=0.36, head_h=0.36, size=11.5)
txt(s, 7.55, y + 3.05, 5.05, 0.9,
    [[("Architecture — the signal we actually study — is 0.98% of EA variance and 1.46% of IP. "
      "Monomer identity is ~90%.", {})]],
    size=12, color=INK, font=BODY, line=1.25)
footer(s, "Conclusion: EA chemistry claims belong on a B split; IP chemistry claims belong on the A split.")

# ============================================================ 10. B SPLIT BUILD
s = slide()
y = heading(s, "How the monomer-B split is built", kicker="the new split",
            sub="Nine folds over the 682 B monomers, grouped by Murcko scaffold so that "
                "chemically related monomers are held out together.")
bullets(s, M, y, 6.5, 3.4, [
    [("Group the 682 B monomers by Murcko scaffold.", {"bold": True}),
     (" Side chains are stripped; monomers sharing a ring system land in the same family. "
      "This gives 112 families.", {})],
    [("Pack whole families into nine folds of 75–76.", {"bold": True}),
     (" Families are placed largest-first into the fold with the most room, so related "
      "chemistry is held out together rather than scattered across folds.", {})],
    [("Test = fold k, validation = fold k+1, cyclically.", {"bold": True}),
     (" Validation monomers are B-disjoint from both train and test, so nothing in the "
      "stopping signal appears in the test set.", {})],
    [("33,390 / 4,788 / 4,788 rows.", {"bold": True}),
     (" Test size deliberately matched to the A split's 4,774, so fold-level metrics from "
      "the two splits sit on comparable footing.", {})],
    [("Frozen and asserted.", {"bold": True}),
     (" Split seed 42, independent of the model seed. Fold membership is written to JSON "
      "and re-checked at the start of every run, so folds cannot drift across seeds.", {})],
], size=12.5, gap=10)

txt(s, 7.35, y - 0.02, 5.25, 0.3, "DIFFICULTY, AGAINST THE A SPLIT", size=11, bold=True,
    color=MUTED, font=BODY)
rows = [
    ["split", "median NN Tanimoto", "held out per fold"],
    ["A-heldout (reference)", "0.31 – 0.47", "1 monomer"],
    [("B-heldout, clustered", {"bold": True}), ("0.48 – 0.50", {"bold": True}), ("75–76 monomers", {"bold": True})],
]
plain_table(s, 7.35, y + 0.36, 5.25, rows, [2.2, 1.75, 1.3], row_h=0.46, head_h=0.5, size=11.5)

card(s, 7.35, y + 1.85, 5.25, 1.55, CREAM)
txt(s, 7.6, y + 2.04, 4.8, 1.2,
    "Grouping by scaffold keeps close analogues out of the training set, which brings the "
    "B split's difficulty close to the A split's — near enough that the two can be read "
    "side by side, with the caveat that they are not identical.",
    size=12, color=INK, font=BODY, line=1.22)
footer(s, "All 9 A monomers appear in train, validation and test — only the B axis is held out.")

# ============================================================ 11. WHY NOT DISJOINT
s = slide()
y = heading(s, "Why a scaffold-disjoint split is impossible", kicker="constraint",
            sub="Murcko clustering of the 682 B monomers gives 112 scaffolds — but the size "
                "distribution is extreme.")
cd = CategoryChartData()
cd.categories = ["1st", "2nd", "3rd", "4th", "5th", "6th", "7th", "8th", "9th", "10th"]
cd.add_series("members", (317, 109, 17, 13, 12, 11, 10, 5, 5, 5))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(M), Inches(y),
                        Inches(6.7), Inches(3.35), cd)
ch = style_chart(gf.chart, [BERRY], value_fmt='0', legend=False)
txt(s, M, y + 3.4, 6.7, 0.3, "TEN LARGEST SCAFFOLD FAMILIES, BY MEMBER COUNT",
    size=10.5, bold=True, color=MUTED, font=BODY, align=PP_ALIGN.CENTER)

card(s, 7.75, y, 4.85, 1.5, INK)
txt(s, 8.02, y + 0.2, 4.3, 0.3, "TOP TWO FAMILIES", size=11, bold=True, color=ROSE, font=BODY)
txt(s, 8.02, y + 0.55, 4.3, 0.65, "62.5%", size=34, bold=True, color=WHITE, font=HEAD)
txt(s, 8.02, y + 1.12, 4.3, 0.3, "of all 682 B monomers", size=12, color=CREAM, font=BODY)

bullets(s, 7.75, y + 1.75, 4.85, 2.0, [
    "One family has 317 members (46%), a second has 109 (16%).",
    "64 scaffolds are singletons — but they cover only 9% of monomers.",
    [("The largest family alone exceeds any balanced fold capacity", {"bold": True}),
     (", so nine equal scaffold-disjoint folds cannot exist.", {})],
    "So the split is capacity-balanced scaffold packing — close to disjoint for the "
    "110 small families, unavoidably not for the two large ones.",
], size=12, gap=8)
footer(s, "This is a property of the benchmark, not of our split: any random split of it is mostly substituent interpolation.")

# ============================================================ 12. TWO FOLD GROUPS
s = slide()
y = heading(s, "The consequence: two kinds of fold", kicker="interpretation",
            sub="Because the two large families must be split across folds, the nine folds are "
                "not exchangeable. The grouping below is derived from the frozen split, not assumed.")
rows = [
    ["folds", "what is held out", "held-out monomers with a same-core relative in training", "what it tests"],
    [("0 – 3", {"bold": True}), "76 members of the 317-family, each",
     ("76 of 76", {"bold": True, "color": BERRY}), "new substituents on a familiar core"],
    [("5", {"bold": True}), "76 members of the 109-family", "0 of 76   (rest sit in validation)",
     "a genuinely new core — but only one"],
    [("4, 6", {"bold": True}), "29 and 26 small families", "13 of 76  ·  33 of 76", "mixed"],
    [("7, 8", {"bold": True}), "28 and 29 small families", ("0 of 75", {"bold": True, "color": TEAL}),
     "new cores — the cleanest tests"],
]
plain_table(s, M, y, 11.9, rows, [0.9, 3.4, 3.5, 4.1], row_h=0.52, head_h=0.62, size=12,
            head_size=11, left_cols=[1, 3])

card(s, M, y + 2.75, 5.8, 1.5, CREAM)
txt(s, M + 0.25, y + 2.92, 5.3, 1.0,
    [[("Paired tests run within group.", {"bold": True, "color": BERRY}),
     (" Pooling five folds that test substituent interpolation with four that test new "
      "scaffolds would mix two populations.", {})]],
    size=13, color=INK, font=BODY, line=1.25)
card(s, 6.8, y + 2.75, 5.8, 1.5, WHITE, line=RULE)
txt(s, 7.05, y + 2.92, 5.3, 1.0,
    [[("The cost: power.", {"bold": True, "color": BERRY}),
     (" Groups of five and four give minimum attainable two-sided p of 0.063 and 0.125. "
      "We accept that rather than a significant-looking pooled number.", {})]],
    size=13, color=INK, font=BODY, line=1.25)

# ============================================================ 13. INTERPRETATION
s = slide()
y = heading(s, "So we read the two splits differently", kicker="interpretation")
rows = [
    ["", "A-heldout", "B-heldout (clustered)"],
    ["valid chemistry test for", ("IP  — floor ≈ 0", {"bold": True, "color": TEAL}),
     ("EA  — floor 0.38, headroom 0.62", {"bold": True, "color": TEAL})],
    ["meaning of \"unseen chemistry\"", "extrapolation from 7 examples", "generalisation from ~530"],
    ["held-out chemistries per fold", "1", "26–29  (or 1 in folds 0–3, 5)"],
    ["folds exchangeable?", "broadly yes — 9 single donors", ("no — two structural groups", {"bold": True, "color": BERRY})],
    ["paired test", "one 9-fold test, min p = 0.0039", "within group, min p = 0.063 / 0.125"],
    ["novelty analysis", "impossible", "curve across 76 monomers per fold"],
    ["contamination handling", "none needed", "filtered and unfiltered reporting"],
]
plain_table(s, M, y, 11.9, rows, [3.5, 4.2, 4.2], row_h=0.44, head_h=0.44, size=12.5,
            left_cols=[0])
txt(s, M, y + 3.75, 11.9, 0.6,
    [[("No single headline number per split.", {"bold": True, "color": BERRY}),
     ("  EA claims come from the B split, IP claims from the A split, and every group-mean R² "
      "is reported beside its fold-specific null floor and the fraction of available headroom it closes.", {})]],
    size=14, color=INK, font=BODY, line=1.3)

# ============================================================ 14. RESULT 1
figure_slide(
    "The architecture advantage holds on both chemical axes",
    "result",
    "fig_ab_comparison",
    [[("HPG-hier beats wDMPNN on ΔR² in every bar — both targets, the A axis, and both B-fold "
       "groups. Within group it loses 4 of 4 (S) and 5 of 5 (D), the maximum either group can show. "
       , {}),
      ("The gap widened after the bug fix: wDMPNN's EA ΔR² fell from 0.58 to 0.43, so the old "
       "numbers were flattering the baseline, not our model.", {"bold": True, "color": BERRY})]],
    sub="ΔR² median across folds, mean of three seeds, error bars = per-fold seed SD. B folds are "
        "shown as their two groups, never pooled — they test different things.")

# ============================================================ 15. RESULT 2
figure_slide(
    "The octamer gains only where the chemistry is new",
    "result",
    "fig_r3_architecture",
    [[("Group D holds out whole new scaffold families; group S holds out new side chains on cores "
       "the model already knows. The octamer adds +0.019 / +0.032 ΔR² on D (5–0 folds both targets) "
       "and nothing on S (2–2). ", {}),
      ("A representation change that helps only where the chemistry is genuinely new is what a real "
       "encoding improvement looks like — pooling the nine folds would have shown p = 0.18.",
       {"bold": True, "color": BERRY})]],
    sub="B-heldout clustered, split by what each fold actually tests. This is why the fold "
        "grouping mattered.")

# ============================================================ 15b. RESULT 3
figure_slide(
    "Every model clears the floor — but the floor moves",
    "result",
    "fig_r3_chemistry_floor",
    [[("The dashed line is a predictor that ignores the held-out monomer entirely. It ranges from "
       "0.02 to 0.75 across folds, so raw R² is not comparable between them — read the gap, not "
       "the level. ", {}),
      ("All four models sit well above it on the B split, which the A split could not "
       "demonstrate for EA.", {"bold": True, "color": BERRY})]],
    sub="Per-fold group-mean R² against the fold's own B-blind null floor. Shaded folds are "
        "group S.")

# ============================================================ 16. WHATS NEXT
s = slide()
y = heading(s, "Where that leaves the programme", kicker="next")
rows = [
    ["stage", "status", "what it gave us"],
    ["R1  A-heldout, 270 runs", ("complete", {"bold": True, "color": TEAL}),
     "architecture advantage confirmed with error bars; chemistry claims repaired"],
    ["R3  B-heldout clustered, 216 runs", ("complete", {"bold": True, "color": TEAL}),
     "advantage holds on the second chemical axis, 9 of 9 folds"],
    ["R2  B-heldout random", "deferred", "clustered is the difficulty-matched split; little to add"],
    ["R4  ChemArch / GlobalArch / frac", "not run", "their numbers are still from the buggy code — see below"],
]
plain_table(s, M, y, 11.9, rows, [3.4, 1.6, 6.9], row_h=0.5, head_h=0.46, size=12,
            left_cols=[0, 2])
txt(s, M, y + 2.45, 11.9, 0.35,
    "481 of 486 cells landed; the five outstanding move no median measurably. "
    "57 of 266 R1 runs were flagged undertrained and every table is reported with and without them.",
    size=11.5, color=MUTED, font=BODY, italic=True, line=1.2)

txt(s, M, y + 3.0, 11.9, 0.3, "FOR DISCUSSION", size=11, bold=True, color=BERRY, font=BODY)
dis = [
    ("Paper scope", "The measurement contributions — null-floor calibration, the two-axis split "
                    "design, scaffold dominance in the benchmark — are now backed by a model result "
                    "that holds on both axes. Lead with the methodology, HPG-hier as the case study?"),
    ("ChemArch and the original HPG", "Still un-regenerated, so they cannot share a table with these "
                                      "numbers. Either spend 162 runs, or cut them from the comparison."),
    ("Training stability", "21% of runs stopped before epoch 10. Worth fixing properly before the "
                           "next experiment, even though 3-seed averaging absorbs it."),
]
for i, (a, b) in enumerate(dis):
    yy = y + 3.38 + i * 0.8
    txt(s, M, yy, 3.0, 0.3, a, size=13, bold=True, color=BERRY, font=HEAD)
    txt(s, M + 3.15, yy, 8.75, 0.75, b, size=11.5, color=INK, font=BODY, line=1.18)

# ============================================================ 16. CLOSING
s = slide(dark=True)
txt(s, M, 1.5, 11.9, 0.35, "WHERE THIS LEAVES US", size=12, bold=True, color=ROSE, font=BODY)
txt(s, M, 2.05, 11.4, 1.5,
    "A correction week that ended with a stronger result",
    size=36, bold=True, color=WHITE, font=HEAD, line=1.08)
items = [
    ("Found before publication, not after", "A reproduction check most people skip caught a bug "
     "that had been in every runner since its first commit."),
    ("Three findings the bug cannot touch", "The null-floor calibration, the two-axis split design, "
     "and the scaffold dominance in the benchmark are all computed from data, not from model runs."),
    ("The model claim came back stronger", "Architecture recovery is confirmed with error bars and "
     "holds on both chemical axes: ΔR² 0.80/0.81 vs wDMPNN's 0.43/0.45 on the A axis, and 9 of 9 "
     "folds on the B axis. The bug had been understating our own model."),
]
for i, (a, b) in enumerate(items):
    yy = 3.85 + i * 1.05
    circle(s, M, yy + 0.02, 0.42, str(i + 1), ROSE, INK, 14)
    txt(s, M + 0.62, yy, 11.0, 0.3, a, size=16, bold=True, color=WHITE, font=HEAD)
    txt(s, M + 0.62, yy + 0.35, 11.0, 0.5, b, size=12.5, color=CREAM, font=BODY, line=1.2)

# ============================================================ 17. BACKUP: PROVENANCE
s = slide()
y = heading(s, "Backup — how the fold grouping was derived", kicker="provenance",
            sub="Nothing here is hand-assigned. The grouping is recomputed from the frozen "
                "split file every time the analysis runs.")
steps = [
    ("1", "Murcko scaffold per monomer",
     "RDKit MurckoScaffold.MurckoScaffoldSmiles — strips side chains, keeps ring systems and "
     "the linkers between them. Bemis & Murcko, J. Med. Chem. 1996. Deterministic, no free parameters."),
    ("2", "Group the 682 B monomers by identical scaffold",
     "Gives 112 families. Acyclic monomers, if any, form singletons."),
    ("3", "For each fold, ask whether the held-out chemistry is in training",
     "Read held_out_monomer_B and train_monomer_B from the frozen JSON; count held-out monomers "
     "whose scaffold also appears among that fold's training monomers."),
    ("4", "Label S if that share exceeds 50%, D otherwise",
     "Threshold fixed in advance. Also record distinct scaffolds per fold, and Morgan r=2 "
     "2048-bit Tanimoto novelty as an independent continuous check."),
]
for i, (n, head, body) in enumerate(steps):
    yy = y + i * 0.98
    circle(s, M, yy + 0.04, 0.42, n, BERRY, WHITE, 14)
    txt(s, M + 0.62, yy, 6.1, 0.3, head, size=14, bold=True, color=BERRY, font=HEAD)
    txt(s, M + 0.62, yy + 0.32, 6.1, 0.62, body, size=11, color=INK, font=BODY, line=1.18)

txt(s, 7.35, y - 0.02, 5.25, 0.3, "THE THRESHOLD IS NOT LOAD-BEARING", size=11, bold=True,
    color=MUTED, font=BODY)
rows = [
    ["fold", "share with scaffold in training", "group"],
    ["0 – 3", "1.00", "S"],
    ["6", "0.434", "D"],
    ["4", "0.171", "D"],
    ["5, 7, 8", "0.00", "D"],
]
plain_table(s, 7.35, y + 0.36, 5.25, rows, [1.3, 2.6, 0.9], row_h=0.42, head_h=0.5, size=11.5)
card(s, 7.35, y + 2.68, 5.25, 1.15, CREAM)
txt(s, 7.6, y + 2.86, 4.8, 0.85,
    "The nearest fold to the cut is 0.434. Any threshold between 0.44 and 0.99 produces the "
    "identical grouping — the data are not near the boundary.",
    size=11.5, color=INK, font=BODY, line=1.2)

txt(s, M, y + 4.15, 11.9, 0.55,
    [[("Code: ", {"bold": True}),
      ("scaffold_structure() in scripts/python/analyze_regen_v1_r3.py.   ", {"font": "Consolas"}),
      ("Outputs: ", {"bold": True}),
      ("_regen_v1_r3_results_fold_composition.csv", {"font": "Consolas"}),
      ("  ·  originally computed in audit_b_heldout_design.py during split construction.", {})]],
    size=11, color=MUTED, font=BODY, line=1.2)



# ============================================================ 19. BACKUP: MEAN
figure_slide(
    "Backup — the same result using means, not medians",
    "provenance",
    "fig_ab_comparison_mean",
    [[("Medians are robust to the pathological folds; means are not. On the A split HPG-hier's EA "
       "ΔR² is 0.803 median but 0.634 mean, because one fold falls to −0.002. On the B split the "
       "same numbers are 0.778 and 0.746. ", {}),
      ("The A folds are heterogeneous, the B folds are not — a second, independent argument for "
       "the B split.", {"bold": True, "color": BERRY})]])

# ============================================================ 20. BACKUP: BOTTOM LINE
figure_slide(
    "Backup — plain accuracy, no decomposition",
    "provenance",
    "fig_r1_overall_performance",
    [[("Overall R², RMSE and MAE per model on the A-heldout split — the bottom-line numbers with no "
       "null floors and no group/deviation split. The octamer leads on all three; wDMPNN trails on "
       "R² despite competitive MAE, which is the decomposition showing up in aggregate.", {})]])

prs.save(OUT)
print("wrote", OUT, "-", len(prs.slides.__iter__.__self__._sldIdLst), "slides")

