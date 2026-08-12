#!/usr/bin/env python3
"""Supervisor meeting deck — 12 August 2026.

Built from DECK_PLAN_supervisor_2026-08-12.md. Minimal style: white ground, one
accent colour, generous margins, one idea per slide. The "what to say" lines from
the plan live in speaker notes, not on the slides.
"""
import os
from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "Supervisor_update_2026-08-12.pptx")

INK   = C(0x1A, 0x1A, 0x1A)
BODY  = C(0x44, 0x4A, 0x4E)
MUTED = C(0x8A, 0x92, 0x96)
WHITE = C(0xFF, 0xFF, 0xFF)
ACCENT= C(0x0B, 0x6E, 0x7A)
WARM  = C(0xB2, 0x4B, 0x35)
RULE  = C(0xD8, 0xDC, 0xDE)
BAND  = C(0xF4, 0xF6, 0xF7)
DARK  = C(0x1E, 0x2A, 0x2F)

HEAD, BODYF = "Georgia", "Calibri"
W, H = 13.333, 7.5
M = 0.78

prs = Presentation()
prs.slide_width, prs.slide_height = In(W), In(H)
BLANK = prs.slide_layouts[6]


def _tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tf = s.shapes.add_textbox(In(x), In(y), In(w), In(h)).text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf


def _p(tf, text, size=13, color=BODY, bold=False, font=BODYF, space=7,
       first=False, italic=False, align=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = text
    p.space_after = Pt(space)
    if align:
        p.alignment = align
    for r in p.runs:
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = font; r.font.italic = italic
    return p


def _rect(s, x, y, w, h, fill):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, In(x), In(y), In(w), In(h))
    r.fill.solid(); r.fill.fore_color.rgb = fill
    r.line.fill.background(); r.shadow.inherit = False
    return r


def slide(title=None, kicker=None, sub=None, notes=None):
    s = prs.slides.add_slide(BLANK)
    if title:
        y = 0.62
        if kicker:
            tf = _tb(s, M, y, W - 2 * M, 0.26)
            _p(tf, kicker.upper(), size=10, color=ACCENT, bold=True, first=True, space=0)
            y += 0.34
        tf = _tb(s, M, y, W - 2 * M, 0.62)
        _p(tf, title, size=25, color=INK, bold=True, font=HEAD, first=True, space=0)
        y += 0.72
        _rect(s, M, y, 1.5, 0.022, ACCENT)
        y += 0.16
        if sub:
            tf = _tb(s, M, y, W - 2 * M - 0.4, 0.4)
            _p(tf, sub, size=12.5, color=MUTED, first=True, space=0)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def divider(n, title, sub):
    s = prs.slides.add_slide(BLANK)
    _rect(s, 0, 0, W, H, DARK)
    tf = _tb(s, M, 2.9, W - 2 * M, 0.3)
    _p(tf, f"SECTION {n}", size=11, color=C(0x6E, 0xA8, 0xB0), bold=True, first=True, space=0)
    tf = _tb(s, M, 3.35, W - 2 * M - 1.5, 0.9)
    _p(tf, title, size=32, color=WHITE, bold=True, font=HEAD, first=True, space=0)
    tf = _tb(s, M, 4.5, W - 2 * M - 2.4, 0.5)
    _p(tf, sub, size=13.5, color=C(0xA8, 0xB8, 0xBC), first=True, space=0)
    return s


def table(s, x, y, w, rows, colw, fs=10.5, hfs=10, rowh=0.31, bolds=(), hi=()):
    """rows[0] is the header. bolds = row indices to bold. hi = (row,col) to accent."""
    nr, nc = len(rows), len(rows[0])
    t = s.shapes.add_table(nr, nc, In(x), In(y), In(w), In(rowh * nr)).table
    for i, cw in enumerate(colw):
        t.columns[i].width = In(cw)
    for ri, row in enumerate(rows):
        t.rows[ri].height = In(rowh)
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.margin_left = In(0.08); cell.margin_right = In(0.05)
            cell.margin_top = In(0.02); cell.margin_bottom = In(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = DARK
            elif ri in bolds:
                cell.fill.fore_color.rgb = BAND
            else:
                cell.fill.fore_color.rgb = WHITE
            p = cell.text_frame.paragraphs[0]
            p.text = str(val)
            if ci > 0:
                p.alignment = PP_ALIGN.CENTER
            col = WHITE if ri == 0 else (ACCENT if (ri, ci) in hi else INK if ri in bolds else BODY)
            for r in p.runs:
                r.font.size = Pt(hfs if ri == 0 else fs)
                r.font.name = BODYF
                r.font.color.rgb = col
                r.font.bold = (ri == 0) or (ri in bolds) or ((ri, ci) in hi)
    return t


def bullets(s, x, y, w, items, size=13.5, gap=9, h=None):
    tf = _tb(s, x, y, w, h if h is not None else min(3.0, H - y - 0.6))
    for i, it in enumerate(items):
        _p(tf, it, size=size, color=BODY, first=(i == 0), space=gap)
    return tf


def keyline(s, text, y=6.35, color=DARK):
    _rect(s, M, y, W - 2 * M, 0.68, color)
    tf = _tb(s, M + 0.3, y + 0.16, W - 2 * M - 0.6, 0.42)
    _p(tf, text, size=13.5, color=WHITE, bold=True, font=HEAD, first=True, space=0)


def foot(s, text, color=MUTED):
    tf = _tb(s, M, H - 0.52, W - 2 * M, 0.3)
    _p(tf, text, size=9.5, color=color, first=True, space=0, italic=True)


# ══════════════════════════════════════════════════════════════════ TITLE ════
s = prs.slides.add_slide(BLANK)
_rect(s, 0, 0, W, H, DARK)
tf = _tb(s, M, 2.5, W - 2 * M, 0.3)
_p(tf, "RESEARCH UPDATE · 12 AUGUST 2026", size=11, color=C(0x6E, 0xA8, 0xB0),
   bold=True, first=True, space=0)
tf = _tb(s, M, 3.0, 10.4, 1.6)
_p(tf, "Evaluating whether polymer\nrepresentations encode architecture",
   size=32, color=WHITE, bold=True, font=HEAD, first=True, space=0)
tf = _tb(s, M, 5.05, 10.0, 0.9)
_p(tf, "1  The evaluation paper     2  Paper strategy and the dataset\n"
       "3  Why the octamer beats the baseline",
   size=13.5, color=C(0xA8, 0xB8, 0xBC), first=True, space=0)

# ═══════════════════════════════════════════════════════════════ SECTION 1 ════
divider(1, "The evaluation paper", "What it argues, and what is already written")

# ── 1.1 the claim ────────────────────────────────────────────────────────────
s = slide("What the wD-MPNN paper claims", kicker="Starting point",
          sub="Section heading: “The wD-MPNN captures how polymer properties depend on "
              "chain architecture and monomer stoichiometry”",
          notes="Open by conceding. Adding chain architecture on top of stoichiometry halves "
                "the error — 49% on EA, 58% on IP. That is not a marginal result. It looks "
                "small in R2 only because R2 saturates near 1. Their claim is correct and "
                "their evidence supports it. My argument is not that they were wrong.\n\n"
                "Tactically: lead with a critique of their numbers and the rest of the talk "
                "is a fight. Lead with a concession and the room hears the actual point.")
table(s, M, 2.35, 8.6,
      [["Representation", "EA R²", "EA RMSE", "IP R²", "IP RMSE"],
       ["Monomers only", "0.917", "0.173", "0.883", "0.165"],
       ["+ chain architecture", "0.929", "0.159", "0.898", "0.154"],
       ["+ stoichiometry", "0.987", "0.069", "0.982", "0.065"],
       ["+ both  (full wD-MPNN)", "0.997", "0.035", "0.997", "0.027"]],
      [3.4, 1.3, 1.3, 1.3, 1.3], bolds=(4,), rowh=0.36)
tf = _tb(s, 10.0, 2.4, 2.6, 2.0)
_p(tf, "Adding architecture on top of stoichiometry:", size=11, color=MUTED, first=True, space=8)
_p(tf, "−49%", size=26, color=ACCENT, bold=True, font=HEAD, space=0)
_p(tf, "EA error", size=11, color=BODY, space=10)
_p(tf, "−58%", size=26, color=ACCENT, bold=True, font=HEAD, space=0)
_p(tf, "IP error", size=11, color=BODY, space=0)
keyline(s, "Their claim is correct and their evidence supports it. That is not what I am "
           "arguing about.", y=4.55)
foot(s, "Aldeghi & Coley, Chem. Sci. 2022, 13, 10486–10498 — Table 2 and Table S1")

# ── 1.2 different regime ─────────────────────────────────────────────────────
s = slide("But that evidence came from a different regime", kicker="The gap",
          sub="The ablation runs under random 10-fold splits — their caption says so",
          notes="Under random splits architecture is about a seventh of what the model still "
                "has to explain, so an averaged metric picks it up. Under held-out-monomer "
                "splits it is one percent of the target, and a model that captured it "
                "perfectly would differ from one that ignored it by about two points of "
                "aggregate R2. Same benchmark, same property — the metric's ability to see "
                "architecture changes completely with the split.")
bullets(s, M, 2.3, 6.3, [
    "The benchmark has only 9 distinct A monomers. Under a random split every one "
    "of them appears in training, so monomer identity is effectively memorised.",
    "What is left to explain is composition and architecture — and in that reduced "
    "space architecture is large.",
    "Under a held-out-monomer split monomer identity becomes the dominant unknown. "
    "Their own RMSE goes from 0.027–0.035 eV to 0.09–0.10 eV.",
])
tf = _tb(s, M, 4.55, 6.3, 0.5)
_p(tf, "The architecture ablation was never run in that second regime.",
   size=14, color=WARM, bold=True, font=HEAD, first=True, space=0)
table(s, 7.5, 2.35, 5.05,
      [["Architecture as a share of…", "EA", "IP"],
       ["total variance   (held-out-monomer)", "0.98%", "1.46%"],
       ["what remains once monomers known   (random)", "13.9%", "14.1%"]],
      [3.15, 0.95, 0.95], rowh=0.52, fs=10)
keyline(s, "The evidence was obtained where aggregate accuracy can see architecture — "
           "and is cited for work evaluated where it cannot.")

# ── 1.3 metrics I ────────────────────────────────────────────────────────────
s = slide("How the two axes are computed", kicker="Method  ·  1 of 2",
          sub="Everything below comes from one unchanged set of predictions — no retraining",
          notes="A group is polymers sharing monomer A, monomer B and composition, so within "
                "a group only architecture differs. Groups with fewer than two architectures "
                "are discarded.")
tf = _tb(s, M, 2.3, 11.7, 0.5)
_p(tf, "A group g = all test polymers sharing monomer A, monomer B and composition. "
       "Within a group, only architecture differs.", size=13, color=BODY, first=True, space=0)
_rect(s, M, 2.95, 11.77, 1.05, BAND)
tf = _tb(s, M + 0.35, 3.16, 11.0, 0.7)
_p(tf, "yᵢ  =  ȳ_g(i)  +  δᵢ            where δᵢ = yᵢ − ȳ_g(i)",
   size=15, color=INK, bold=True, font="Consolas", first=True, space=4)
_p(tf, "           where the chemistry sits        what architecture does within it",
   size=10.5, color=MUTED, font="Consolas", space=0)

_rect(s, M, 4.35, 5.7, 1.75, BAND)
tf = _tb(s, M + 0.3, 4.55, 5.1, 1.4)
_p(tf, "AXIS 1 — CHEMISTRY PLACEMENT", size=9.5, color=ACCENT, bold=True, first=True, space=8)
_p(tf, "group_mean_R²  =  R²( ȳ_g , ŷ̄_g )", size=13, color=INK, bold=True,
   font="Consolas", space=8)
_p(tf, "Did the model put each chemistry at the right value?", size=11.5, color=BODY, space=0)

_rect(s, 7.0, 4.35, 5.55, 1.75, BAND)
tf = _tb(s, 7.3, 4.55, 4.95, 1.4)
_p(tf, "AXIS 2 — ARCHITECTURE RECOVERY", size=9.5, color=ACCENT, bold=True, first=True, space=8)
_p(tf, "ΔR²  =  R²( δᵢ , δ̂ᵢ )", size=13, color=INK, bold=True, font="Consolas", space=8)
_p(tf, "Cannot be improved by getting the chemistry right — the group mean is "
       "subtracted from both.", size=11.5, color=BODY, space=0)
keyline(s, "ΔR² removes exactly the term that buries architecture under held-out splits. "
           "It puts the model back in the 14% row.")

# ── 1.4 metrics II ───────────────────────────────────────────────────────────
s = slide("The null floor, and two supporting quantities", kicker="Method  ·  2 of 2",
          notes="The null predictor is scored with the identical metric on the identical "
                "fold. Because monomer A is held out by construction it cannot use any "
                "information about monomer A.")
tf = _tb(s, M, 2.15, 5.9, 2.6)
_p(tf, "NULL FLOOR", size=9.5, color=ACCENT, bold=True, first=True, space=8)
_p(tf, "A parameter-free lookup table, never trained. For a test polymer it returns the mean "
       "training value of polymers sharing monomer B, composition and architecture; failing "
       "that, monomer B and architecture; failing that, the global training mean.",
   size=12, color=BODY, space=9)
_p(tf, "skill  =  (R²_model − R²_null) / (1 − R²_null)", size=12, color=INK, bold=True,
   font="Consolas", space=7)
_p(tf, "0 = matched a lookup table.  A model can post R² = 0.98 and skill ≈ 0.",
   size=11.5, color=MUTED, space=0)

tf = _tb(s, 7.1, 2.15, 5.45, 2.6)
_p(tf, "ARCHITECTURE-SPREAD RATIO", size=9.5, color=ACCENT, bold=True, first=True, space=8)
_p(tf, "ratioₘ = (max ŷ − min ŷ) / (max y − min y)", size=11.5, color=INK, bold=True,
   font="Consolas", space=8)
_p(tf, "Of the true spread between architectures, how much does the model reproduce? "
       "1.00 is perfect; 0 = one value for all; above 1 = exaggeration.",
   size=12, color=BODY, space=9)
_p(tf, "Two-sided — 1.4 is as wrong as 0.71. Summarise with median |log₂ ratio|, "
       "where 0 is perfect.", size=11.5, color=WARM, space=0)

_rect(s, M, 5.05, 11.77, 0.028, RULE)
tf = _tb(s, M, 5.3, 11.7, 0.8)
_p(tf, "PROTOCOL", size=9.5, color=ACCENT, bold=True, first=True, space=7)
_p(tf, "Three seeds (42/43/44). Predictions averaged first, then the metric computed once — "
       "not the other way round. Best-checkpoint predictions, never the final "
       "patience-expired model. S and D folds never pooled.",
   size=12.5, color=BODY, space=0)

# ── 1.5 blind baselines ──────────────────────────────────────────────────────
s = slide("Baselines blind to the tested property should be standard",
          kicker="Proposal", sub="And the wD-MPNN paper already ran one — that is the point",
          notes="They ran the right comparison. They reported it honestly. What did not "
                "happen is anyone drawing the conclusion from it — it appears once, in "
                "passing, as a remark about that one dataset. I want to argue this is not "
                "optional. It should be computed on every fold and printed next to every "
                "score.\n\nFraming: not 'we found something they missed' — they didn't miss "
                "it — but 'we are proposing that a comparison they already ran becomes part "
                "of the protocol'. Much easier argument to win.")
tf = _tb(s, M, 2.3, 5.9, 1.2)
_p(tf, "Whenever a paper claims a representation encodes property X, it should report a "
       "baseline that cannot encode X — scored identically, on every fold.",
   size=13.5, color=BODY, first=True, space=0)
table(s, 6.9, 2.25, 5.65,
      [["Diblock phase dataset", "PRC"],
       ["Full wD-MPNN", "0.68"],
       ["Random forest, mole fractions alone", "0.69"],
       ["Random forest, volume fractions alone", "0.71"],
       ["Best RF (fingerprints + stoich + size)", "0.74"]],
      [4.15, 1.5], bolds=(2, 3), rowh=0.36)
tf = _tb(s, M, 3.9, 5.9, 1.4)
_p(tf, "A baseline with no architectural information at all scores above the "
       "architecture-aware model — on a dataset used to demonstrate architecture awareness.",
   size=13, color=WARM, bold=True, first=True, space=0)
keyline(s, "They ran the comparison and reported it. Nobody drew the conclusion. "
           "We propose making it mandatory.")

# ── 1.6 null floor result ────────────────────────────────────────────────────
s = slide("The null floor across every fold", kicker="Result",
          sub="Group-mean R², A split (held-out monomer A)",
          notes="Same split, same benchmark. On EA a lookup table gets 0.68; on IP it gets "
                "−0.03. You cannot know which situation you are in without computing it.\n\n"
                "On EA folds 2 and 3 it scores 0.96 and 0.95 — a model reporting 0.98 there "
                "has learned much less than the number suggests.\n\n"
                "And on IP fold 0 the lookup table beats both trained models.\n\n"
                "Negative R2 is unbounded below — −19.07 and −7.53 are real, not errors.")
rows = [["Fold", "EA null", "EA octamer", "EA wD-MPNN", "IP null", "IP octamer", "IP wD-MPNN"],
        ["0", "0.694", "0.958", "0.882", "0.969", "0.748", "0.496"],
        ["1", "0.487", "0.991", "0.982", "0.509", "0.970", "0.982"],
        ["2", "0.961", "0.986", "0.981", "−1.019", "0.995", "0.976"],
        ["3", "0.953", "0.995", "0.979", "−3.206", "0.982", "0.962"],
        ["4", "0.884", "0.962", "0.947", "−0.251", "0.960", "0.918"],
        ["5", "0.676", "0.994", "0.970", "−7.528", "0.977", "0.755"],
        ["6", "−19.069", "0.937", "0.878", "0.569", "0.996", "0.986"],
        ["7", "0.098", "0.982", "0.846", "0.410", "0.988", "0.482"],
        ["8", "0.428", "0.989", "0.992", "−0.034", "0.995", "0.987"],
        ["Median", "0.676", "0.984", "0.967", "−0.034", "0.978", "0.971"]]
table(s, M, 2.25, 8.4, rows, [1.2, 1.2, 1.2, 1.2, 1.2, 1.2, 1.2],
      bolds=(10,), rowh=0.335, fs=10,
      hi=((1, 4), (3, 1), (4, 1)))
tf = _tb(s, 9.5, 2.4, 3.1, 3.4)
_p(tf, "EA folds 2–3", size=11, color=ACCENT, bold=True, first=True, space=4)
_p(tf, "Lookup table scores 0.96 and 0.95 knowing nothing about the held-out monomer.",
   size=11.5, color=BODY, space=14)
_p(tf, "IP fold 0", size=11, color=WARM, bold=True, space=4)
_p(tf, "The lookup table (0.969) beats both trained models.",
   size=11.5, color=BODY, space=14)
_p(tf, "Median", size=11, color=ACCENT, bold=True, space=4)
_p(tf, "0.676 on EA against −0.034 on IP. Same split. You cannot tell without computing it.",
   size=11.5, color=BODY, space=0)

# ── 1.7 folds not exchangeable ───────────────────────────────────────────────
s = slide("The folds are not exchangeable", kicker="Split design",
          sub="62.5% of the 682 B monomers sit in just two scaffold families (317 and 109)",
          notes="Folds 0–3 are interpolation. Folds 5, 7 and 8 are extrapolation. Averaging "
                "them gives a number that describes neither. So S and D folds are reported "
                "separately throughout.")
table(s, M, 2.6, 11.77,
      [["Fold", "0", "1", "2", "3", "4", "5", "6", "7", "8"],
       ["Held-out monomers with a scaffold relative in training",
        "1.00", "1.00", "1.00", "1.00", "0.17", "0.00", "0.43", "0.00", "0.00"],
       ["Group", "S", "S", "S", "S", "D", "D", "D", "D", "D"]],
      [4.17, 0.845, 0.845, 0.845, 0.845, 0.845, 0.845, 0.845, 0.845, 0.845],
      rowh=0.42, fs=10.5)
tf = _tb(s, M, 4.4, 11.7, 1.2)
_p(tf, "Folds 0–3 test interpolation inside a known scaffold family. Folds 5, 7 and 8 test "
       "extrapolation to scaffolds never seen. These are different experiments.",
   size=13.5, color=BODY, first=True, space=0)
keyline(s, "A balanced scaffold-disjoint split cannot be built on this benchmark. "
           "So S and D folds are never pooled.")

# ── 1.8 headline A split ─────────────────────────────────────────────────────
s = slide("Under interpolation, every metric agrees", kicker="Headline  ·  1 of 3",
          sub="A split, 9 folds. Octamer minus wD-MPNN at its published configuration",
          notes="On the A split all five metrics favour the octamer on 8 or 9 of 9 folds. If "
                "that were the only regime we looked at, none of this machinery would be "
                "needed.")
table(s, M, 2.5, 11.77,
      [["Metric", "EA median", "EA mean", "EA wins", "p", "IP median", "IP mean", "IP wins", "p"],
       ["Overall R²", "+0.017", "+0.030", "9/9", "0.004", "+0.021", "+0.070", "8/9", "0.039"],
       ["MAE", "−0.018", "−0.023", "8/9", "0.039", "−0.026", "−0.026", "8/9", "0.039"],
       ["RMSE", "−0.024", "−0.026", "9/9", "0.004", "−0.033", "−0.030", "8/9", "0.039"],
       ["Group-mean R²", "+0.008", "+0.027", "8/9", "0.039", "+0.019", "+0.060", "7/9", "0.180"],
       ["ΔR²  (architecture)", "+0.190", "+0.236", "8/9", "0.039", "+0.278", "+0.267", "9/9", "0.004"]],
      [2.65, 1.14, 1.14, 1.14, 1.14, 1.14, 1.14, 1.14, 1.14],
      bolds=(5,), rowh=0.42)
keyline(s, "Every metric points the same way. Metric choice is immaterial here.", y=5.3)
foot(s, "Baseline at the authors' own settings — batch 50, 30 epochs. Minimum attainable "
        "p at 9 folds is 0.004.")

# ── 1.9 headline D folds ─────────────────────────────────────────────────────
s = slide("Under extrapolation, they come apart", kicker="Headline  ·  2 of 3",
          sub="B split, cross-scaffold D folds, 5 folds. Baseline at our harmonised configuration",
          notes="Four accuracy metrics split three-to-two, which is what a coin does. "
                "Architecture recovery splits five-to-zero on both properties. Same "
                "predictions.\n\nAnd mean and median disagree in sign for four of the four "
                "accuracy metrics. It is not just that they are insensitive — the answer "
                "depends on which average you pick. ΔR² does not do that.\n\nWith 5 folds "
                "the smallest achievable p is 0.0625, so this cannot reach 0.05 by design.")
table(s, M, 2.5, 11.77,
      [["Metric", "EA median", "EA mean", "EA wins", "IP median", "IP mean", "IP wins"],
       ["Overall R²", "+0.003", "−0.006", "3/5", "+0.010", "−0.014", "3/5"],
       ["MAE", "−0.000", "+0.001", "3/5", "−0.010", "−0.003", "3/5"],
       ["RMSE", "−0.004", "+0.005", "3/5", "−0.012", "+0.005", "3/5"],
       ["Group-mean R²", "+0.002", "−0.008", "3/5", "+0.007", "−0.016", "3/5"],
       ["ΔR²  (architecture)", "+0.260", "+0.225", "5/5", "+0.152", "+0.198", "5/5"]],
      [2.87, 1.48, 1.48, 1.48, 1.48, 1.48, 1.48], bolds=(5,), rowh=0.42)
tf = _tb(s, M, 5.15, 11.7, 0.9)
_p(tf, "Mean and median disagree in sign on four of the four accuracy metrics (EA) and "
       "three of four (IP). ΔR² does not.", size=13, color=WARM, bold=True, first=True, space=0)
keyline(s, "Same predictions. Four metrics say tie; architecture recovery says 5 of 5.", y=5.95)
foot(s, "Minimum attainable p at 5 folds is 0.0625 — this cannot reach 0.05 by design.")

# ── 1.10 the missing arm ─────────────────────────────────────────────────────
s = slide("What the one missing run should show", kicker="Headline  ·  3 of 3",
          sub="The published-config baseline has not been run on the B split — 54 jobs, cheap",
          notes="We expect the accuracy metrics essentially unchanged — still around three of "
                "five — because the two configurations are equivalent there. We expect the "
                "architecture gap to shrink but survive: roughly +0.19 on EA and +0.07 on IP, "
                "still five of five.\n\nIP is the thin case. Fold 5 projects to +0.023, close "
                "enough to zero that 5/5 could become 4/5. That is a reason to run it rather "
                "than assume.\n\nThis is a projection from an offset measured on another "
                "split, not a result.")
tf = _tb(s, M, 2.3, 5.5, 0.5)
_p(tf, "Both configurations exist on the A split, so the offset between them is measurable:",
   size=12.5, color=BODY, first=True, space=0)
table(s, M, 2.8, 5.5,
      [["Published minus 300-epoch config", "EA", "IP"],
       ["Overall R²", "−0.002", "+0.006"],
       ["MAE", "+0.000", "−0.003"],
       ["RMSE", "+0.004", "−0.005"],
       ["Group-mean R²", "+0.001", "+0.001"],
       ["ΔR²", "+0.067", "+0.081"]],
      [3.1, 1.2, 1.2], bolds=(5,), rowh=0.335, fs=10)
tf = _tb(s, 6.75, 2.3, 5.8, 0.5)
_p(tf, "Applying that offset to the D-fold ΔR² we already have:",
   size=12.5, color=BODY, first=True, space=0)
table(s, 6.75, 2.8, 5.8,
      [["ΔR², D folds", "median", "wins"],
       ["EA  observed", "+0.260", "5/5"],
       ["EA  projected", "+0.193", "5/5"],
       ["IP  observed", "+0.152", "5/5"],
       ["IP  projected", "+0.071", "5/5"]],
      [3.0, 1.4, 1.4], bolds=(2, 4), rowh=0.335, fs=10)
tf = _tb(s, M, 4.9, 11.7, 0.9)
_p(tf, "Accuracy metrics should be unchanged — the two configurations are equivalent there. "
       "The architecture gap should shrink but survive. IP fold 5 projects to +0.023, so 5/5 "
       "could become 4/5 — which is the reason to run it rather than assume.",
   size=12.5, color=BODY, first=True, space=0)
keyline(s, "The published configuration is the stronger baseline on the axis we care about.",
        y=5.95, color=WARM)
foot(s, "Projection from an offset measured on a different split. Not a result.")

# ═══════════════════════════════════════════════════════════════ SECTION 2 ════
divider(2, "Paper strategy", "Four papers, and the dataset the argument requires")

s = slide("Four papers", kicker="Sequence",
          sub="The through-line comes from the accepted review's own Future Perspectives section",
          notes="Why Paper 1 first: it is writable today; it does not depend on which model "
                "wins, which insulates it from the protocol-matching confound; and the "
                "measurement work invalidates most earlier model comparisons, so it is the "
                "stronger contribution.\n\nWhat Paper 1 does not claim: that our model is "
                "best. That belongs to Paper 3.")
table(s, M, 2.45, 11.77,
      [["", "Paper", "Status", "Note"],
       ["0", "Review — accepted, Digital Discovery", "done",
        "Argues in the literature that representation quality needs more than accuracy"],
       ["1", "Evaluation framework", "write now",
        "Demonstrates it empirically. One cheap run outstanding"],
       ["2", "Architecture-aware benchmark dataset", "costed",
        "Builds the data the argument requires — and resolves our own confound"],
       ["3", "The model — octamer vs baseline", "evidence in hand",
        "Held until external validity is tested"]],
      [0.5, 3.6, 1.7, 5.97], bolds=(2,), rowh=0.52, fs=11)
bullets(s, M, 5.05, 11.7, [
    "Paper 1 goes first because it is writable today, it does not depend on which model "
    "wins, and the measurement work invalidates most of the earlier model comparisons.",
], size=13)
keyline(s, "Paper 1 deliberately does not claim our model is best. That belongs to Paper 3.")

s = slide("What Paper 2 builds", kicker="The ask",
          sub="The existing benchmark has three architecture settings — literally three "
              "distinct transition matrices across 42,966 rows",
          notes="The authors published their generation code; it encodes a real Suzuki "
                "coupling and our monomers are the right type. We have read it and verified "
                "it runs on monomers we already hold.\n\nCompositions are quarters, so 6 and "
                "10 units cannot represent them exactly and would look worse for purely "
                "arithmetic reasons. 8, 12, 16 and 24 are the clean comparisons.\n\nThis is "
                "CPU, not GPU — a different queue from model training.")
table(s, M, 2.55, 11.77,
      [["", "Change", "Purpose"],
       ["1", "Blockiness varies continuously at fixed chemistry and composition",
        "Turns a 3-way classification into a real measurement axis"],
       ["2", "Publish the un-averaged per-chain values, not only the ensemble average",
        "Costs nothing extra — the calculation already produces them"],
       ["3", "Sweep the label chain length: 8, 12, 16 and 24 units",
        "Resolves our own protocol-matching confound"]],
      [0.5, 5.7, 5.57], bolds=(3,), rowh=0.56, fs=11)
tf = _tb(s, M, 4.9, 5.7, 1.2)
_p(tf, "COST", size=9.5, color=ACCENT, bold=True, first=True, space=7)
_p(tf, "~2,000 polymers · ≤512,000 structures · ~30 CPU-seconds each  →  ≈ 8–9 kSU CPU. "
       "That is a ceiling; a one-day pilot would measure the real figure.",
   size=12, color=BODY, space=0)
tf = _tb(s, 7.1, 4.9, 5.45, 1.2)
_p(tf, "FOR SCALE", size=9.5, color=ACCENT, bold=True, first=True, space=7)
_p(tf, "Reproducing the full 42,966 would be ~150 kSU. We are proposing about 5% of that — "
       "and it is CPU, not GPU.", size=12, color=BODY, space=0)
keyline(s, "The pipeline is published, verified portable, and runs on monomers we already hold.")

# ═══════════════════════════════════════════════════════════════ SECTION 3 ════
divider(3, "Why does the octamer win?", "What we have excluded, and what remains")

s = slide("The result being explained", kicker="Starting point",
          sub="Median over 9 A-split folds, baseline at its published configuration",
          notes="The accuracy gap is real but modest. The architecture-recovery gap is more "
                "than double.")
table(s, M, 2.6, 11.77,
      [["", "EA overall R²", "EA MAE", "EA ΔR²", "IP overall R²", "IP MAE", "IP ΔR²"],
       ["HPG-octamer", "0.984", "0.055", "0.849", "0.978", "0.035", "0.886"],
       ["wD-MPNN (published cfg)", "0.967", "0.070", "0.397", "0.971", "0.050", "0.565"]],
      [3.17, 1.43, 1.43, 1.43, 1.43, 1.44, 1.44], bolds=(1,), rowh=0.46, fs=12)
tf = _tb(s, M, 4.5, 11.7, 0.6)
_p(tf, "The accuracy gap is real but modest. The architecture-recovery gap is more than double.",
   size=14, color=BODY, first=True, space=0)
keyline(s, "So the question is: which part of the octamer's design is doing that?")

s = slide("What the difference looks like mechanically", kicker="Diagnosis",
          sub="Median across all 682-group folds",
          notes="On EA the two models place the chemistry equally well — 0.047 against 0.048. "
                "The whole difference is in how much architectural variation they reproduce. "
                "And wD-MPNN is not systematically flattening: it flattens to 0.25 on some "
                "folds and exaggerates to 1.48 on others. Unreliable in both directions "
                "rather than biased in one.")
table(s, M, 2.5, 7.2,
      [["", "Placement error (eV)", "", "Spread ratio", ""],
       ["", "octamer", "wD-MPNN", "octamer", "wD-MPNN"],
       ["EA", "0.047", "0.048", "0.77", "0.61"],
       ["IP", "0.029", "0.045", "0.95", "0.55"]],
      [1.2, 1.5, 1.5, 1.5, 1.5], rowh=0.36, fs=11)
table(s, 8.6, 2.5, 3.95,
      [["Distance from 1  (0 = perfect)", "oct", "wD"],
       ["EA", "0.42", "0.71"],
       ["IP", "0.15", "0.86"]],
      [2.15, 0.9, 0.9], rowh=0.36, fs=11)
tf = _tb(s, M, 4.5, 11.7, 1.3)
_p(tf, "On EA the two models place the chemistry equally well — 0.047 against 0.048. "
       "The whole difference is in how much architectural variation they reproduce.",
   size=13.5, color=BODY, first=True, space=9)
_p(tf, "wD-MPNN is not systematically flattening: it flattens to 0.25 on some folds and "
       "exaggerates to 1.48 on others — unreliable in both directions rather than biased "
       "in one.", size=13, color=WARM, space=0)

s = slide("Five differences, tested one at a time", kicker="Attribution",
          sub="Each elimination pre-registered before running",
          notes="The octamer differs from our own HPG-hier in five ways at once.")
table(s, M, 2.6, 11.77,
      [["", "Factor", "Status"],
       ["1", "8-slot chain instead of a 2-node graph", "open"],
       ["2", "Learned position embeddings", "excluded — next slide but one"],
       ["3", "Attention readout instead of stoichiometry-weighted", "open"],
       ["4", "Discards the 16-d port-pair edge features", "open"],
       ["5", "16 sampled sequences averaged, instead of 1", "excluded — next slide"]],
      [0.5, 7.3, 3.97], bolds=(2, 5), rowh=0.46, fs=12)
keyline(s, "Two excluded, three open — and two of the three are confounded with each other.")

s = slide("Factor 5 — the 16-sequence ensemble is not doing the work", kicker="Excluded",
          sub="Retrained with K = 1: one fixed sequence per polymer, never resampled. B split, "
              "both targets, three seeds",
          notes="Dropping from sixteen sampled sequences to one changed nothing measurable on "
                "any metric — accuracy, error, chemistry placement or architecture recovery. "
                "The averaging is not where the advantage comes from. That is a negative "
                "result and I am reporting it as one.\n\nDisclosed: the seed-stability half "
                "used a criterion not defined in advance, so that part is inconclusive. IP "
                "fold 7 is a 2-seed cell — one run missing, being completed now.")
table(s, M, 2.7, 11.77,
      [["Metric  (K=1 minus K=16) — full range across folds",
        "EA · S folds", "EA · D folds", "IP · S folds", "IP · D folds"],
       ["Overall R²", "−0.010, +0.001", "−0.006, +0.018", "−0.003, +0.003", "−0.006, +0.010"],
       ["RMSE", "−0.002, +0.011", "−0.015, +0.007", "−0.004, +0.003", "−0.006, +0.008"],
       ["MAE", "−0.001, +0.008", "−0.010, +0.002", "−0.000, +0.002", "−0.004, +0.005"],
       ["Group-mean R²", "−0.010, +0.001", "−0.006, +0.019", "−0.004, +0.003", "−0.006, +0.010"],
       ["ΔR²", "−0.002, +0.029", "−0.024, +0.032", "−0.019, +0.005", "−0.062, +0.025"]],
      [3.37, 2.1, 2.1, 2.1, 2.1], bolds=(5,), rowh=0.4, fs=10)
tf = _tb(s, M, 5.35, 11.7, 0.8)
_p(tf, "Ranges, not medians — a median can hide a wide spread. The largest change on any "
       "fold, on any metric, is 0.062. Pre-registered threshold ±0.024; pre-registered "
       "outcome C. No sign test reaches significance.",
   size=12.5, color=BODY, first=True, space=0)
keyline(s, "The biggest single-fold change anywhere is 0.062. The averaging is not where the "
           "advantage comes from.")

s = slide("Factor 2 — position embeddings are not doing the work either", kicker="Excluded",
          sub="Retrained with the 8 learned slot vectors removed. A split, 9 folds, both "
              "targets, three seeds — 54 runs, completed this week",
          notes="Removing the position embeddings moved nothing — every metric's median is "
                "within four thousandths of zero except ΔR², which moved by one hundredth "
                "against a threshold of five hundredths fixed before running.\n\nCaveats: "
                "scope is R1 only, the B split has not been run. EA fold 6's +0.105 is the "
                "largest number in the arm but that cell's baseline ΔR² is −0.14 with a seed "
                "SD of 0.81 — it is noise. And EA fold 0 moves in opposite directions on the "
                "two axes; on one fold at these magnitudes that is not a finding.\n\n"
                "Say the mechanism precisely: without position embeddings, slots holding the "
                "same monomer get identical embeddings, but end and interior slots still "
                "differ through path structure. This is a REDUCTION of positional "
                "information, not its elimination. Do not say 'position-blind'.")
table(s, M, 2.75, 11.77,
      [["Metric  (ablated minus baseline)", "EA median", "EA range across folds", "EA p",
        "IP median", "IP range across folds", "IP p"],
       ["Overall R²", "−0.002", "−0.008, +0.035", "0.51", "+0.002", "−0.032, +0.088", "0.51"],
       ["RMSE", "+0.001", "−0.060, +0.012", "0.51", "−0.002", "−0.032, +0.013", "0.51"],
       ["MAE", "+0.001", "−0.051, +0.011", "0.51", "−0.004", "−0.027, +0.013", "0.51"],
       ["Group-mean R²", "−0.001", "−0.008, +0.037", "0.51", "+0.001", "−0.032, +0.092", "0.51"],
       ["ΔR²", "−0.010", "−0.053, +0.105", "0.51", "−0.010", "−0.108, +0.051", "1.00"]],
      [2.85, 1.35, 2.1, 0.9, 1.35, 2.1, 1.12], bolds=(5,), rowh=0.4, fs=10)
tf = _tb(s, M, 5.4, 11.7, 0.9)
_p(tf, "The medians sit well inside the pre-registered ±0.051 band — that was the criterion, "
       "fixed before running, and it gives outcome 3. But individual folds swing by up to "
       "0.11, so the honest claim is “no consistent effect”, not “no effect on any fold”.",
   size=12.5, color=BODY, first=True, space=0)
keyline(s, "No consistent effect. Two known-unstable cells account for most of the spread.")

s = slide("Why we report all five metrics on the ablations too", kicker="Defence",
          notes="This is the answer to the obvious objection — that I invented a metric that "
                "flatters my own model. If ΔR² simply reported differences wherever you "
                "pointed it, it would have found something in these two ablations. It "
                "didn't. It agrees with accuracy when there is genuinely nothing there, and "
                "separates from it only when there is something accuracy cannot see.")
_rect(s, M, 2.4, 5.7, 1.9, BAND)
tf = _tb(s, M + 0.32, 2.65, 5.1, 1.5)
_p(tf, "ON BOTH ABLATIONS", size=9.5, color=ACCENT, bold=True, first=True, space=8)
_p(tf, "Every metric agrees — conventional and new alike, all say “no change”.",
   size=14, color=INK, bold=True, font=HEAD, space=0)
_rect(s, 7.1, 2.4, 5.45, 1.9, BAND)
tf = _tb(s, 7.42, 2.65, 4.85, 1.5)
_p(tf, "ON THE CROSS-SCAFFOLD FOLDS", size=9.5, color=WARM, bold=True, first=True, space=8)
_p(tf, "Accuracy says “tie”. Architecture recovery says 5 of 5.",
   size=14, color=INK, bold=True, font=HEAD, space=0)
tf = _tb(s, M, 4.65, 11.7, 1.2)
_p(tf, "If ΔR² simply reported differences wherever you pointed it, it would have found "
       "something in those two ablations. It did not. It agrees with accuracy when there is "
       "nothing there, and separates from it only when there is something accuracy cannot see.",
   size=13.5, color=BODY, first=True, space=0)
keyline(s, "That is the answer to “you invented a metric that flatters your own model”.")

s = slide("A correction to flag", kicker="Honesty",
          sub="Our own pre-registration undercounts what remains open",
          notes="The pre-registration's outcome-3 text says the remaining candidates are "
                "factors 1 and 4. Factor 3 — the attention readout — has never been tested. "
                "I checked every prediction directory; only two of the four cells exist.\n\n"
                "This needs a dated addendum, not an edit.")
table(s, M, 2.6, 9.0,
      [["", "Stoichiometry-weighted readout", "Attention readout"],
       ["2-node graph", "HPG-hier — have", "arm C — never run"],
       ["8-slot chain", "arm D — never run", "octamer — have"]],
      [2.4, 3.3, 3.3], rowh=0.48, fs=11.5)
tf = _tb(s, M, 4.5, 11.7, 1.3)
_p(tf, "So factors 1, 3 and 4 are open — not two — and factors 1 and 3 are confounded with "
       "each other. Arms C and D separate them; only the pair is informative.",
   size=13.5, color=BODY, first=True, space=9)
_p(tf, "Arm C needs no code change. Arm D needs a small patch. Pilot is 12 runs; the full "
       "R1 arm would be 108 runs, ~3.9 kSU.", size=13, color=BODY, space=0)
keyline(s, "This needs a dated addendum to the pre-registration — not an edit.")

s = slide("The confound, and how Paper 2 resolves it", kicker="The honest limit",
          sub="Labels were computed on 8-unit chains averaged over up to 32 sequences. "
              "The octamer uses an 8-slot chain, 16 sequences, averaged",
          notes="Two readings, both legitimate. Favourable: we encoded the right physics. "
                "Sceptical: the advantage is alignment with this dataset's recipe.\n\nNo "
                "ablation on this dataset can separate them — every label was made at 8 "
                "units, so there is no variation to exploit. But a dataset can, and it is "
                "the one we already want to build.\n\nCheaper things first: the "
                "glass-transition and block-copolymer datasets are already in the repository "
                "and their labels have nothing to do with 8-unit averaging. And freezing the "
                "model to test whether architecture is linearly decodable separates "
                "'predicts better' from 'represents better'.")
table(s, M, 2.75, 11.77,
      [["If the octamer's advantage is…", "then as label chain length moves away from 8…"],
       ["protocol matching", "performance peaks sharply at 8 and degrades at 12, 16, 24"],
       ["correct physics", "performance is flat, or improves with longer chains"]],
      [4.4, 7.37], rowh=0.5, fs=12)
tf = _tb(s, M, 4.6, 11.7, 1.4)
_p(tf, "This is change 3 on the Paper 2 slide. The same monomers, the same pipeline, a subset "
       "run at four chain lengths instead of one — a clean, pre-registerable discriminator "
       "at no extra cost.", size=13.5, color=BODY, first=True, space=9)
_p(tf, "Cheaper first: test on the glass-transition and block-copolymer datasets we already "
       "hold, whose labels have nothing to do with 8-unit averaging.",
   size=13, color=MUTED, space=0)
keyline(s, "Paper 2 gives the field a continuous architecture axis — and settles our own "
           "confound.")

# ── closing ──────────────────────────────────────────────────────────────────
s = slide("What I am asking for", kicker="Decisions")
items = [
    ("1", "Approval to run the published-config baseline on the B split",
     "54 runs at 30 epochs. Closes the last hole in Paper 1."),
    ("2", "A decision on the Paper 2 dataset",
     "~8–9 kSU CPU, not GPU. Also resolves our protocol-matching confound."),
    ("3", "A view on arms C and D — now, or after Paper 1 is submitted?",
     "12-run pilot first; the full arm is 108 runs, ~3.9 kSU."),
]
y = 2.45
for n, h, b in items:
    _rect(s, M, y, 11.77, 1.12, BAND)
    tf = _tb(s, M + 0.35, y + 0.2, 0.4, 0.5)
    _p(tf, n, size=19, color=ACCENT, bold=True, font=HEAD, first=True, space=0)
    tf = _tb(s, M + 0.95, y + 0.19, 10.6, 0.8)
    _p(tf, h, size=14.5, color=INK, bold=True, font=HEAD, first=True, space=5)
    _p(tf, b, size=12, color=BODY, space=0)
    y += 1.28
keyline(s, "Nothing else blocks Paper 1. The draft is written.", y=6.4)

prs.save(OUT)
print(f"saved {OUT}")
