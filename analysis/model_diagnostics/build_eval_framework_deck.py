#!/usr/bin/env python3
"""Evaluation-framework deck — the six components of the proposed Paper 1.

Every figure is an existing artefact from 29-07-2026 supervisor_update/figures/
(post-fix, three-seed regenerated). Every number traces to a repository report.
"""
import os
from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

FIG = "figs"

DARK  = C(0x2B, 0x3A, 0x42)
INK   = C(0x1A, 0x24, 0x2A)
BODY  = C(0x3D, 0x4C, 0x54)
MUTED = C(0x77, 0x86, 0x8E)
WHITE = C(0xFF, 0xFF, 0xFF)
TEAL  = C(0x02, 0x80, 0x90)
AMBER = C(0xC4, 0x8B, 0x0E)
TERRA = C(0xB8, 0x50, 0x42)
SAGE  = C(0xA7, 0xBE, 0xAE)
TINT  = C(0xF2, 0xF5, 0xF6)
TINT2 = C(0xE8, 0xEF, 0xF1)

HEAD, BODYF, MONO = "Cambria", "Calibri", "Courier New"
W, H, M = 13.333, 7.5, 0.55

prs = Presentation()
prs.slide_width, prs.slide_height = In(W), In(H)
BLANK = prs.slide_layouts[6]


def slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    if dark:
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, In(W), In(H))
        bg.fill.solid(); bg.fill.fore_color.rgb = DARK
        bg.line.fill.background(); bg.shadow.inherit = False
    return s


def tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf


def para(tf, text, size=13, color=BODY, bold=False, font=BODYF, space=6,
         first=False, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = text; p.space_after = Pt(space)
    for r in p.runs:
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = font; r.font.italic = italic
    return p


def card(s, x, y, w, h, fill=TINT):
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, In(x), In(y), In(w), In(h))
    r.fill.solid(); r.fill.fore_color.rgb = fill
    r.line.fill.background(); r.shadow.inherit = False
    r.adjustments[0] = 0.05
    return r


def numdot(s, x, y, d, n, fill):
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, In(x), In(y), In(d), In(d))
    o.fill.solid(); o.fill.fore_color.rgb = fill
    o.line.fill.background(); o.shadow.inherit = False
    tf = o.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = str(n); p.alignment = PP_ALIGN.CENTER
    for r in p.runs:
        r.font.size = Pt(14); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = BODYF
    return o


def piece_title(s, n, text, sub):
    numdot(s, M, 0.44, 0.5, n, TEAL)
    tf = tb(s, M + 0.7, 0.42, W - 2 * M - 0.7, 0.6)
    para(tf, text, size=25, color=INK, bold=True, font=HEAD, space=2, first=True)
    tf2 = tb(s, M + 0.7, 1.12, W - 2 * M - 0.7, 0.4)
    para(tf2, sub, size=12.5, color=MUTED, first=True, space=0)


def fig(s, name, x, y, w):
    p = os.path.join(FIG, name)
    if os.path.exists(p):
        return s.shapes.add_picture(p, In(x), In(y), width=In(w))
    ph = card(s, x, y, w, 2.2, TINT2)
    tf = tb(s, x + 0.2, y + 0.9, w - 0.4, 0.4)
    para(tf, f"[FIGURE MISSING: {name}]", size=12, color=TERRA, first=True, space=0)
    return ph


def caption(s, x, y, w, text):
    tf = tb(s, x, y, w, 0.35)
    para(tf, text, size=9.5, color=MUTED, first=True, space=0, italic=True)


def defbox(s, x, y, w, h, label, lines, fill=TINT):
    card(s, x, y, w, h, fill)
    tf = tb(s, x + 0.28, y + 0.2, w - 0.56, 0.3)
    para(tf, label, size=10, color=TEAL, bold=True, first=True, space=9)
    tf2 = tb(s, x + 0.28, y + 0.58, w - 0.56, h - 0.75)
    for i, (t, kw) in enumerate(lines):
        para(tf2, t, first=(i == 0), **kw)


def foot(s, text, color=MUTED):
    tf = tb(s, M, H - 0.48, W - 2 * M, 0.3)
    para(tf, text, size=9.5, color=color, first=True, space=0, italic=True)


B = dict(size=12, color=BODY, space=7)
BB = dict(size=12, color=INK, space=7, bold=True)
MONOS = dict(size=12.5, color=INK, space=8, font=MONO, bold=True)

# ── S1 title ─────────────────────────────────────────────────────────────────
s = slide(dark=True)
tf = tb(s, M, 2.0, W - 2 * M, 0.4)
para(tf, "PAPER 1 — PROPOSED FRAMING", size=12, color=TEAL, bold=True, first=True, space=0)
tf = tb(s, M, 2.55, 11.2, 2.0)
para(tf, "Current evaluation cannot see\nthe property the field claims to model",
     size=33, color=WHITE, bold=True, font=HEAD, first=True, space=0)
tf = tb(s, M, 4.9, 10.6, 1.0)
para(tf, "Six components of an evaluation framework for polymer representation "
         "learning. All six are built, measured, and demonstrated on data we already hold.",
     size=14, color=SAGE, first=True, space=0)
tf = tb(s, M, 6.4, 11.5, 0.4)
para(tf, "No new compute required to write this paper.", size=12, color=MUTED,
     first=True, space=0, italic=True)
s.notes_slide.notes_text_frame.text = (
    "Framing: this is not 'here are better metrics'. It is 'the metric everyone "
    "reports cannot answer the question everyone is asking', and here is the evidence.")

# ── S2 the problem ───────────────────────────────────────────────────────────
s = slide()
tf = tb(s, M, 0.42, W - 2 * M, 0.7)
para(tf, "The problem, and the six things needed to fix it", size=26, color=INK,
     bold=True, font=HEAD, first=True, space=0)
tf = tb(s, M, 1.18, W - 2 * M, 0.4)
para(tf, "Chain architecture is what we claim to model. It is ~1% of the variance "
         "in the target.", size=13, color=MUTED, first=True, space=0)

items = [("Variance decomposition", "Establishes that the signal is invisible to "
          "the standard metric"),
         ("Two-axis metric", "Separates chemistry placement from architecture recovery"),
         ("Null floor", "The best achievable without the thing being tested"),
         ("Split design", "Two chemistry-extrapolation regimes, folds not exchangeable"),
         ("Noise floor", "What difference is measurable at all"),
         ("The demonstration", "Two metrics, same predictions, opposite conclusions")]
x, y = M, 1.85
for i, (h, b) in enumerate(items, 1):
    card(s, x, y, 3.94, 1.55, TINT if i % 2 else TINT2)
    numdot(s, x + 0.28, y + 0.26, 0.42, i, TEAL if i < 6 else TERRA)
    tf = tb(s, x + 0.28, y + 0.82, 3.4, 0.65)
    para(tf, h, size=13.5, color=INK, bold=True, font=HEAD, first=True, space=4)
    para(tf, b, size=11, color=BODY, space=0)
    x += 4.14
    if i == 3:
        x, y = M, y + 1.75

card(s, M, 5.42, 12.23, 1.15, DARK)
tf = tb(s, M + 0.35, 5.62, 11.6, 0.8)
para(tf, "Together they answer a question the field currently cannot: is a "
         "representation better, or does it merely score higher?",
     size=14.5, color=WHITE, bold=True, font=HEAD, first=True, space=0)

# ── S3 piece 1 ───────────────────────────────────────────────────────────────
s = slide()
piece_title(s, 1, "Variance decomposition — the signal is 1% of the target",
            "Why the standard metric cannot see architecture")
defbox(s, M, 1.75, 4.85, 2.35, "WHAT IS MEASURED",
       [("Partition target variance into the part explained by monomer identity "
         "and composition, and the part attributable to chain architecture "
         "at fixed (A, B, fracA).", B)])
card(s, M, 4.25, 4.85, 2.35, TINT2)
tf = tb(s, M + 0.28, 4.45, 4.3, 2.0)
para(tf, "Result", size=10, color=TEAL, bold=True, first=True, space=9)
para(tf, "Architecture is 0.98% (EA) and 1.46% (IP) of total variance —", **BB)
para(tf, "but 51–60% of what remains after composition is accounted for.", **BB)
para(tf, "So a model can score 0.98 R² while being blind to architecture. "
         "Reporting overall R² cannot distinguish the two.", **B)
fig(s, "fig2_variance_by_axis.png", 5.75, 1.9, 7.05)
caption(s, 5.75, 5.55, 7.05,
        "fig2_variance_by_axis · corroborated independently four months apart "
        "(_dataset_design_audit.md, and the May architecture diagnostic)")
foot(s, "This is the premise for everything that follows: the quantity of interest is buried under a 90% chemistry term.")

# ── S4 piece 2 ───────────────────────────────────────────────────────────────
s = slide()
piece_title(s, 2, "Two-axis metric — separate placement from recovery",
            "One number per axis, always reported together")
defbox(s, M, 1.75, 4.85, 2.6, "DEFINITION",
       [("Group  g = (monomer A, monomer B, fracA).", B),
        ("group-mean R²  —  R² of the per-group means.", MONOS),
        ("Can the model put a chemistry in the right place?", B),
        ("ΔR²  —  R² of the within-group deviations,", MONOS),
        ("y − ȳ_g. Can it rank architectures inside a group?", B)])
card(s, M, 4.5, 4.85, 2.1, TINT2)
tf = tb(s, M + 0.28, 4.7, 4.3, 1.75)
para(tf, "Why both", size=10, color=TEAL, bold=True, first=True, space=9)
para(tf, "The chemistry term cancels inside a group, so ΔR² isolates architecture. "
         "Reported alone, either axis is misleading — a model can win one and lose "
         "the other, which is exactly what we observe.", **B)
fig(s, "fig_ab_comparison.png", 5.75, 2.2, 7.05)
caption(s, 5.75, 5.6, 7.05,
        "fig_ab_comparison · both axes, three evaluation regimes, five models")
foot(s, "A third quantity, pairwise ordering accuracy, is also reported — it is rank-based and less sensitive to scale.")

# ── S5 piece 3 — the centrepiece ─────────────────────────────────────────────
s = slide()
piece_title(s, 3, "Null floor — the best achievable without the thing being tested",
            "The component I would lead the paper with")
defbox(s, M, 1.75, 4.6, 2.5, "DEFINITION",
       [("An A-blind (or B-blind) group-mean lookup table. No parameters, "
         "no fitting, and it never sees the held-out monomer.", B),
        ("skill = (R²_model − R²_null) / (1 − R²_null)", MONOS),
        ("0 = no better than knowing nothing about the held-out chemistry. "
         "1 = perfect. Scale-free, so comparable across folds and targets.", B)])
card(s, M, 4.4, 4.6, 2.2, TINT2)
tf = tb(s, M + 0.28, 4.6, 4.05, 1.85)
para(tf, "Result", size=10, color=TERRA, bold=True, first=True, space=9)
para(tf, "On the A split, EA: the null reaches median group-mean R² 0.676 — "
         "and on one fold it beats the trained model, 0.961 vs 0.922.", **BB)
para(tf, "The same null on IP scores −0.034. The EA metric on this split is "
         "near-degenerate, and you cannot tell without a null.", **B)
fig(s, "fig_r1_skill_vs_null.png", 5.5, 1.78, 7.3)
caption(s, 5.5, 6.4, 7.3,
        "fig_r1_skill_vs_null · skill collapses to 0.32–0.55 on folds 2–4 — "
        "those models are barely better than knowing nothing")
foot(s, "Nobody in this literature reports a null floor. This is a finding about the field, not about our model.")

# ── S6 piece 4 ───────────────────────────────────────────────────────────────
s = slide()
piece_title(s, 4, "Split design — two extrapolation regimes, folds not exchangeable",
            "And an honest account of what each split can support")
defbox(s, M, 1.75, 4.85, 2.35, "THE TWO SPLITS",
       [("R1 · monomer-A-held-out — 9 folds, one A monomer held out each. "
         "Reproduces the original paper's design.", B),
        ("R3 · monomer-B-held-out, scaffold-clustered — 9 folds over 682 B "
         "monomers, capacity-balanced Murcko scaffold packing.", B)])
card(s, M, 4.25, 4.85, 2.35, TINT2)
tf = tb(s, M + 0.28, 4.45, 4.3, 2.0)
para(tf, "The finding that forces the fold grouping", size=10, color=TERRA,
     bold=True, first=True, space=9)
para(tf, "112 Murcko scaffolds, the largest two covering 62.5% of all B monomers.", **BB)
para(tf, "A balanced scaffold-disjoint split is therefore impossible and the nine "
         "folds are not exchangeable. We derive S (within-scaffold, folds 0–3) and "
         "D (cross-scaffold, 4–8) and never pool them.", **B)
fig(s, "fig3_scaffold_cluster_sizes.png", 5.9, 1.95, 6.6)
caption(s, 5.9, 6.1, 6.6,
        "fig3_scaffold_cluster_sizes · two families dominate the monomer-B space")
foot(s, "Pooling nine non-exchangeable folds into one mean is the error this component exists to prevent.")

# ── S7 piece 5 ───────────────────────────────────────────────────────────────
s = slide()
piece_title(s, 5, "Noise floor — what difference is measurable at all",
            "Measured, not assumed")
defbox(s, M, 1.75, 4.85, 2.2, "PROTOCOL",
       [("Repeat runs identical in model, seed, split, code and GPU. "
         "Report the spread of every metric across repeats, and the per-cell "
         "across-seed SD for every model.", B)])
card(s, M, 4.1, 4.85, 2.5, TINT2)
tf = tb(s, M + 0.28, 4.3, 4.3, 2.15)
para(tf, "Result", size=10, color=TERRA, bold=True, first=True, space=9)
para(tf, "Three identical runs, EA fold 1:", **B)
para(tf, "group-mean R²  =  0.450 / 0.790 / 0.978", **MONOS)
para(tf, "SD 0.268. Single-run benchmark numbers are not measurable at this scale.", **BB)
para(tf, "We therefore average three seeds at the prediction level as the "
         "replicate unit, and quote a per-split materiality threshold before "
         "running any ablation.", **B)
fig(s, "fig1_run_to_run_variance.png", 5.9, 1.85, 6.35)
caption(s, 5.9, 6.35, 6.35,
        "fig1_run_to_run_variance · same model, same seed, same code")
foot(s, "No polymer representation paper we are aware of publishes a measured noise floor. Shipping one is itself a contribution.")

# ── S8 piece 6 — the demonstration ───────────────────────────────────────────
s = slide()
piece_title(s, 6, "The demonstration — same predictions, opposite conclusions",
            "B split, cross-scaffold folds: octamer versus wDMPNN")

rows = [("Overall R²", "+0.003", "3 / 5", "p = 1.00", MUTED),
        ("MAE", "−0.000", "3 / 5", "p = 1.00", MUTED),
        ("group-mean R²", "+0.002", "3 / 5", "p = 1.00", MUTED),
        ("ΔR²  (architecture)", "+0.260 EA / +0.152 IP", "5 / 5", "p = 0.062", TERRA)]
card(s, M, 1.8, 6.2, 2.9, TINT)
tf = tb(s, M + 0.3, 2.0, 5.6, 0.3)
para(tf, "MEDIAN PAIRED DIFFERENCE, D FOLDS", size=10, color=MUTED, bold=True,
     first=True, space=10)
yy = 2.44
for lab, med, wins, p, col in rows:
    tf = tb(s, M + 0.3, yy, 2.3, 0.3)
    para(tf, lab, size=12, color=col, bold=(col == TERRA), first=True, space=0)
    tf = tb(s, M + 2.6, yy, 2.0, 0.3)
    para(tf, med, size=12, color=col, bold=(col == TERRA), first=True, space=0)
    tf = tb(s, M + 4.6, yy, 0.7, 0.3)
    para(tf, wins, size=12, color=col, bold=(col == TERRA), first=True, space=0)
    tf = tb(s, M + 5.3, yy, 0.9, 0.3)
    para(tf, p, size=11, color=col, first=True, space=0)
    yy += 0.42

card(s, M, 4.95, 6.2, 1.65, DARK)
tf = tb(s, M + 0.3, 5.15, 5.6, 1.3)
para(tf, "On genuinely new chemistry the two models are statistically "
         "indistinguishable on every conventional measure — and separated by "
         "0.15–0.26 on architecture recovery.", size=13, color=WHITE, bold=True,
     font=HEAD, first=True, space=6)
para(tf, "The metric everyone reports cannot see the effect being studied.",
     size=12, color=SAGE, space=0)

fig(s, "fig_r3_architecture.png", 7.0, 2.5, 5.8)
caption(s, 7.0, 5.15, 5.8,
        "fig_r3_architecture · S and D fold groups, four models. The conventional "
        "metrics for these same runs are in the table on the left.")
foot(s, "[ADD: a single side-by-side figure of these two panels would make the point in one image — worth building for the paper.]", TERRA)

# ── S9 what is missing ───────────────────────────────────────────────────────
s = slide()
tf = tb(s, M, 0.42, W - 2 * M, 0.7)
para(tf, "What the framework still cannot test — and what that motivates",
     size=26, color=INK, bold=True, font=HEAD, first=True, space=0)
tf = tb(s, M, 1.18, W - 2 * M, 0.4)
para(tf, "Two of the three axes have an extrapolation split. The third does not.",
     size=13, color=MUTED, first=True, space=0)

axes = [("Monomer A", "held-out chemistry", "R1 · 9 folds", TEAL, "BUILT"),
        ("Monomer B", "held-out scaffold family", "R3 · 9 folds, S/D grouped", TEAL, "BUILT"),
        ("Architecture", "held-out arrangement", "impossible on this dataset", TERRA, "MISSING")]
x = M
for name, what, how, col, tag in axes:
    card(s, x, 1.9, 3.94, 2.15, TINT if col == TEAL else TINT2)
    tf = tb(s, x + 0.3, 2.12, 3.35, 0.3)
    para(tf, tag, size=10, color=col, bold=True, first=True, space=9)
    tf = tb(s, x + 0.3, 2.5, 3.35, 1.4)
    para(tf, name, size=16, color=INK, bold=True, font=HEAD, first=True, space=5)
    para(tf, what, size=12, color=BODY, space=7)
    para(tf, how, size=11.5, color=col, bold=True, space=0)
    x += 4.14

card(s, M, 4.3, 12.23, 1.5, TINT)
tf = tb(s, M + 0.32, 4.52, 11.6, 1.15)
para(tf, "Why it is impossible today", size=13.5, color=INK, bold=True,
     font=HEAD, first=True, space=6)
para(tf, "There are only three distinct chain-arrangement descriptions in all "
         "42,966 rows — one per architecture type, and identical across "
         "compositions. Hold one out and the model has nothing to interpolate "
         "from; it would be extrapolating to a category it has never seen, with "
         "no intermediate cases in between.", size=12, color=BODY, space=0)

card(s, M, 6.0, 12.23, 0.85, DARK)
tf = tb(s, M + 0.32, 6.18, 11.6, 0.55)
para(tf, "A continuous blockiness axis would make the third split possible — "
         "and that is the argument for the dataset, independent of any model.",
     size=13.5, color=WHITE, bold=True, font=HEAD, first=True, space=0)

# ── S10 framing + questions ──────────────────────────────────────────────────
s = slide(dark=True)
tf = tb(s, M, 0.6, W - 2 * M, 0.7)
para(tf, "Proposed framing, and what I would like your view on", size=27,
     color=WHITE, bold=True, font=HEAD, first=True, space=0)

card(s, M, 1.6, 6.05, 2.5, C(0x3A, 0x4A, 0x53))
tf = tb(s, M + 0.32, 1.82, 5.4, 2.1)
para(tf, "THE CLAIM", size=10, color=TEAL, bold=True, first=True, space=10)
para(tf, "Not \"here are better metrics\" — but \"current evaluation cannot "
         "answer the question the field is asking, and here is the evidence.\"",
     size=13, color=WHITE, bold=True, font=HEAD, space=8)
para(tf, "The metrics are the remedy, not the contribution.", size=12, color=SAGE, space=0)

card(s, 6.95, 1.6, 5.83, 2.5, C(0x3A, 0x4A, 0x53))
tf = tb(s, 7.27, 1.82, 5.2, 2.1)
para(tf, "WHY NOW", size=10, color=TEAL, bold=True, first=True, space=10)
para(tf, "All six components are built and measured. No new compute is needed "
         "to write it.", size=13, color=WHITE, space=8)
para(tf, "It also does not depend on the octamer being the best model — the "
         "framework stands whichever model wins, which insulates it from the "
         "protocol-matching confound.", size=12, color=SAGE, space=0)

card(s, M, 4.35, 12.23, 2.3, C(0x3A, 0x4A, 0x53))
tf = tb(s, M + 0.32, 4.58, 11.6, 0.3)
para(tf, "QUESTIONS FOR YOU", size=10, color=TEAL, bold=True, first=True, space=12)
tf = tb(s, M + 0.32, 4.98, 11.55, 1.6)
for i, q in enumerate([
    "Is \"the evaluation cannot see the property\" a strong enough central claim, "
    "or does it need a model result attached to be publishable?",
    "Which of the six carries the most weight? My view is the null floor — "
    "a parameter-free predictor beating a trained model is the most striking single result.",
    "Venue — a methods/benchmark venue, or a chemistry venue where the audience "
    "is the one making the mistake?",
    "Should the architecture-held-out gap be stated in this paper as motivation "
    "for the next one, or held back?"]):
    para(tf, f"{i+1}.   {q}", size=12.5, color=WHITE, first=(i == 0), space=11)


# ── S11 paper strategy ───────────────────────────────────────────────────────
s = slide()
tf = tb(s, M, 0.42, W - 2 * M, 0.7)
para(tf, "Where this sits — four papers, and why this one goes first",
     size=26, color=INK, bold=True, font=HEAD, first=True, space=0)
tf = tb(s, M, 1.18, W - 2 * M, 0.4)
para(tf, "The through-line comes from the accepted review's own Future "
         "Perspectives section", size=13, color=MUTED, first=True, space=0)

papers = [
    ("0", "Review — accepted, Digital Discovery", TEAL, "DONE",
     "Argues in the literature that representation quality needs evaluating "
     "beyond predictive accuracy. Contains a section titled exactly that."),
    ("1", "Evaluation framework  —  this deck", TERRA, "WRITE NOW",
     "Demonstrates that argument empirically on the field's standard copolymer "
     "benchmark. All six components built; no new compute needed."),
    ("2", "Architecture-aware benchmark dataset", AMBER, "COSTED",
     "Enables the architecture-held-out split that Paper 1 shows is missing. "
     "~8–9 kSU CPU, pipeline verified portable."),
    ("3", "The model — octamer versus baseline", AMBER, "EVIDENCE IN HAND",
     "Significant on two comparisons at p = 0.004. Held until external validity "
     "is tested, so it does not carry the protocol-matching confound alone."),
]
y = 1.85
for n, h, col, tag, b in papers:
    card(s, M, y, 12.23, 1.02, TINT if n in "13" else TINT2)
    numdot(s, M + 0.3, y + 0.26, 0.42, n, col)
    tf = tb(s, M + 0.92, y + 0.14, 5.1, 0.8)
    para(tf, h, size=13.5, color=INK, bold=True, font=HEAD, first=True, space=4)
    para(tf, tag, size=10, color=col, bold=True, space=0)
    tf = tb(s, 6.4, y + 0.22, 6.3, 0.75)
    para(tf, b, size=11, color=BODY, first=True, space=0)
    y += 1.12

card(s, M, 6.4, 12.23, 0.75, DARK)
tf = tb(s, M + 0.32, 6.58, 11.6, 0.5)
para(tf, "Year 1 argues it in the literature. Year 2 demonstrates it on the "
         "benchmark everyone uses. Year 3 builds the dataset the argument requires.",
     size=13.5, color=WHITE, bold=True, font=HEAD, first=True, space=0)

# ── S12 why this order ───────────────────────────────────────────────────────
s = slide()
tf = tb(s, M, 0.42, W - 2 * M, 0.7)
para(tf, "Why Paper 1 first — three reasons, one of them defensive",
     size=26, color=INK, bold=True, font=HEAD, first=True, space=0)
tf = tb(s, M, 1.18, W - 2 * M, 0.4)
para(tf, "The alternative was to lead with the model result", size=13,
     color=MUTED, first=True, space=0)

reasons = [
    ("It is writable today", TEAL,
     "Every component is built and measured. The model paper depends on "
     "external-validity tests that have not been run; the dataset paper "
     "depends on compute not yet approved."),
    ("It does not depend on which model wins", TEAL,
     "The framework stands whether the octamer, HPG-hier or something later "
     "turns out best. That insulates it entirely from the protocol-matching "
     "confound, which is the largest threat to the model paper."),
    ("The measurement work invalidates the modelling work", TERRA,
     "The noise floor, the checkpoint bug, the null floors and the split "
     "design collectively void most of the earlier model comparisons. The "
     "stronger contribution is the one that did the invalidating."),
]
x = M
for h, col, b in reasons:
    card(s, x, 1.85, 3.94, 2.5, TINT)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, In(x + 0.3), In(2.1), In(0.18), In(0.18))
    dot.fill.solid(); dot.fill.fore_color.rgb = col
    dot.line.fill.background(); dot.shadow.inherit = False
    tf = tb(s, x + 0.3, 2.48, 3.35, 1.7)
    para(tf, h, size=14, color=INK, bold=True, font=HEAD, first=True, space=8)
    para(tf, b, size=11.5, color=BODY, space=0)
    x += 4.14

card(s, M, 4.6, 12.23, 1.35, TINT2)
tf = tb(s, M + 0.32, 4.82, 11.6, 1.0)
para(tf, "What Paper 1 deliberately does not claim", size=13.5, color=TERRA,
     bold=True, font=HEAD, first=True, space=7)
para(tf, "That our model is the best. The octamer appears in the paper only as "
         "one of five models used to demonstrate that two metrics can disagree "
         "on the same predictions. Its own significance result belongs to Paper 3, "
         "after external validity has been tested.", size=12, color=BODY, space=0)

card(s, M, 6.2, 12.23, 0.9, DARK)
tf = tb(s, M + 0.32, 6.4, 11.6, 0.55)
para(tf, "Target: submit in the first quarter of Year 3. No experiments block it.",
     size=13.5, color=WHITE, bold=True, font=HEAD, first=True, space=0)

prs.save("eval_framework_deck.pptx")
print("saved eval_framework_deck.pptx")
