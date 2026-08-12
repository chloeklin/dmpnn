#!/usr/bin/env python3
"""Build the research-update deck: Panel Chair (Mon) + Main Supervisor (Wed).

All figures are drawn from verified repository artefacts. Placeholders are marked
[ADD: ...] and are the only content requiring input before presenting.
"""
from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── palette ──────────────────────────────────────────────────────────────────
DARK   = C(0x2B, 0x3A, 0x42)   # charcoal-slate
INK    = C(0x1A, 0x24, 0x2A)
BODY   = C(0x3D, 0x4C, 0x54)
MUTED  = C(0x77, 0x86, 0x8E)
WHITE  = C(0xFF, 0xFF, 0xFF)
TEAL   = C(0x02, 0x80, 0x90)   # established / positive
AMBER  = C(0xC4, 0x8B, 0x0E)   # in progress
TERRA  = C(0xB8, 0x50, 0x42)   # decision needed / ruled out
SAGE   = C(0xA7, 0xBE, 0xAE)
TINT   = C(0xF2, 0xF5, 0xF6)   # card fill
TINT2  = C(0xE8, 0xEF, 0xF1)

HEAD = "Cambria"
BODYF = "Calibri"

W, H = 13.333, 7.5
M = 0.55

prs = Presentation()
prs.slide_width, prs.slide_height = In(W), In(H)
BLANK = prs.slide_layouts[6]


# ── helpers ──────────────────────────────────────────────────────────────────
def slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    if dark:
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, In(W), In(H))
        bg.fill.solid(); bg.fill.fore_color.rgb = DARK
        bg.line.fill.background(); bg.shadow.inherit = False
    return s


def tb(s, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return tf


def para(tf, text, size=14, color=BODY, bold=False, font=BODYF, space=6,
         first=False, align=None, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = text
    p.space_after = Pt(space)
    if align: p.alignment = align
    for r in p.runs:
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = font; r.font.italic = italic
    return p


def title(s, text, sub=None, dark=False):
    tf = tb(s, M, 0.42, W - 2 * M, 0.85)
    para(tf, text, size=27, color=WHITE if dark else INK, bold=True, font=HEAD,
         space=3, first=True)
    if sub:
        t2 = tb(s, M, 1.28, W - 2 * M, 0.4)
        para(t2, sub, size=13.5, color=SAGE if dark else MUTED, first=True, space=0)


def card(s, x, y, w, h, fill=TINT, line=None):
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, In(x), In(y), In(w), In(h))
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if line:
        r.line.color.rgb = line; r.line.width = Pt(1.25)
    else:
        r.line.fill.background()
    r.shadow.inherit = False
    r.adjustments[0] = 0.045
    return r


def dot(s, x, y, d, color):
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, In(x), In(y), In(d), In(d))
    o.fill.solid(); o.fill.fore_color.rgb = color
    o.line.fill.background(); o.shadow.inherit = False
    return o


def numdot(s, x, y, d, n, fill, txt=WHITE):
    o = dot(s, x, y, d, fill)
    tf = o.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = str(n); p.alignment = PP_ALIGN.CENTER
    for r in p.runs:
        r.font.size = Pt(13); r.font.bold = True
        r.font.color.rgb = txt; r.font.name = BODYF
    return o


def table(s, x, y, w, rows, colw, header=True, fs=11.5, rowh=0.34, hfs=11):
    nr, nc = len(rows), len(rows[0])
    shp = s.shapes.add_table(nr, nc, In(x), In(y), In(w), In(rowh * nr))
    t = shp.table
    for i, cw in enumerate(colw):
        t.columns[i].width = In(cw)
    for ri, row in enumerate(rows):
        t.rows[ri].height = In(rowh)
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.margin_left = In(0.09); cell.margin_right = In(0.06)
            cell.margin_top = In(0.03); cell.margin_bottom = In(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if header and ri == 0:
                cell.fill.fore_color.rgb = DARK
            else:
                cell.fill.fore_color.rgb = WHITE if ri % 2 else TINT
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            txt, col, bold = val, BODY, False
            if isinstance(val, tuple):
                txt, col, bold = val
            p.text = str(txt)
            for r in p.runs:
                r.font.size = Pt(hfs if (header and ri == 0) else fs)
                r.font.name = BODYF
                r.font.color.rgb = WHITE if (header and ri == 0) else col
                r.font.bold = True if (header and ri == 0) else bold
    return t


def stat(s, x, y, w, big, label, color=TEAL, bigsize=34):
    tf = tb(s, x, y, w, 0.62)
    para(tf, big, size=bigsize, color=color, bold=True, font=HEAD, space=0, first=True)
    t2 = tb(s, x, y + 0.56, w, 0.55)
    para(t2, label, size=10.5, color=MUTED, first=True, space=0)


def foot(s, text, color=MUTED):
    tf = tb(s, M, H - 0.52, W - 2 * M, 0.32)
    para(tf, text, size=9.5, color=color, first=True, space=0, italic=True)


def divider(s, kicker, text, sub):
    tf = tb(s, M, 2.5, W - 2 * M, 0.4)
    para(tf, kicker, size=13, color=TEAL, bold=True, first=True, space=0)
    t2 = tb(s, M, 3.0, W - 2.4 * M, 1.3)
    para(t2, text, size=36, color=WHITE, bold=True, font=HEAD, first=True, space=0)
    t3 = tb(s, M, 4.42, W - 3.5 * M, 0.6)
    para(t3, sub, size=14, color=SAGE, first=True, space=0)


# ═════════════════════════ SECTION 1 — PANEL CHAIR ═══════════════════════════

# S1 title
s = slide(dark=True)
tf = tb(s, M, 2.15, W - 2 * M, 0.4)
para(tf, "RESEARCH UPDATE  ·  10 AUGUST 2026", size=12, color=TEAL, bold=True,
     first=True, space=0)
tf = tb(s, M, 2.68, 10.6, 1.9)
para(tf, "Architecture-aware copolymer\nrepresentation learning", size=34,
     color=WHITE, bold=True, font=HEAD, first=True, space=0)
tf = tb(s, M, 4.85, 10.2, 1.0)
para(tf, "Two questions this month: can we reuse the EA/IP generation method, "
         "and why does the octamer model beat the published baseline?",
     size=14.5, color=SAGE, first=True, space=0)
tf = tb(s, M, 6.35, 11, 0.4)
para(tf, "Section 1 — Panel Chair (today)   ·   Section 2 — Main Supervisor (Wednesday)",
     size=11, color=MUTED, first=True, space=0)
s.notes_slide.notes_text_frame.text = (
    "Opening line: I have a model that beats the published state of the art, the "
    "improvement is statistically significant which is rare here, but I've found a "
    "reason to be cautious about WHY it wins — and most of this month went into "
    "testing that rather than chasing more accuracy.")

# S2 — the two questions
s = slide()
title(s, "Two questions — one answered, one now more interesting",
      "Where the month went")

card(s, M, 1.85, 6.05, 1.95, TINT2)
dot(s, M + 0.32, 2.12, 0.2, TEAL)
tf = tb(s, M + 0.62, 2.06, 5.3, 0.35)
para(tf, "ANSWERED", size=10.5, color=TEAL, bold=True, first=True, space=0)
tf = tb(s, M + 0.32, 2.5, 5.5, 1.5)
para(tf, "Can we reuse the EA/IP method to generate\nour own dataset?",
     size=15.5, color=INK, bold=True, font=HEAD, first=True, space=8)
para(tf, "Yes. The authors' generation code is published, it runs on the "
         "monomers we already have, and I have it costed.", size=12.5, color=BODY, space=0)

card(s, 6.95, 1.85, 5.83, 1.95, TINT2)
dot(s, 7.27, 2.12, 0.2, TERRA)
tf = tb(s, 7.57, 2.06, 5.0, 0.35)
para(tf, "NOW THE REAL QUESTION", size=10.5, color=TERRA, bold=True, first=True, space=0)
tf = tb(s, 7.27, 2.5, 5.25, 1.5)
para(tf, "What dataset should we actually\ngenerate?", size=15.5, color=INK,
     bold=True, font=HEAD, first=True, space=8)
para(tf, "Regenerating the same thing, bigger, adds nothing. I have a proposal "
         "and a demonstrated gap it fills.", size=12.5, color=BODY, space=0)

card(s, M, 4.15, 12.23, 1.6, TINT)
tf = tb(s, M + 0.32, 4.38, 11.6, 0.4)
para(tf, "Running in parallel: why does the octamer model beat the published baseline?",
     size=14.5, color=INK, bold=True, font=HEAD, first=True, space=8)
para(tf, "Five candidate explanations. One ruled out by experiment, one under test now, "
         "three untested. Plus a confound I found in my own work and want to raise "
         "before anyone else does.", size=12.5, color=BODY, space=0)
foot(s, "Detail on the octamer investigation from slide 7.")

# S3 — yes, reusable
s = slide()
title(s, "Yes — the EA/IP generation pipeline can be reused",
      "Three things checked, all clear")

items = [
    ("Generation code is published",
     "The authors released data and scripts. All eight steps of the documented "
     "protocol have code; the DFT calibration constants are in the repository, "
     "not just cited."),
    ("Our monomers are already compatible",
     "The pipeline encodes a real reaction — Suzuki coupling — so it only joins "
     "boronic acids to bromides. Checked: 9 of 9 monomer A are boronic acids, "
     "682 of 682 monomer B are bromides. No chemistry rewrite needed."),
    ("Cost is bounded and affordable",
     "≤ 256 structures per polymer. For ~2,000 polymers that is a ceiling of "
     "≈ 8–9 kSU of CPU — about 5% of what full replication would cost."),
]
y = 1.92
for i, (h, b) in enumerate(items, 1):
    card(s, M, y, 8.35, 1.5, TINT)
    numdot(s, M + 0.3, y + 0.32, 0.42, i, TEAL)
    tf = tb(s, M + 0.92, y + 0.24, 7.2, 0.4)
    para(tf, h, size=14.5, color=INK, bold=True, font=HEAD, first=True, space=5)
    para(tf, b, size=11.5, color=BODY, space=0)
    y += 1.66

card(s, 9.2, 1.92, 3.58, 4.66, TINT2)
tf = tb(s, 9.5, 2.18, 3.0, 0.35)
para(tf, "WHAT IT WOULD NEED", size=10, color=TERRA, bold=True, first=True, space=10)
for t in ["Pin the quantum-chemistry version — the repo does not record it",
          "Add random seeding — the original has none, so its own sequences "
          "aren't reproducible",
          "Fix one broken committed command and a notebook needing manual edits",
          "Job script requests 8 CPUs but uses 1 — an 8× overcharge if copied"]:
    para(tf, "• " + t, size=11, color=BODY, space=9)
para(tf, "A few days of careful work — not a rewrite.", size=11.5, color=TEAL,
     bold=True, space=0)
foot(s, "Sources: paper SI p. S4; coleygroup/polymer-chemprop-data; monomer check run against data/ea_ip.csv.")

# S4 — how the labels were made
s = slide()
title(s, "How the existing labels were computed — and what that costs",
      "The method we would be reusing")

steps = [("Build the chains",
          "Join the two monomers into an 8-unit chain matching the requested "
          "composition and arrangement. Up to 32 different chains per polymer."),
         ("Make 3D shapes",
          "8 conformations per chain — up to 256 structures per polymer."),
         ("Optimise geometry",
          "Rough force-field pass, then quantum chemistry (GFN2-xTB)."),
         ("Compute properties",
          "Electron affinity and ionisation potential for every structure."),
         ("Average",
          "Over the 8 shapes, then over all chains. One number per polymer.")]
y = 1.9
for i, (h, b) in enumerate(steps, 1):
    numdot(s, M + 0.05, y + 0.02, 0.38, i, DARK)
    tf = tb(s, M + 0.62, y - 0.02, 7.3, 0.75)
    para(tf, h, size=13.5, color=INK, bold=True, font=HEAD, first=True, space=3)
    para(tf, b, size=11.5, color=BODY, space=0)
    y += 0.94

card(s, 8.6, 1.9, 4.18, 2.62, TINT2)
tf = tb(s, 8.9, 2.12, 3.6, 0.35)
para(tf, "COST — BOUNDED, NOT MEASURED", size=10, color=MUTED, bold=True,
     first=True, space=10)
stat(s, 8.9, 2.62, 3.6, "≤ 256", "structures per polymer (worst case)")
stat(s, 8.9, 3.42, 3.6, "≈ 8–9 kSU", "CPU for ~2,000 polymers — a ceiling", TEAL)

card(s, 8.6, 4.72, 4.18, 1.86, TINT)
tf = tb(s, 8.9, 4.95, 3.6, 1.5)
para(tf, "Two things to flag", size=12.5, color=INK, bold=True, font=HEAD,
     first=True, space=7)
para(tf, "This is CPU, not GPU — a different queue from model training.",
     size=11, color=BODY, space=7)
para(tf, "Full replication of the original 42,966 would be ≈ 150 kSU. "
         "We would be proposing about 5% of that.", size=11, color=BODY, space=0)
foot(s, "Steps 3–4 dominate the cost; everything before is fast. A one-day pilot would replace the ceiling with a measured figure.")

# S5 — the gap
s = slide()
title(s, "The harder question: the existing benchmark has only three "
         "architecture settings",
      "Why more of the same data would not help")

card(s, M, 1.9, 6.2, 2.6, TINT)
tf = tb(s, M + 0.3, 2.12, 5.6, 0.4)
para(tf, "What we can measure today", size=14.5, color=INK, bold=True,
     font=HEAD, first=True, space=8)
table(s, M + 0.3, 2.66, 5.6,
      [["Composition", "Architectures present", "Task"],
       ["fracA = 0.25", "block, random", "2-way"],
       ["fracA = 0.50", "block, random, alternating", "3-way"],
       ["fracA = 0.75", "block, random", "2-way"]],
      [1.5, 2.8, 1.3], fs=10.5, hfs=10, rowh=0.36)

card(s, 6.95, 1.9, 5.83, 2.6, TINT2)
tf = tb(s, 7.25, 2.12, 5.25, 2.2)
para(tf, "Verified against all 42,966 rows", size=14.5, color=INK, bold=True,
     font=HEAD, first=True, space=8)
para(tf, "There are only three distinct chain-arrangement descriptions in the "
         "entire dataset — one per architecture type. They do not vary with "
         "composition.", size=12, color=BODY, space=8)
para(tf, "So a model only ever has to tell two or three discrete cases apart. "
         "That is the whole architecture task in the only benchmark that has one.",
     size=12, color=BODY, space=0)

card(s, M, 4.75, 12.23, 1.85, TINT)
dot(s, M + 0.32, 5.06, 0.2, TERRA)
tf = tb(s, M + 0.64, 5.0, 11.4, 1.4)
para(tf, "The gap, stated as evidence rather than opinion", size=14, color=INK,
     bold=True, font=HEAD, first=True, space=7)
para(tf, "The only architecture-varying polymer benchmark offers a three-way "
         "classification. The obvious alternative — a 5,400-row block copolymer "
         "phase database — is 99.7% diblock across 61 chemistry pairs, so it has "
         "no architecture axis at all. Architecture-aware representation learning "
         "currently has nowhere to be evaluated properly.",
     size=12, color=BODY, space=0)
foot(s, "Block copolymer database inspected 8 Aug 2026: 5,400 rows, 5,382 diblock, 18-class phase label driven by molecular weight and temperature.")

# S6 — the proposal
s = slide()
title(s, "Proposal: a continuous architecture axis, not three labels",
      "Same chemistry, same composition — only the arrangement varies")

card(s, M, 1.88, 7.55, 2.55, TINT)
tf = tb(s, M + 0.3, 2.08, 6.95, 0.4)
para(tf, "The core idea", size=14, color=INK, bold=True, font=HEAD, first=True, space=7)
para(tf, "Hold composition at 50/50 — four A units and four B units, always — "
         "and vary only where they sit along the chain:", size=11.5, color=BODY, space=8)
for lbl, seq in [("perfectly alternating", "A B A B A B A B"),
                 ("mostly mixed", "A B A A B A B B"),
                 ("random", "A A B A B B A B"),
                 ("fully blocky", "A A A A B B B B")]:
    p = para(tf, f"{seq}        {lbl}", size=11.5, color=BODY, space=4,
             font="Courier New")
para(tf, "Three points on that line today. The proposal makes it continuous.",
     size=11.5, color=TEAL, bold=True, space=0)

card(s, 8.3, 1.88, 4.48, 2.55, TINT2)
tf = tb(s, 8.6, 2.08, 3.9, 2.2)
para(tf, "Also worth shipping", size=14, color=INK, bold=True, font=HEAD,
     first=True, space=7)
para(tf, "Un-averaged labels — one value per chain, not only the average. "
         "The calculation already produces them; the original discards them.",
     size=11.5, color=BODY, space=7)
para(tf, "Costs nothing extra, is strictly more information than any existing "
         "benchmark, and lets others test the averaging question themselves.",
     size=11.5, color=BODY, space=0)

y = 4.7
hdrs = [("ESTABLISHED", TEAL,
         "Pipeline is reusable · our monomers fit · cost bounded at ≈8–9 kSU · "
         "the three-setting limitation is measured, not asserted"),
        ("I PROPOSE", AMBER,
         "Continuous blockiness axis at fixed chemistry and composition · publish "
         "un-averaged labels · ship the evaluation protocol and measured noise floor"),
        ("NEEDS A DECISION", TERRA,
         "Is this the right dataset? How many polymers and which monomers? "
         "Does it belong to this paper or the next one?")]
x = M
for h, col, b in hdrs:
    card(s, x, y, 3.94, 1.9, TINT)
    dot(s, x + 0.3, y + 0.28, 0.18, col)
    tf = tb(s, x + 0.58, y + 0.22, 3.1, 0.3)
    para(tf, h, size=10, color=col, bold=True, first=True, space=0)
    tf = tb(s, x + 0.3, y + 0.66, 3.36, 1.1)
    para(tf, b, size=11, color=BODY, first=True, space=0)
    x += 4.14
foot(s, "This slide is also the basis for Wednesday's discussion — see Section 2.")

# S7 — octamer result
s = slide()
title(s, "The octamer model beats the published baseline — significantly",
      "Chemistry-extrapolation split, 9 folds, three seeds averaged")

table(s, M, 1.95, 7.6,
      [["Model", "Overall R²  EA / IP", "MAE eV  EA / IP", "Architecture recovery  EA / IP"],
       ["wDMPNN (published config)", "0.967 / 0.971", "0.070 / 0.050", "0.397 / 0.565"],
       ["HPG-hier", "0.966 / 0.890", "0.067 / 0.068", "0.776 / 0.808"],
       [("HPG-octamer", INK, True), ("0.984 / 0.978", INK, True),
        ("0.055 / 0.035", INK, True), ("0.849 / 0.886", TEAL, True)]],
      [2.6, 1.75, 1.55, 1.7], fs=11, hfs=10, rowh=0.44)

card(s, M, 3.95, 7.6, 1.35, TINT2)
tf = tb(s, M + 0.3, 4.15, 7.0, 1.0)
para(tf, "The obvious objection is closed off", size=13, color=INK, bold=True,
     font=HEAD, first=True, space=6)
para(tf, "I re-ran the baseline at the exact settings from the original paper. "
         "It got better — noticeably better at the chemistry part — and the "
         "octamer still won on architecture. \"The baseline was under-tuned\" "
         "is no longer available.", size=11.5, color=BODY, space=0)

card(s, 8.4, 1.95, 4.38, 3.35, TINT)
tf = tb(s, 8.7, 2.15, 3.8, 0.35)
para(tf, "PAIRED PER-FOLD TESTS", size=10, color=MUTED, bold=True, first=True, space=10)
stat(s, 8.7, 2.62, 3.8, "9 / 9 folds", "EA overall R², p = 0.004", TEAL)
stat(s, 8.7, 3.44, 3.8, "9 / 9 folds", "IP architecture recovery, p = 0.004", TEAL)
tf = tb(s, 8.7, 4.35, 3.8, 0.8)
para(tf, "Four further comparisons at p = 0.039. p = 0.004 is the best achievable "
         "with nine folds.", size=11, color=BODY, first=True, space=0)

card(s, M, 5.5, 12.23, 1.1, TINT)
dot(s, M + 0.32, 5.79, 0.18, TERRA)
tf = tb(s, M + 0.62, 5.72, 11.4, 0.8)
para(tf, "On the hardest split the two models are statistically indistinguishable on "
         "accuracy (3 of 5 folds, p = 1.0) yet separated by 0.15–0.26 on architecture "
         "recovery (5 of 5 folds). Same predictions, two measures pointing opposite ways.",
     size=12, color=BODY, first=True, space=0)
foot(s, "Significance is rare in this literature. Comparison family to be fixed in advance before any p-value is quoted in writing.")

# S8 — hypothesis status table
s = slide()
title(s, "Why does it win? Five candidate explanations, one eliminated",
      "The octamer differs from the simpler model in five ways at once")

rows = [["#", "Candidate explanation", "Status", "Evidence", "What it implies"],
        ["5", "Averaging over 16 sampled chains",
         ("RULED OUT", TERRA, True),
         "Pre-registered K=1 ablation, 53 cells. No material change; seed-SD sign tests 5/9, 3/8, 8/17",
         "The advantage is not from approximating the label-averaging step"],
        ["2", "Learned position information",
         ("TESTING NOW", AMBER, True),
         "Pilot, 2 cells: architecture recovery fell 0.173 and 0.069 vs threshold 0.051",
         "Both exceed threshold, same direction — leading candidate"],
        ["3", "How chain information is pooled",
         ("NOT TESTED", MUTED, True), "—",
         "Effect likely inside the noise floor"],
        ["1", "The 8-unit chain structure itself",
         ("NOT TESTED", MUTED, True), "—",
         "Also the factor implicated in the confound (next slide)"],
        ["4", "Discards junction chemistry the simpler model uses",
         ("PLAUSIBLE, UNTESTED", MUTED, True), "—",
         "Wins while seeing less information"]]
table(s, M, 1.9, 12.23, rows, [0.42, 2.85, 1.62, 3.75, 3.59], fs=10, hfs=9.8, rowh=0.60)

tf = tb(s, M, 5.85, 12.23, 0.6)
para(tf, "A limit worth stating: the remaining differences are small enough to sit inside "
         "this dataset's run-to-run noise. Beyond a point, no further ablation on this data "
         "can separate them — a property of the benchmark, not of effort.",
     size=11, color=BODY, first=True, space=0)
foot(s, "[ADD: full 54-cell positional-embedding result once the arm completes — pilot is 2 cells, seed 42 only.]", TERRA)

# S9 — the confound
s = slide()
title(s, "A confound I found in my own work",
      "Raising it before anyone else does")

card(s, M, 1.9, 6.05, 2.15, TINT2)
tf = tb(s, M + 0.3, 2.1, 5.45, 1.8)
para(tf, "How the labels were made", size=13.5, color=INK, bold=True,
     font=HEAD, first=True, space=8)
for t in ["Computed on 8-unit chains",
          "Averaged over up to 32 sampled arrangements",
          "One number per polymer"]:
    para(tf, "• " + t, size=12, color=BODY, space=5)

card(s, 6.95, 1.9, 5.83, 2.15, TINT2)
tf = tb(s, 7.25, 2.1, 5.25, 1.8)
para(tf, "How my best model works", size=13.5, color=INK, bold=True,
     font=HEAD, first=True, space=8)
for t in ["Uses an 8-slot chain",
          "Samples 16 arrangements",
          "Averages its predictions"]:
    para(tf, "• " + t, size=12, color=BODY, space=5)

card(s, M, 4.3, 12.23, 0.75, DARK)
tf = tb(s, M + 0.32, 4.48, 11.5, 0.45)
para(tf, "My model's structure mirrors how the labels were generated.",
     size=15, color=WHITE, bold=True, font=HEAD, first=True, space=0)

y = 5.3
for h, col, b in [("Favourable reading", TEAL,
                   "I encoded the right physics. The property genuinely is an average "
                   "over chain arrangements, so a model built that way should win. "
                   "Matching a model to the data-generating process is good science."),
                  ("Sceptical reading", TERRA,
                   "The advantage may be \"this matches the dataset's recipe\" rather "
                   "than \"this is a better polymer representation.\" Real polymers are "
                   "hundreds of units long; measured properties are not 8-unit averages.")]:
    x = M if col == TEAL else 6.95
    card(s, x, y, 5.83 if col != TEAL else 6.05, 1.32, TINT)
    dot(s, x + 0.3, y + 0.26, 0.18, col)
    tf = tb(s, x + 0.58, y + 0.19, 4.9, 0.3)
    para(tf, h, size=11.5, color=col, bold=True, first=True, space=0)
    tf = tb(s, x + 0.3, y + 0.58, (5.45 if col == TEAL else 5.25), 0.7)
    para(tf, b, size=11, color=BODY, first=True, space=0)
foot(s, "These cannot be distinguished on this dataset by any ablation. It is the single biggest threat to the claim.")

# S10 — status
s = slide()
title(s, "Status: on track, with the milestone document as the priority",
      "Twelve-month view")

card(s, M, 1.9, 6.05, 2.7, TINT)
tf = tb(s, M + 0.3, 2.12, 5.45, 2.3)
para(tf, "Delivered this year", size=14, color=INK, bold=True, font=HEAD,
     first=True, space=8)
for t in ["Review paper accepted, Digital Discovery",
          "Two-axis evaluation framework with a null-model floor",
          "~590 training runs under a frozen, pre-registered protocol",
          "A measured run-to-run noise floor — not published elsewhere in this literature",
          "A model that significantly beats the published baseline at that baseline's own settings"]:
    para(tf, "• " + t, size=11.5, color=BODY, space=6)

card(s, 6.95, 1.9, 5.83, 2.7, TINT)
tf = tb(s, 7.25, 2.12, 5.25, 2.3)
para(tf, "Found and fixed this year", size=14, color=INK, bold=True, font=HEAD,
     first=True, space=8)
for t in ["A model-selection bug affecting every experiment",
          "An input-parsing bug that silently disabled the baseline's polymer features",
          "A mis-specified quality filter that favoured our own models",
          "A performance defect that made all timing comparisons wrong"]:
    para(tf, "• " + t, size=11.5, color=BODY, space=6)
para(tf, "All documented, corrected and disclosed.", size=11.5, color=TEAL,
     bold=True, space=0)

card(s, M, 4.85, 12.23, 1.75, TINT2)
tf = tb(s, M + 0.32, 5.08, 11.5, 0.35)
para(tf, "Next", size=14, color=INK, bold=True, font=HEAD, first=True, space=8)
x = M + 0.32
for n, t in [("Now → 29 Aug", "Milestone document. Writing to a fixed outline; no new "
              "experiments for it."),
             ("Sep – Nov", "Complete the octamer diagnostics; cheap external-validity tests."),
             ("Dec – Feb", "Submit the measurement paper; decide on dataset generation.")]:
    tf2 = tb(s, x, 5.5, 3.7, 0.95)
    para(tf2, n, size=11, color=TEAL, bold=True, first=True, space=4)
    para(tf2, t, size=11, color=BODY, space=0)
    x += 3.9
foot(s, "Nothing needed today. One decision is pending with my main supervisor — whether to build the dataset — and it is costed and designed.")


# ═══════════════════════ SECTION 2 — MAIN SUPERVISOR ═════════════════════════

# S11 divider
s = slide(dark=True)
divider(s, "SECTION 2  ·  WEDNESDAY",
        "Discussion with main supervisor",
        "Same two topics, more technical detail — and the points where I need your "
        "decision rather than your agreement.")
tf = tb(s, M, 5.5, 11.5, 1.0)
para(tf, "1.  What dataset should we generate — and should we commit compute to it?",
     size=13.5, color=WHITE, first=True, space=7)
para(tf, "2.  Is my approach to diagnosing the octamer advantage the right one?",
     size=13.5, color=WHITE, space=0)

# S12 — dataset design detail
s = slide()
title(s, "Dataset design: four principles, and what I would fix",
      "For discussion — I have a proposal, not a decision")

principles = [
    ("Architecture varies at fixed chemistry",
     "The property under test must be isolable. This is what no other dataset has, "
     "and what makes our within-group metric computable at all."),
    ("No representation is privileged by the label protocol",
     "State plainly how labels are generated and what inductive bias that favours. "
     "Our current confound is exactly this problem."),
    ("Publish the evaluation, not just the numbers",
     "The two-axis metrics, null-floor predictor, scaffold-aware fold design — all "
     "reusable and all currently ours alone."),
    ("Publish the measured noise floor",
     "Three identical runs gave R² of 0.45, 0.79 and 0.98. Nobody reports this. It "
     "would stop others quoting single-run numbers."),
]
y = 1.9
for i, (h, b) in enumerate(principles, 1):
    card(s, M, y, 7.4, 1.12, TINT)
    numdot(s, M + 0.28, y + 0.3, 0.36, i, DARK)
    tf = tb(s, M + 0.82, y + 0.2, 6.3, 0.75)
    para(tf, h, size=12.5, color=INK, bold=True, font=HEAD, first=True, space=4)
    para(tf, b, size=10.5, color=BODY, space=0)
    y += 1.22

card(s, 8.25, 1.9, 4.53, 4.44, TINT2)
tf = tb(s, 8.55, 2.12, 3.95, 4.0)
para(tf, "Scoping questions I have not settled", size=13.5, color=TERRA,
     bold=True, font=HEAD, first=True, space=10)
for t in ["How many polymers? I costed ~2,000 as a ceiling of 8–9 kSU, but the "
          "number is a design choice, not a constraint.",
          "Existing monomers or new ones? My inclination is existing — the "
          "contribution is the architecture axis, not chemical diversity. New "
          "monomers would add a held-out chemistry test set.",
          "How many blockiness levels, and spaced how? Uniform in transition "
          "probability, or uniform in a physical measure of run length?",
          "One property or two? A second property would stop it being a "
          "single-task leaderboard."]:
    para(tf, "• " + t, size=11, color=BODY, space=9)
foot(s, "[ADD: any constraint from your side on total compute available for CPU generation this year.]", TERRA)

# S13 — decisions needed, dataset
s = slide()
title(s, "Decisions I need from you — dataset", "Before I commit compute or time")

qs = [("Is this the right dataset to generate at all?",
       "The alternative is to finish the current paper with a stated single-dataset "
       "limitation and not build anything. That is a legitimate choice and it is "
       "cheaper."),
      ("Is the scope right?",
       "~2,000 polymers, existing monomers, continuous blockiness at fixed "
       "composition, un-averaged labels published alongside averages."),
      ("What am I missing?",
       "I have looked at architecture, composition and chemistry. I have not "
       "considered chain length, dispersity, or tacticity — are any of those "
       "more important than a finer architecture axis?"),
      ("Does this belong to this paper or the next one?",
       "My read is that it is a separate contribution on a separate timeline, and "
       "that the current paper should ship with the limitation stated. I would "
       "like to test that view.")]
y = 1.9
for i, (q, b) in enumerate(qs, 1):
    card(s, M, y, 12.23, 1.12, TINT if i % 2 else TINT2)
    numdot(s, M + 0.3, y + 0.3, 0.4, i, TERRA)
    tf = tb(s, M + 0.9, y + 0.2, 11.1, 0.8)
    para(tf, q, size=14, color=INK, bold=True, font=HEAD, first=True, space=5)
    para(tf, b, size=11.5, color=BODY, space=0)
    y += 1.22
foot(s, "A benchmark lives or dies on adoption — that judgement about the field is yours more than mine.")

# S14 — five factors technical
s = slide()
title(s, "The five factors in detail — and why three may be unresolvable",
      "For discussion: is this the right diagnostic strategy?")

rows = [["#", "What differs", "Mechanism", "Status"],
        ["1", "Topology", "2-node transition graph → 8-slot chain",
         ("Untested", MUTED, False)],
        ["2", "Positional embeddings", "8 learned vectors added to slot embeddings",
         ("Arm running", AMBER, True)],
        ["3", "Readout", "Stoichiometry-weighted sum → learned attention pooling",
         ("Untested", MUTED, False)],
        ["4", "Edge features", "17-dim port-pair + transition weight → none",
         ("Untested", MUTED, False)],
        ["5", "Replicas", "1 → 16, averaged inside the loss",
         ("Ruled out", TERRA, True)]]
table(s, M, 1.95, 7.55, rows, [0.42, 2.05, 3.6, 1.48], fs=10.5, hfs=10, rowh=0.5)

card(s, M, 5.05, 7.55, 1.55, TINT2)
tf = tb(s, M + 0.3, 5.26, 6.95, 1.2)
para(tf, "The power problem", size=13, color=INK, bold=True, font=HEAD,
     first=True, space=6)
para(tf, "The octamer-vs-HPG-hier gap is +0.053 (EA) and +0.068 (IP). Per-cell noise "
         "after three-seed averaging is roughly 0.03–0.05. Splitting that gap across "
         "four factors means resolving ~0.015 against ~0.04 of noise. Factors 1 and 3 "
         "are almost certainly unresolvable here.", size=11, color=BODY, space=0)

card(s, 8.4, 1.95, 4.38, 4.65, TINT)
tf = tb(s, 8.7, 2.16, 3.8, 4.2)
para(tf, "Factor 4 is the interesting one", size=13.5, color=INK, bold=True,
     font=HEAD, first=True, space=8)
para(tf, "The octamer's message passing takes no edge features at all. It "
         "discards the 17-dimensional port-pair features and the transition "
         "weight that the simpler model uses.", size=11.5, color=BODY, space=8)
para(tf, "So it wins while seeing less junction chemistry.", size=11.5,
     color=TEAL, bold=True, space=8)
para(tf, "Either that information was not helping, or the explicit sequence more "
         "than compensates. Neither the pooling nor the topology arm tests this — "
         "which is my argument against running them.", size=11.5, color=BODY, space=8)
para(tf, "Is that the right call, or would you run the 2×2 anyway?", size=11.5,
     color=TERRA, bold=True, space=0)
foot(s, "One measured limit: with depth 2, a slot exchanges messages with at most ±2 positions before pooling — so the model reads local blockiness, not global arrangement.")

# S15 — posemb pilot
s = slide()
title(s, "Positional embeddings: pilot points one way, arm is running",
      "Pre-registered before submission; threshold fixed in advance at 0.051")

table(s, M, 2.0, 7.55,
      [["Cell", "With position embeddings", "Without", "Change"],
       ["EA, fold 0, seed 42", "0.779", "0.606", ("−0.173", TERRA, True)],
       ["EA, fold 4, seed 42", "0.426", "0.357", ("−0.069", TERRA, True)]],
      [2.35, 2.15, 1.35, 1.7], fs=11.5, hfs=10.5, rowh=0.48)

card(s, M, 3.75, 7.55, 1.35, TINT2)
tf = tb(s, M + 0.3, 3.95, 6.95, 1.0)
para(tf, "Both cells exceed the pre-registered threshold, in the same direction. "
         "That is the first pre-registered outcome: position embeddings are a "
         "principal source of the advantage.", size=12, color=BODY, first=True, space=0)

card(s, M, 5.3, 7.55, 1.3, TINT)
tf = tb(s, M + 0.3, 5.5, 6.95, 1.0)
para(tf, "Worth noting: on fold 0 overall R² went up (+0.028) while architecture "
         "recovery fell 0.173. Same predictions, two metrics in opposite directions — "
         "our own measurement argument appearing inside our own ablation.",
     size=11.5, color=BODY, first=True, space=0)

card(s, 8.4, 2.0, 4.38, 4.6, TINT2)
tf = tb(s, 8.7, 2.22, 3.8, 4.1)
para(tf, "If this holds at 54 cells", size=13.5, color=INK, bold=True,
     font=HEAD, first=True, space=8)
para(tf, "We have an answer to \"why does the octamer win\" — and it is the learned "
         "position information, not the chain length.", size=11.5, color=BODY, space=8)
para(tf, "That also weakens the protocol-matching worry considerably, since "
         "positional embeddings have nothing to do with how the labels were averaged.",
     size=11.5, color=TEAL, space=10)
para(tf, "Caveats I would keep", size=12.5, color=INK, bold=True, font=HEAD, space=7)
for t in ["Two cells, one seed, one target — not a result",
          "R1 per-cell noise is right-tailed: median 0.051, mean 0.123",
          "Judge on a paired sign test across cells, not the median"]:
    para(tf, "• " + t, size=11, color=BODY, space=5)
foot(s, "[ADD: full 54-cell result, paired sign test count, and whether the B-split arm was gated in.]", TERRA)

# S16 — testing the confound
s = slide()
title(s, "Testing whether the octamer is dataset-specific — four options",
      "Ranked by cost. My inclination is 1 and 3 first.")

opts = [("1", "Test on a different dataset", TEAL,
         "We already hold a glass-transition dataset and a block copolymer phase "
         "dataset. Neither has labels built from 8-unit averaging. If the octamer "
         "still wins there, the sceptical reading weakens sharply.",
         "Days of GPU · no new data · data confirmed present in the repository"),
        ("2", "Sweep the model's chain length", AMBER,
         "8 vs 12 vs 16 vs 24 slots. A sharp peak at exactly 8 would indicate "
         "protocol-matching; flat or improving would not.",
         "One caution already checked: compositions are quarters, so 6 and 10 slots "
         "cannot represent them exactly and would look worse for purely arithmetic "
         "reasons. 8, 12, 16 and 24 are the clean comparisons."),
        ("3", "Probe what the representation contains", TEAL,
         "Freeze the model and test whether architecture can be read out of its "
         "internal representation by a linear classifier.",
         "Separates \"predicts better\" from \"represents better\" — which is the "
         "actual claim. One afternoon."),
        ("4", "Generate labels differently", TERRA,
         "The dataset proposal in the first half of this deck.",
         "Definitive, but months and ~9 kSU, and it needs your sign-off.")]
y = 1.88
for n, h, col, b, note in opts:
    card(s, M, y, 12.23, 1.05, TINT if n in "13" else TINT2)
    numdot(s, M + 0.3, y + 0.28, 0.4, n, col)
    tf = tb(s, M + 0.9, y + 0.14, 4.85, 0.88)
    para(tf, h, size=13, color=INK, bold=True, font=HEAD, first=True, space=4)
    para(tf, b, size=10.5, color=BODY, space=0)
    tf = tb(s, 6.75, y + 0.22, 5.95, 0.8)
    para(tf, note, size=10.5, color=MUTED, first=True, space=0)
    y += 1.16
foot(s, "Options 1 and 3 together would move the claim a long way for very little compute. Do you agree with that ordering?")

# S17 — questions
s = slide(dark=True)
tf = tb(s, M, 0.7, W - 2 * M, 0.7)
para(tf, "Questions for discussion", size=30, color=WHITE, bold=True,
     font=HEAD, first=True, space=0)

left = [("Dataset", ["Is this the right dataset to generate, or should we ship the "
                     "current paper with the limitation stated?",
                     "Is ~2,000 polymers with existing monomers the right scope?",
                     "Am I missing a dimension — chain length, dispersity, tacticity?",
                     "This paper or the next one?"])]
right = [("Octamer diagnosis", ["Do you agree with eliminating factors one at a time, "
                                "given three of them are probably unresolvable at this "
                                "noise level?",
                                "Should I run the external-validity tests (options 1 and 3) "
                                "before any more ablations?",
                                "Are there hypotheses I have not considered?",
                                "How should I frame the protocol-matching confound in the paper?"])]
for col_items, x in [(left, M), (right, 6.95)]:
    for h, items in col_items:
        card(s, x, 1.75, 6.05 if x == M else 5.83, 3.15, C(0x3A, 0x4A, 0x53))
        tf = tb(s, x + 0.35, 2.0, 5.2, 0.35)
        para(tf, h.upper(), size=11, color=TEAL, bold=True, first=True, space=12)
        tf = tb(s, x + 0.35, 2.42, (5.35 if x == M else 5.15), 2.4)
        for i, t in enumerate(items):
            para(tf, f"{i+1}.  {t}", size=12.5, color=WHITE, first=(i == 0), space=13)

tf = tb(s, M, 5.35, 11.8, 0.4)
para(tf, "What I would like to leave with: a decision on the dataset, and agreement "
         "on whether external validity comes before further ablation.",
     size=12, color=SAGE, first=True, space=0, italic=True)

prs.save("update_deck.pptx")
print("saved update_deck.pptx —", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
