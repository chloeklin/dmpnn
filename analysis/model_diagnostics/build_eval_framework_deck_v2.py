#!/usr/bin/env python3
"""Evaluation-framework deck, rebuild — Paper 1.

Every figure is a verified artefact from analysis/paper1_figures/ (F1-F6), each
with a manifest recording its prediction sources, selection criteria and cell
counts. Every number below traces to a manifest, a CSV in that directory, or an
independently reproduced calculation recorded in the 11 August verification.

Corrections carried in this rebuild:
  * the null beats a trained model on IP fold 0, not on an EA fold
  * demonstration numbers are reported per target, not pooled
  * F2 is a two-model paired panel, labelled a selected example with its base rate
  * architecture-spread recovery reported with its two-sided caveat

Run from anywhere:  python build_eval_framework_deck_v2.py
"""
import os
from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "..", "paper1_figures")
OUT = os.path.join(HERE, "Eval_framework_paper1_v2.pptx")

DARK  = C(0x2B, 0x3A, 0x42)
INK   = C(0x1A, 0x24, 0x2A)
BODY  = C(0x3D, 0x4C, 0x54)
MUTED = C(0x77, 0x86, 0x8E)
WHITE = C(0xFF, 0xFF, 0xFF)
TEAL  = C(0x02, 0x80, 0x90)
AMBER = C(0xC4, 0x8B, 0x0E)
TERRA = C(0xB8, 0x50, 0x42)
SAGE  = C(0xA7, 0xBE, 0xAE)
TINT  = C(0xF2, 0xF5, 0xF6)
TINT2 = C(0xE8, 0xEF, 0xF1)

HEAD, BODYF, MONO = "Cambria", "Calibri", "Courier New"
W, H, M = 13.333, 7.5, 0.55

prs = Presentation()
prs.slide_width, prs.slide_height = In(W), In(H)
BLANK = prs.slide_layouts[6]


def slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    if dark:
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, In(W), In(H))
        bg.fill.solid(); bg.fill.fore_color.rgb = DARK
        bg.line.fill.background(); bg.shadow.inherit = False
    return s


def tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf


def para(tf, text, size=13, color=BODY, bold=False, font=BODYF, space=6,
         first=False, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = text; p.space_after = Pt(space)
    for r in p.runs:
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = font; r.font.italic = italic
    return p


def card(s, x, y, w, h, fill=TINT):
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, In(x), In(y), In(w), In(h))
    r.fill.solid(); r.fill.fore_color.rgb = fill
    r.line.fill.background(); r.shadow.inherit = False
    r.adjustments[0] = 0.05
    return r


def numdot(s, x, y, d, n, fill):
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, In(x), In(y), In(d), In(d))
    o.fill.solid(); o.fill.fore_color.rgb = fill
    o.line.fill.background(); o.shadow.inherit = False
    tf = o.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = str(n); p.alignment = PP_ALIGN.CENTER
    for r in p.runs:
        r.font.size = Pt(14); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = BODYF
    return o


def piece_title(s, n, text, sub):
    numdot(s, M, 0.44, 0.5, n, TEAL)
    tf = tb(s, M + 0.7, 0.42, W - 2 * M - 0.7, 0.6)
    para(tf, text, size=24, color=INK, bold=True, font=HEAD, space=2, first=True)
    tf2 = tb(s, M + 0.7, 1.14, W - 2 * M - 0.7, 0.4)
    para(tf2, sub, size=12.5, color=MUTED, first=True, space=0)


def fig(s, name, x, y, w):
    p = os.path.join(FIG, name)
    if os.path.exists(p):
        return s.shapes.add_picture(p, In(x), In(y), width=In(w))
    ph = card(s, x, y, w, 2.4, TINT2)
    tf = tb(s, x + 0.2, y + 1.0, w - 0.4, 0.4)
    para(tf, f"[FIGURE MISSING: {name}]", size=12, color=TERRA, first=True, space=0)
    return ph


def caption(s, x, y, w, text):
    tf = tb(s, x, y, w, 0.4)
    para(tf, text, size=9, color=MUTED, first=True, space=0, italic=True)


def foot(s, text, color=MUTED):
    tf = tb(s, M, H - 0.46, W - 2 * M, 0.3)
    para(tf, text, size=9, color=color, first=True, space=0, italic=True)


def statcard(s, x, y, w, h, label, big, sub, col=TEAL, fill=TINT):
    card(s, x, y, w, h, fill)
    tf = tb(s, x + 0.26, y + 0.18, w - 0.5, 0.3)
    para(tf, label, size=9.5, color=col, bold=True, first=True, space=6)
    tf2 = tb(s, x + 0.26, y + 0.52, w - 0.5, 0.5)
    para(tf2, big, size=23, color=INK, bold=True, font=HEAD, first=True, space=4)
    tf3 = tb(s, x + 0.26, y + 1.06, w - 0.5, h - 1.18)
    para(tf3, sub, size=10.5, color=BODY, first=True, space=0)


B  = dict(size=11.5, color=BODY, space=6)
BB = dict(size=11.5, color=INK, space=6, bold=True)

# ══ S1 title ═════════════════════════════════════════════════════════════════
s = slide(dark=True)
tf = tb(s, M, 2.0, W - 2 * M, 0.4)
para(tf, "PAPER 1 — EVALUATION FRAMEWORK", size=12, color=TEAL, bold=True,
     first=True, space=0)
tf = tb(s, M, 2.55, 11.2, 2.0)
para(tf, "Current evaluation cannot see\nthe property the field claims to model",
     size=33, color=WHITE, bold=True, font=HEAD, first=True, space=0)
tf = tb(s, M, 4.9, 10.6, 1.0)
para(tf, "Six components of an evaluation framework for polymer representation "
         "learning. All six are built, measured, and demonstrated on data already held.",
     size=14, color=SAGE, first=True, space=0)
tf = tb(s, M, 6.3, 11.5, 0.6)
para(tf, "Every figure carries a manifest recording its sources and selection rule. "
         "No new compute is required to write this paper.",
     size=11.5, color=MUTED, first=True, space=0, italic=True)
s.notes_slide.notes_text_frame.text = (
    "Framing: this is not 'here are better metrics'. It is 'the metric everyone "
    "reports cannot answer the question everyone is asking', and here is the "
    "evidence. Rebuilt 11 August against six purpose-built, verified figures.")

# ══ S2 the six components ════════════════════════════════════════════════════
s = slide()
tf = tb(s, M, 0.42, W - 2 * M, 0.7)
para(tf, "The problem, and the six things needed to fix it", size=26, color=INK,
     bold=True, font=HEAD, first=True, space=0)
tf = tb(s, M, 1.18, W - 2 * M, 0.4)
para(tf, "Chain architecture is what we claim to model. It is ~1% of the variance "
         "in the target.", size=13, color=MUTED, first=True, space=0)

items = [("Variance decomposition", "Establishes that the signal is invisible to "
          "the standard metric", "F1"),
         ("Two-axis metric", "Separates chemistry placement from architecture "
          "recovery", "F2"),
         ("Null floor", "The best achievable without the thing being tested", "F3"),
         ("Split design", "Two chemistry-extrapolation regimes, folds not "
          "exchangeable", "F4"),
         ("Noise floor", "What difference is measurable at all", "F5"),
         ("The demonstration", "Two metrics, same predictions, opposite "
          "conclusions", "F6")]
x, y = M, 1.85
for i, (h, b, f) in enumerate(items, 1):
    card(s, x, y, 3.94, 1.55, TINT if i % 2 else TINT2)
    numdot(s, x + 0.28, y + 0.26, 0.42, i, TEAL if i < 6 else TERRA)
    tf = tb(s, x + 0.28, y + 0.82, 3.4, 0.65)
    para(tf, h, size=13.5, color=INK, bold=True, font=HEAD, first=True, space=4)
    para(tf, b, size=10.5, color=BODY, space=0)
    tf2 = tb(s, x + 3.35, y + 0.28, 0.5, 0.3)
    para(tf2, f, size=10, color=MUTED, bold=True, first=True, space=0)
    x += 4.14
    if i == 3:
        x, y = M, y + 1.75

card(s, M, 5.42, 12.23, 1.15, DARK)
tf = tb(s, M + 0.35, 5.62, 11.6, 0.8)
para(tf, "Together they answer a question the field currently cannot: is a "
         "representation better, or does it merely score higher?",
     size=14.5, color=WHITE, bold=True, font=HEAD, first=True, space=0)

# ══ S3 component 1 — variance decomposition ══════════════════════════════════
s = slide()
piece_title(s, 1, "Variance decomposition — the signal is 1% of the target",
            "Why the standard metric cannot see architecture")
fig(s, "f1_variance_decomposition.png", M, 1.72, 7.7)
caption(s, M, 5.55, 7.7,
        "F1 · data/ea_ip.csv, 42,966 rows · sequential R² on group-transform means · no model predictions")

statcard(s, 8.5, 1.78, 4.28, 1.6, "ARCHITECTURE, SHARE OF TOTAL VARIANCE",
         "0.98%  /  1.46%", "EA / IP. Monomer identity alone accounts for "
         "92.9% (EA) and 89.7% (IP).", TERRA)
statcard(s, 8.5, 3.52, 4.28, 1.6, "ARCHITECTURE, SHARE OF POST-AB RESIDUAL",
         "13.9%  /  14.1%", "Once monomer identity is removed, architecture is "
         "a seventh of what remains.", TEAL)

card(s, 8.5, 5.26, 4.28, 1.55, TINT2)
tf = tb(s, 8.76, 5.44, 3.8, 1.2)
para(tf, "Why this matters", size=10, color=TEAL, bold=True, first=True, space=7)
para(tf, "A metric dominated by a 90% term cannot resolve a 1% term. Overall R² "
         "is not a weak measure of architecture — it is not a measure of it.",
     size=11, color=BODY, space=0)
foot(s, "Component shares sum to 1.000 on both targets; residual is 0.000 — the four factors are exhaustive.")

# ══ S4 component 2 — two-axis metric ═════════════════════════════════════════
s = slide()
piece_title(s, 2, "Two-axis metric — separate placement from recovery",
            "One group of three polymers: same chemistry, same composition, different architecture")
# F2 has a taller aspect ratio than the other five; width tuned so the rendered
# height clears the caption line.
fig(s, "f2_worked_example.png", M, 1.72, 7.6)
caption(s, M, 5.70, 7.6,
        "F2 · selected example · R1 fold 2, EA · three seeds averaged at prediction level · "
        "1 of 61 eligible group-folds (1.0%)")

card(s, 8.7, 1.78, 4.08, 2.35, TINT)
tf = tb(s, 8.96, 1.96, 3.6, 2.0)
para(tf, "THE TWO AXES", size=10, color=TEAL, bold=True, first=True, space=8)
para(tf, "Group-mean R² — did the model put the chemistry in the right place?", **B)
para(tf, "ΔR² — within a group, did it recover the architecture ordering?", **B)
para(tf, "The second is computed on the residual after the first is removed.",
     size=11, color=MUTED, space=0, italic=True)

card(s, 8.7, 4.27, 4.08, 1.5, TINT2)
tf = tb(s, 8.96, 4.45, 3.6, 1.15)
para(tf, "ON THIS GROUP", size=10, color=TERRA, bold=True, first=True, space=7)
para(tf, "Both models place the chemistry within 0.05 eV. The octamer recovers "
         "99% of the true architecture range; wDMPNN recovers 7%.",
     size=11, color=BODY, space=0)

card(s, 8.7, 5.91, 4.08, 0.95, DARK)
tf = tb(s, 8.96, 6.08, 3.6, 0.7)
para(tf, "wDMPNN does not rank the architectures wrongly — it predicts nearly "
         "one value for all three.", size=11, color=WHITE, bold=True, first=True, space=0)
s.notes_slide.notes_text_frame.text = (
    "Say 'selected example' out loud. Base rate is in the manifest: 61 of 6,138 "
    "group-folds on EA (1.0%), 54 on IP (0.9%), present in 8 of 9 folds on each "
    "target. Median architecture-spread recovery across all groups: octamer 0.768 "
    "EA / 0.947 IP, wDMPNN 0.610 / 0.551. That ratio is TWO-SIDED — closer to 1 is "
    "better, not larger — so do not quote it as a ranking without the caveat.")

# ══ S5 component 3 — null floor ══════════════════════════════════════════════
s = slide()
piece_title(s, 3, "Null floor — the best achievable without the thing being tested",
            "A parameter-free group-mean lookup table that never sees the held-out monomer")
fig(s, "f3_null_floor.png", M, 1.72, 8.0)
caption(s, M, 5.5, 8.0,
        "F3 · A split, 9 folds × 2 targets × 3 series · null reused from "
        "aggregate_lomo_seeds.null_floor · no missing cells")

statcard(s, 8.8, 1.78, 3.98, 1.62, "NULL MEDIAN GROUP-MEAN R²",
         "0.676   vs   −0.034", "EA / IP. Same benchmark, same split design, "
         "opposite conclusions about difficulty.", TEAL)

card(s, 8.8, 3.54, 3.98, 1.72, TINT2)
tf = tb(s, 9.06, 3.72, 3.5, 1.4)
para(tf, "THE HEADLINE CASE", size=10, color=TERRA, bold=True, first=True, space=7)
para(tf, "On IP fold 0 the null scores 0.969 — beating both trained models.", **BB)
para(tf, "On EA folds 2 and 3 it reaches 0.961 and 0.953 while knowing nothing "
         "about the held-out monomer.", size=11, color=BODY, space=0)

card(s, 8.8, 5.4, 3.98, 1.46, DARK)
tf = tb(s, 9.06, 5.58, 3.5, 1.1)
para(tf, "A reported group-mean R² of 0.95 on those folds is a number a lookup "
         "table matches. Without a null floor you cannot tell.",
     size=11.5, color=WHITE, bold=True, first=True, space=0)
foot(s, "Values below the axis limit are clipped — EA fold 6 null = −19.07; IP folds 2, 3, 5 = −1.02, −3.21, −7.53.", TERRA)

# ══ S6 component 4 — split design ════════════════════════════════════════════
s = slide()
piece_title(s, 4, "Split design — two extrapolation regimes, folds not exchangeable",
            "The benchmark's B-monomer space is dominated by two scaffold families")
fig(s, "f4_split_design.png", M, 1.72, 8.0)
caption(s, M, 5.4, 8.0,
        "F4 · Murcko scaffolds over 682 B monomers · split monomer_b_heldout_clustered")

statcard(s, 8.8, 1.78, 3.98, 1.6, "TOP TWO SCAFFOLD FAMILIES",
         "62.5%", "317 and 109 of 682 B monomers. A balanced scaffold-disjoint "
         "split is therefore impossible.", TERRA)

card(s, 8.8, 3.52, 3.98, 1.9, TINT2)
tf = tb(s, 9.06, 3.7, 3.5, 1.6)
para(tf, "FOLDS ARE NOT EXCHANGEABLE", size=10, color=TEAL, bold=True,
     first=True, space=7)
para(tf, "S folds (0–3): every held-out B monomer has a same-scaffold relative "
         "in training.", **B)
para(tf, "D folds (5, 7, 8): none do.", **BB)
para(tf, "Pooling them averages two different experiments.", size=11,
     color=MUTED, space=0, italic=True)

card(s, 8.8, 5.56, 3.98, 1.3, DARK)
tf = tb(s, 9.06, 5.74, 3.5, 0.95)
para(tf, "So S and D folds are reported separately throughout, never pooled.",
     size=11.5, color=WHITE, bold=True, first=True, space=0)

# ══ S7 component 5 — noise floor ═════════════════════════════════════════════
s = slide()
piece_title(s, 5, "Noise floor — what difference is measurable at all",
            "Identical configuration, identical seed, three repeats")
fig(s, "f5_noise_floor.png", M, 1.72, 8.0)
caption(s, M, 5.45, 8.0,
        "F5 · left: 2 folds × 3 repeats, HPG-octamer A-split EA · right: across-seed "
        "SD of ΔR², all 9 folds, both models")

statcard(s, 8.8, 1.78, 3.98, 1.62, "GROUP-MEAN R², EA FOLD 1, THREE REPEATS",
         "0.450 / 0.790 / 0.978", "SD 0.268. Same configuration, same seed.", TERRA)

card(s, 8.8, 3.54, 3.98, 1.72, TINT2)
tf = tb(s, 9.06, 3.72, 3.5, 1.4)
para(tf, "WHAT THIS BOUNDS", size=10, color=TEAL, bold=True, first=True, space=7)
para(tf, "Any single-run comparison smaller than ~0.27 on this metric is "
         "unmeasurable.", **BB)
para(tf, "MAE is far more stable — SD 0.091 eV on the same cell, 0.018 eV on "
         "fold 0. Stability is metric-specific.", size=11, color=BODY, space=0)

card(s, 8.8, 5.4, 3.98, 1.46, DARK)
tf = tb(s, 9.06, 5.58, 3.5, 1.1)
para(tf, "Nobody in this literature publishes a noise floor. Most reported "
         "improvements are smaller than this one.",
     size=11.5, color=WHITE, bold=True, first=True, space=0)
s.notes_slide.notes_text_frame.text = (
    "If asked why the variance is so large on one fold and not the other: we do "
    "not fully know. It is a property of that fold's group structure, not of the "
    "model. That is itself worth reporting — and it is why every comparison in "
    "this work uses three seeds averaged at the prediction level.")

# ══ S8 component 6 — the demonstration ═══════════════════════════════════════
s = slide()
piece_title(s, 6, "The demonstration — same predictions, opposite conclusions",
            "B split, cross-scaffold D folds: octamer against wDMPNN at its published configuration")
fig(s, "f6_demonstration.png", M, 1.72, 8.3)
caption(s, M, 5.35, 8.3,
        "F6 · paired per-fold differences, D folds [4,5,6,7,8] · three seeds averaged "
        "at prediction level")

card(s, 9.1, 1.78, 3.68, 2.35, TINT)
tf = tb(s, 9.36, 1.96, 3.2, 2.0)
para(tf, "FOUR METRICS SAY 'TIE'", size=10, color=MUTED, bold=True, first=True, space=8)
para(tf, "Overall R², MAE, RMSE and group-mean R² each split 3 of 5 folds — "
         "exactly what a coin does.", **B)
para(tf, "Medians: EA +0.003, IP +0.010 on overall R².", size=11, color=MUTED,
     space=0, italic=True)

card(s, 9.1, 4.27, 3.68, 1.75, TINT2)
tf = tb(s, 9.36, 4.45, 3.2, 1.45)
para(tf, "ΔR² SAYS OTHERWISE", size=10, color=TERRA, bold=True, first=True, space=7)
para(tf, "5 of 5 folds, both targets.", **BB)
para(tf, "Median +0.2597 (EA) and +0.1518 (IP) — two orders of magnitude larger "
         "than the accuracy differences.", size=11, color=BODY, space=0)

card(s, 9.1, 6.05, 3.68, 0.86, DARK)
tf = tb(s, 9.36, 6.21, 3.2, 0.62)
para(tf, "One set of predictions. Two defensible metrics. Opposite answers.",
     size=11.5, color=WHITE, bold=True, first=True, space=0)
foot(s, "5 folds gives a minimum attainable two-sided sign-test p of 0.0625 — 5/5 is the strongest result the design permits, and is reported as such.")

# ══ S9 what the framework establishes ════════════════════════════════════════
s = slide()
tf = tb(s, M, 0.42, W - 2 * M, 0.7)
para(tf, "What the framework establishes", size=26, color=INK, bold=True,
     font=HEAD, first=True, space=0)
tf = tb(s, M, 1.18, W - 2 * M, 0.4)
para(tf, "Four claims, each tied to one figure", size=13, color=MUTED,
     first=True, space=0)

claims = [
    ("The standard metric cannot see the property", "F1",
     "Architecture is 0.98% (EA) and 1.46% (IP) of total variance. Overall R² "
     "is dominated by a term two orders of magnitude larger."),
    ("A split can be degenerate for one target and not another", "F3",
     "A-blind null median group-mean R² is 0.676 on EA against −0.034 on IP. "
     "On IP fold 0 the null beats every trained model."),
    ("Benchmark numbers at this scale are not measurable single-run", "F5",
     "Three identical repeats give group-mean R² of 0.450, 0.790 and 0.978 — "
     "SD 0.268, larger than any improvement in the literature."),
    ("Two defensible metrics disagree on the same predictions", "F6",
     "Four accuracy metrics split 3/5 folds. ΔR² splits 5/5 on both targets, "
     "with medians of +0.26 and +0.15."),
]
y = 1.85
for h, f, b in claims:
    card(s, M, y, 12.23, 1.16, TINT2 if f in ("F3", "F6") else TINT)
    tf = tb(s, M + 0.34, y + 0.2, 5.0, 0.8)
    para(tf, h, size=13, color=INK, bold=True, font=HEAD, first=True, space=4)
    para(tf, f, size=10, color=TEAL, bold=True, space=0)
    tf2 = tb(s, 6.0, y + 0.28, 6.6, 0.8)
    para(tf2, b, size=11, color=BODY, first=True, space=0)
    y += 1.26

card(s, M, 6.5, 12.23, 0.72, DARK)
tf = tb(s, M + 0.34, 6.66, 11.55, 0.5)
para(tf, "None of these depends on which model wins.", size=13.5, color=WHITE,
     bold=True, font=HEAD, first=True, space=0)

# ══ S10 paper strategy ═══════════════════════════════════════════════════════
s = slide()
tf = tb(s, M, 0.42, W - 2 * M, 0.7)
para(tf, "Where this sits — four papers, and why this one goes first", size=26,
     color=INK, bold=True, font=HEAD, first=True, space=0)
tf = tb(s, M, 1.18, W - 2 * M, 0.4)
para(tf, "The through-line comes from the accepted review's own Future "
         "Perspectives section", size=13, color=MUTED, first=True, space=0)

papers = [
    ("0", "Review — accepted, Digital Discovery", TEAL, "DONE",
     "Argues in the literature that representation quality needs evaluating "
     "beyond predictive accuracy. Contains a section titled exactly that."),
    ("1", "Evaluation framework  —  this deck", TERRA, "WRITE NOW",
     "Demonstrates that argument empirically on the field's standard copolymer "
     "benchmark. All six components built and verified; no new compute needed."),
    ("2", "Architecture-aware benchmark dataset", AMBER, "COSTED",
     "Enables the architecture-held-out split that Paper 1 shows is missing. "
     "~8–9 kSU CPU, pipeline verified portable."),
    ("3", "The model — octamer versus baseline", AMBER, "EVIDENCE IN HAND",
     "Significant on two comparisons at p = 0.004. Held until external validity "
     "is tested, so it does not carry the protocol-matching confound alone."),
]
y = 1.85
for n, h, col, tag, b in papers:
    card(s, M, y, 12.23, 1.02, TINT if n in "13" else TINT2)
    numdot(s, M + 0.3, y + 0.26, 0.42, n, col)
    tf = tb(s, M + 0.92, y + 0.14, 5.1, 0.8)
    para(tf, h, size=13.5, color=INK, bold=True, font=HEAD, first=True, space=4)
    para(tf, tag, size=10, color=col, bold=True, space=0)
    tf = tb(s, 6.4, y + 0.22, 6.3, 0.75)
    para(tf, b, size=11, color=BODY, first=True, space=0)
    y += 1.12

card(s, M, 6.4, 12.23, 0.75, DARK)
tf = tb(s, M + 0.32, 6.58, 11.6, 0.5)
para(tf, "Year 1 argues it in the literature. Year 2 demonstrates it on the "
         "benchmark everyone uses. Year 3 builds the dataset the argument requires.",
     size=13.5, color=WHITE, bold=True, font=HEAD, first=True, space=0)

# ══ S11 why this order ═══════════════════════════════════════════════════════
s = slide()
tf = tb(s, M, 0.42, W - 2 * M, 0.7)
para(tf, "Why Paper 1 first — three reasons, one of them defensive", size=26,
     color=INK, bold=True, font=HEAD, first=True, space=0)
tf = tb(s, M, 1.18, W - 2 * M, 0.4)
para(tf, "The alternative was to lead with the model result", size=13,
     color=MUTED, first=True, space=0)

reasons = [
    ("It is writable today", TEAL,
     "Every component is built, measured and verified against a manifest. The "
     "model paper depends on external-validity tests not yet run; the dataset "
     "paper depends on compute not yet approved."),
    ("It does not depend on which model wins", TEAL,
     "The framework stands whether the octamer, HPG-hier or something later "
     "turns out best. That insulates it entirely from the protocol-matching "
     "confound, which is the largest threat to the model paper."),
    ("The measurement work invalidates the modelling work", TERRA,
     "The noise floor, the checkpoint bug, the null floors and the split design "
     "collectively void most of the earlier model comparisons. The stronger "
     "contribution is the one that did the invalidating."),
]
x = M
for h, col, b in reasons:
    card(s, x, 1.85, 3.94, 2.5, TINT)
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, In(x + 0.3), In(2.1), In(0.18), In(0.18))
    o.fill.solid(); o.fill.fore_color.rgb = col
    o.line.fill.background(); o.shadow.inherit = False
    tf = tb(s, x + 0.3, 2.48, 3.35, 1.7)
    para(tf, h, size=14, color=INK, bold=True, font=HEAD, first=True, space=8)
    para(tf, b, size=11, color=BODY, space=0)
    x += 4.14

card(s, M, 4.6, 12.23, 1.35, TINT2)
tf = tb(s, M + 0.32, 4.82, 11.6, 1.0)
para(tf, "What Paper 1 deliberately does not claim", size=13.5, color=TERRA,
     bold=True, font=HEAD, first=True, space=7)
para(tf, "That our model is the best. The octamer appears only as one of the "
         "models used to demonstrate that two metrics can disagree on the same "
         "predictions. Its own significance result belongs to Paper 3, after "
         "external validity has been tested.", size=12, color=BODY, space=0)

card(s, M, 6.2, 12.23, 0.9, DARK)
tf = tb(s, M + 0.32, 6.4, 11.6, 0.55)
para(tf, "Target: submit in the first quarter of Year 3. No experiments block it.",
     size=13.5, color=WHITE, bold=True, font=HEAD, first=True, space=0)

# ══ S12 provenance ═══════════════════════════════════════════════════════════
s = slide()
tf = tb(s, M, 0.42, W - 2 * M, 0.7)
para(tf, "Provenance — what stands behind each figure", size=26, color=INK,
     bold=True, font=HEAD, first=True, space=0)
tf = tb(s, M, 1.18, W - 2 * M, 0.4)
para(tf, "Every figure has a manifest in analysis/paper1_figures/ recording "
         "sources, selection rule and cell count", size=13, color=MUTED,
     first=True, space=0)

rows = [
    ("F1", "Variance decomposition", "data/ea_ip.csv, 42,966 rows. No model predictions. Shares sum to 1.000."),
    ("F2", "Worked example", "Selected example, criterion stated. 1 of 61 eligible group-folds; base rate reported."),
    ("F3", "Null floor", "9 folds × 2 targets × 3 series = 54 cells, none missing. Null reused, not rebuilt."),
    ("F4", "Split design", "682 B monomers, Murcko scaffolds. Fold composition verified against split metadata."),
    ("F5", "Noise floor", "2 folds × 3 repeats, plus across-seed SD over all 9 folds and both models."),
    ("F6", "Demonstration", "5 D folds × 2 targets. Every median and win count matches expected in f6_summary.csv."),
]
y = 1.8
for ri, (tag, name, detail) in enumerate(rows):
    card(s, M, y, 12.23, 0.76, TINT if ri % 2 == 0 else TINT2)
    tf = tb(s, M + 0.3, y + 0.2, 0.6, 0.4)
    para(tf, tag, size=12, color=TEAL, bold=True, font=HEAD, first=True, space=0)
    tf = tb(s, M + 1.0, y + 0.2, 3.0, 0.4)
    para(tf, name, size=11.5, color=INK, bold=True, first=True, space=0)
    tf = tb(s, 4.9, y + 0.2, 7.8, 0.5)
    para(tf, detail, size=10.5, color=BODY, first=True, space=0)
    y += 0.86

card(s, M, 6.75, 12.23, 0.6, TINT2)
tf = tb(s, M + 0.32, 6.88, 11.55, 0.4)
para(tf, "Standing protocol: 3 seeds (42/43/44) averaged at the prediction level, "
         "best-checkpoint predictions, paired per-fold sign tests, S and D folds never pooled.",
     size=10.5, color=BODY, first=True, space=0, italic=True)

prs.save(OUT)
print(f"saved {OUT}")
print(f"slides: {len(prs.slides.__iter__.__self__._sldIdLst)}")
