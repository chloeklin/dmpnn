#!/usr/bin/env python3
"""
Stage 2D Supervisor Presentation Generator
==========================================
Generates publication-quality figures, summary table, slide deck,
and one-paragraph summary from final corrected Stage 2D results.

Outputs (all saved to postrerun_output/):
  fig1_model_progression.{pdf,png}
  fig2_overall_performance.{pdf,png}
  fig3_architecture_deviation.{pdf,png}
  fig4_head_to_head.{pdf,png}
  fig5_stage2d_summary.{pdf,png}
  table1_stage2d_summary.{csv,md}
  stage2d_supervisor_summary.pptx
  stage2d_supervisor_summary.md
"""

import numpy as np
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
from pathlib import Path

# ── Paths ────────────────────────────────────────────────────────
OUT = Path(__file__).parent / "postrerun_output"
OUT.mkdir(exist_ok=True)

# ── Load data ────────────────────────────────────────────────────
master = pd.read_csv(OUT / "stage2d_master_results.csv")
archdev = pd.read_csv(OUT / "stage2d_architecture_deviation_results.csv")

# Build lookup dicts
def lookup(df, col):
    d = {}
    for _, row in df.iterrows():
        d[(row['variant'], row['target'])] = row[col]
    return d

r2 = lookup(master, 'R2')
mae = lookup(master, 'MAE')
r2d = lookup(archdev, 'R2_dev')

VARIANTS = ['frac', '2d0_fixed', '2d0_arch', '2d0_gate',
            '2d1_fixed', '2d1_arch', '2d1_gate']
LABELS = ['Frac', '2D0-fixed', '2D0-arch', '2D0-gate',
          '2D1-fixed', '2D1-arch', '2D1-gate']
TARGETS = ['EA', 'IP']
TARGET_FULL = {'EA': 'EA', 'IP': 'IP'}

# Best models from analysis
BEST_2D0 = '2d0_arch'
BEST_2D1_EA = '2d1_arch'
BEST_2D1_IP = '2d1_fixed'

# Significance data (from significance_analysis.md)
SIG = {
    'EA_overall':  {'delta': +0.0010, 'std': 0.0028, 'p': 0.5111, 'sig': False},
    'EA_archdev':  {'delta': +0.0192, 'std': 0.0112, 'p': 0.0266, 'sig': True},
    'IP_overall':  {'delta': +0.0014, 'std': 0.0014, 'p': 0.1157, 'sig': False},
    'IP_archdev':  {'delta': +0.0145, 'std': 0.0057, 'p': 0.0071, 'sig': True},
}

# ── Style constants ──────────────────────────────────────────────
plt.rcParams.update({
    'font.family': 'sans-serif',
    'font.size': 11,
    'axes.titlesize': 13,
    'axes.labelsize': 12,
    'xtick.labelsize': 10,
    'ytick.labelsize': 10,
    'legend.fontsize': 10,
    'figure.dpi': 150,
    'savefig.dpi': 300,
    'savefig.bbox': 'tight',
    'savefig.pad_inches': 0.15,
})

# Color palette
C_FRAC = '#8c8c8c'       # grey
C_2D0  = '#4C72B0'       # blue family
C_2D1  = '#DD8452'       # orange family
C_2D0_LIGHT = '#7DA1D4'
C_2D1_LIGHT = '#E8AD85'

VARIANT_COLORS = {
    'frac':      C_FRAC,
    '2d0_fixed': '#7DA1D4',
    '2d0_arch':  '#4C72B0',
    '2d0_gate':  '#2B5090',
    '2d1_fixed': '#E8AD85',
    '2d1_arch':  '#DD8452',
    '2d1_gate':  '#C46A30',
}

def sig_star(p):
    if p < 0.001: return '***'
    if p < 0.01:  return '**'
    if p < 0.05:  return '*'
    return 'n.s.'

def save_fig(fig, name):
    fig.savefig(OUT / f"{name}.pdf")
    fig.savefig(OUT / f"{name}.png")
    plt.close(fig)
    print(f"  → {name}.pdf/png")

# ====================================================================
# FIGURE 1: Model Progression Overview
# ====================================================================
def make_fig1():
    fig, ax = plt.subplots(figsize=(10, 5))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 6)
    ax.axis('off')
    fig.suptitle('From Composition to Architecture-Aware Modeling',
                 fontsize=15, fontweight='bold', y=0.97)

    box_props = dict(boxstyle='round,pad=0.6', linewidth=2)
    title_kw = dict(fontsize=12, fontweight='bold', ha='center', va='center')
    eq_kw = dict(fontsize=11, ha='center', va='center', family='monospace',
                 style='italic')
    desc_kw = dict(ha='center', va='center', color='#555555')

    # Stage 2A — Frac
    ax.add_patch(FancyBboxPatch((0.5, 3.6), 2.4, 1.8,
                                facecolor='#E8E8E8', edgecolor=C_FRAC, **box_props))
    ax.text(1.7, 5.0, 'Stage 2A', **title_kw)
    ax.text(1.7, 4.5, 'Frac', fontsize=11, ha='center', va='center',
            color=C_FRAC, fontweight='bold')
    ax.text(1.7, 4.0, r'$h_{poly} = f_A h_A + f_B h_B$', **eq_kw)

    # Arrow 1
    ax.annotate('', xy=(3.55, 4.5), xytext=(2.95, 4.5),
                arrowprops=dict(arrowstyle='->', lw=2.5, color='#333'))
    ax.text(3.25, 4.85, '+arch info', fontsize=8, ha='center', color='#333')

    # Stage 2D0 — Global Architecture
    ax.add_patch(FancyBboxPatch((3.8, 3.6), 2.4, 1.8,
                                facecolor='#DAE5F5', edgecolor=C_2D0, **box_props))
    ax.text(5.0, 5.0, 'Stage 2D0', **title_kw)
    ax.text(5.0, 4.5, 'Global Architecture', fontsize=10, ha='center',
            va='center', color=C_2D0, fontweight='bold')
    ax.text(5.0, 4.0, r'$h_{mix} + \alpha_{arch} \cdot e_{arch}$', **eq_kw)

    # Arrow 2
    ax.annotate('', xy=(6.85, 4.5), xytext=(6.25, 4.5),
                arrowprops=dict(arrowstyle='->', lw=2.5, color='#333'))
    ax.text(6.55, 4.85, '+chemistry', fontsize=8, ha='center', color='#333')

    # Stage 2D1 — Chemistry-Conditioned
    ax.add_patch(FancyBboxPatch((7.1, 3.6), 2.5, 1.8,
                                facecolor='#FDEBD0', edgecolor=C_2D1, **box_props))
    ax.text(8.35, 5.0, 'Stage 2D1', **title_kw)
    ax.text(8.35, 4.5, 'Chemistry-Conditioned', fontsize=10, ha='center',
            va='center', color=C_2D1, fontweight='bold')
    ax.text(8.35, 3.95, r'$h_{mix} + \alpha_{arch} \cdot r_{arch}(h_A, h_B, ...)$',
            **eq_kw)

    # Bottom descriptions
    ax.text(1.7, 2.8, 'Baseline:\nWeighted average\nof monomer reps',
            **desc_kw, fontsize=9)
    ax.text(5.0, 2.8, 'Architecture offset:\nLearned embedding\nper arch type',
            **desc_kw, fontsize=9)
    ax.text(8.35, 2.8, 'Architecture + chemistry:\nInteraction MLP conditions\non monomer identity',
            **desc_kw, fontsize=9)

    save_fig(fig, 'fig1_model_progression')


# ====================================================================
# FIGURE 2: Overall Performance
# ====================================================================
def make_fig2():
    fig, axes = plt.subplots(1, 2, figsize=(12, 5), sharey=False)
    fig.suptitle('Overall Prediction Performance', fontsize=14, fontweight='bold',
                 y=1.02)

    for ax_i, tgt in enumerate(TARGETS):
        ax = axes[ax_i]
        tfull = TARGET_FULL[tgt]
        vals = [r2[(v, tfull)] for v in VARIANTS]
        colors = [VARIANT_COLORS[v] for v in VARIANTS]

        bars = ax.bar(range(len(VARIANTS)), vals, color=colors, edgecolor='white',
                      linewidth=0.8, width=0.7)

        # Frac baseline line
        frac_val = r2[('frac', tfull)]
        ax.axhline(frac_val, color=C_FRAC, linestyle='--', linewidth=1,
                   alpha=0.7, zorder=0)
        ax.text(len(VARIANTS) - 0.3, frac_val + 0.0003, 'Frac baseline',
                fontsize=8, color=C_FRAC, ha='right', va='bottom')

        # Find best
        best_idx = np.argmax(vals)
        bars[best_idx].set_edgecolor('#222222')
        bars[best_idx].set_linewidth(2)

        # Annotate improvement over frac
        for i, v in enumerate(vals):
            delta = v - frac_val
            if delta > 0:
                ax.text(i, v + 0.0003, f'+{delta:.4f}', ha='center', va='bottom',
                        fontsize=7.5, fontweight='bold' if i == best_idx else 'normal')

        # Star the best
        ax.text(best_idx, vals[best_idx] + 0.0012, '★',
                ha='center', va='bottom', fontsize=14, color='#DAA520')

        ax.set_xticks(range(len(VARIANTS)))
        ax.set_xticklabels(LABELS, rotation=35, ha='right', fontsize=9)
        ax.set_ylabel('R²')
        ax.set_title(f'{tgt} vs SHE (eV)', fontweight='bold')
        ax.set_ylim(0.960, max(vals) + 0.003)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)

    fig.tight_layout()
    save_fig(fig, 'fig2_overall_performance')


# ====================================================================
# FIGURE 3: Architecture-Deviation Performance
# ====================================================================
def make_fig3():
    fig, axes = plt.subplots(1, 2, figsize=(12, 5), sharey=False)
    fig.suptitle('Architecture-Deviation Prediction Performance',
                 fontsize=14, fontweight='bold', y=1.02)

    for ax_i, tgt in enumerate(TARGETS):
        ax = axes[ax_i]
        tfull = TARGET_FULL[tgt]
        vals = [r2d[(v, tfull)] for v in VARIANTS]
        colors = [VARIANT_COLORS[v] for v in VARIANTS]

        bars = ax.bar(range(len(VARIANTS)), vals, color=colors, edgecolor='white',
                      linewidth=0.8, width=0.7)

        # Best 2D0 and 2D1
        d0_vals = {v: r2d[(v, tfull)] for v in VARIANTS if v.startswith('2d0')}
        d1_vals = {v: r2d[(v, tfull)] for v in VARIANTS if v.startswith('2d1')}
        best_d0 = max(d0_vals, key=d0_vals.get)
        best_d1 = max(d1_vals, key=d1_vals.get)
        best_d0_idx = VARIANTS.index(best_d0)
        best_d1_idx = VARIANTS.index(best_d1)

        bars[best_d0_idx].set_edgecolor(C_2D0)
        bars[best_d0_idx].set_linewidth(2.5)
        bars[best_d1_idx].set_edgecolor(C_2D1)
        bars[best_d1_idx].set_linewidth(2.5)

        # Delta annotation between best 2D0 and best 2D1
        delta = d1_vals[best_d1] - d0_vals[best_d0]
        mid_x = (best_d0_idx + best_d1_idx) / 2
        top_y = max(d0_vals[best_d0], d1_vals[best_d1])
        ax.annotate('', xy=(best_d1_idx, top_y + 0.008),
                    xytext=(best_d0_idx, top_y + 0.008),
                    arrowprops=dict(arrowstyle='<->', color='#333', lw=1.5))
        ax.text(mid_x, top_y + 0.012, f'Δ = +{delta:.4f}',
                ha='center', va='bottom', fontsize=9, fontweight='bold',
                color='#333')

        # Annotate Frac baseline
        frac_val = r2d[('frac', tfull)]
        ax.axhline(0, color='#ccc', linestyle='-', linewidth=0.5, zorder=0)

        ax.set_xticks(range(len(VARIANTS)))
        ax.set_xticklabels(LABELS, rotation=35, ha='right', fontsize=9)
        ax.set_ylabel('R²(Δy)')
        ax.set_title(f'{tgt} Architecture-Deviation R²', fontweight='bold')

        ymin = min(vals) - 0.05
        ymax = max(vals) + 0.04
        ax.set_ylim(ymin, ymax)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)

    fig.tight_layout()
    save_fig(fig, 'fig3_architecture_deviation')


# ====================================================================
# FIGURE 4: Head-to-Head Comparison
# ====================================================================
def make_fig4():
    fig, ax = plt.subplots(figsize=(10, 5.5))
    fig.suptitle('Best 2D0 vs Best 2D1: Head-to-Head',
                 fontsize=14, fontweight='bold', y=1.0)

    # 4 metric groups
    metrics = [
        ('Overall EA R²', 'EA_overall',
         r2[(BEST_2D0, TARGET_FULL['EA'])],
         r2[(BEST_2D1_EA, TARGET_FULL['EA'])]),
        ('Overall IP R²', 'IP_overall',
         r2[(BEST_2D0, TARGET_FULL['IP'])],
         r2[(BEST_2D1_IP, TARGET_FULL['IP'])]),
        ('Arch-dev EA R²', 'EA_archdev',
         r2d[(BEST_2D0, TARGET_FULL['EA'])],
         r2d[(BEST_2D1_EA, TARGET_FULL['EA'])]),
        ('Arch-dev IP R²', 'IP_archdev',
         r2d[(BEST_2D0, TARGET_FULL['IP'])],
         r2d[(BEST_2D1_IP, TARGET_FULL['IP'])]),
    ]

    x = np.arange(len(metrics))
    w = 0.32

    bars_d0 = ax.bar(x - w/2, [m[2] for m in metrics], w, label='Best 2D0 (2D0-arch)',
                      color=C_2D0, edgecolor='white', linewidth=0.8)
    bars_d1 = ax.bar(x + w/2, [m[3] for m in metrics], w, label='Best 2D1',
                      color=C_2D1, edgecolor='white', linewidth=0.8)

    # Annotate deltas + p-values
    for i, (label, sig_key, v0, v1) in enumerate(metrics):
        s = SIG[sig_key]
        delta = s['delta']
        p = s['p']
        star = sig_star(p)
        top = max(v0, v1)

        # Delta text
        color = '#C0392B' if s['sig'] else '#555555'
        weight = 'bold' if s['sig'] else 'normal'
        ax.text(i, top + 0.012, f'Δ = +{delta:.4f}\np = {p:.4f} {star}',
                ha='center', va='bottom', fontsize=8.5, color=color,
                fontweight=weight, linespacing=1.3)

    ax.set_xticks(x)
    ax.set_xticklabels([m[0] for m in metrics], fontsize=10)
    ax.set_ylabel('R²')
    ax.legend(loc='lower right', framealpha=0.9)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)

    # Adjust y limits: overall metrics ~0.96-0.99, archdev ~0.82-0.90
    # Use broken axis effect by just setting a reasonable range
    ax.set_ylim(0.78, 1.02)
    ax.axhline(0.96, color='#eee', linewidth=0.5, zorder=0)

    fig.tight_layout()
    save_fig(fig, 'fig4_head_to_head')


# ====================================================================
# FIGURE 5: The Main Takeaway (4-panel summary)
# ====================================================================
def make_fig5():
    fig, axes = plt.subplots(2, 2, figsize=(10, 8))
    fig.suptitle('Stage 2D Summary: Frac → 2D0 → 2D1',
                 fontsize=15, fontweight='bold', y=1.01)

    panels = [
        (axes[0, 0], 'Overall EA R²', 'EA', r2, TARGET_FULL['EA'],
         BEST_2D0, BEST_2D1_EA, 'EA_overall'),
        (axes[0, 1], 'Overall IP R²', 'IP', r2, TARGET_FULL['IP'],
         BEST_2D0, BEST_2D1_IP, 'IP_overall'),
        (axes[1, 0], 'Arch-dev EA R²', 'EA', r2d, TARGET_FULL['EA'],
         BEST_2D0, BEST_2D1_EA, 'EA_archdev'),
        (axes[1, 1], 'Arch-dev IP R²', 'IP', r2d, TARGET_FULL['IP'],
         BEST_2D0, BEST_2D1_IP, 'IP_archdev'),
    ]

    for ax, title, tgt, data, tfull, d0, d1, sig_key in panels:
        vals = [data[('frac', tfull)], data[(d0, tfull)], data[(d1, tfull)]]
        labels_short = ['Frac', f'2D0-arch', f'Best 2D1']
        colors = [C_FRAC, C_2D0, C_2D1]

        bars = ax.bar(range(3), vals, color=colors, edgecolor='white',
                      linewidth=0.8, width=0.55)

        # Value labels on bars
        for j, v in enumerate(vals):
            ax.text(j, v + 0.001, f'{v:.4f}', ha='center', va='bottom',
                    fontsize=9, fontweight='bold')

        # Arrows with delta annotations
        for j in range(2):
            delta = vals[j+1] - vals[j]
            mid_y = (vals[j] + vals[j+1]) / 2
            sign = '+' if delta >= 0 else ''
            s = SIG.get(sig_key, {})
            p = s.get('p', None)

            # Arrow at top
            y_arrow = max(vals[j], vals[j+1]) + 0.006
            ax.annotate('', xy=(j+1, y_arrow), xytext=(j, y_arrow),
                        arrowprops=dict(arrowstyle='->', lw=1.5,
                                        color='#333'))

            txt = f'{sign}{delta:.4f}'
            if j == 1 and p is not None:
                star = sig_star(p)
                color = '#C0392B' if s.get('sig', False) else '#555'
                txt += f'\n({star})'
            else:
                color = '#333'

            ax.text(j + 0.5, y_arrow + 0.001, txt,
                    ha='center', va='bottom', fontsize=8, color=color,
                    fontweight='bold')

        ax.set_xticks(range(3))
        ax.set_xticklabels(labels_short, fontsize=9)
        ax.set_title(title, fontweight='bold', fontsize=12)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)

        # Y limits
        ymin = min(vals) - 0.015
        ymax = max(vals) + 0.025
        ax.set_ylim(ymin, ymax)
        ax.set_ylabel('R²')

    # Panel labels
    for i, (ax, label) in enumerate(zip(axes.flat, ['A', 'B', 'C', 'D'])):
        ax.text(-0.12, 1.10, label, transform=ax.transAxes,
                fontsize=16, fontweight='bold', va='top')

    fig.tight_layout()
    save_fig(fig, 'fig5_stage2d_summary')


# ====================================================================
# TABLE 1: Supervisor Summary Table
# ====================================================================
def make_table1():
    rows = []
    for label, variant_ea, variant_ip in [
        ('Frac', 'frac', 'frac'),
        ('Best 2D0 (2D0-arch)', BEST_2D0, BEST_2D0),
        ('Best 2D1 (EA: 2D1-arch, IP: 2D1-fixed)', BEST_2D1_EA, BEST_2D1_IP),
    ]:
        rows.append({
            'Model': label,
            'EA_R2': f"{r2[(variant_ea, TARGET_FULL['EA'])]:.4f}",
            'IP_R2': f"{r2[(variant_ip, TARGET_FULL['IP'])]:.4f}",
            'EA_ArchDev_R2': f"{r2d[(variant_ea, TARGET_FULL['EA'])]:.4f}",
            'IP_ArchDev_R2': f"{r2d[(variant_ip, TARGET_FULL['IP'])]:.4f}",
        })

    df = pd.DataFrame(rows)
    df.to_csv(OUT / 'table1_stage2d_summary.csv', index=False)

    # Markdown
    md = '# Stage 2D Summary Table\n\n'
    md += '| Model | EA R² | IP R² | EA Arch-Dev R² | IP Arch-Dev R² |\n'
    md += '|-------|-------|-------|----------------|----------------|\n'
    for _, row in df.iterrows():
        md += f"| {row['Model']} | {row['EA_R2']} | {row['IP_R2']} | {row['EA_ArchDev_R2']} | {row['IP_ArchDev_R2']} |\n"
    md += '\n**Significance (2D1 vs 2D0):**\n\n'
    md += f"- EA overall R²: Δ = +{SIG['EA_overall']['delta']:.4f}, p = {SIG['EA_overall']['p']:.4f} ({sig_star(SIG['EA_overall']['p'])})\n"
    md += f"- IP overall R²: Δ = +{SIG['IP_overall']['delta']:.4f}, p = {SIG['IP_overall']['p']:.4f} ({sig_star(SIG['IP_overall']['p'])})\n"
    md += f"- EA arch-dev R²: Δ = +{SIG['EA_archdev']['delta']:.4f}, p = {SIG['EA_archdev']['p']:.4f} ({sig_star(SIG['EA_archdev']['p'])})\n"
    md += f"- IP arch-dev R²: Δ = +{SIG['IP_archdev']['delta']:.4f}, p = {SIG['IP_archdev']['p']:.4f} ({sig_star(SIG['IP_archdev']['p'])})\n"

    (OUT / 'table1_stage2d_summary.md').write_text(md)
    print(f"  → table1_stage2d_summary.csv/md")


# ====================================================================
# SLIDE DECK
# ====================================================================
def make_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_title_slide(title, subtitle=''):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5),
                                         Inches(11.7), Inches(2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
        p.alignment = PP_ALIGN.LEFT

        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(18)
            p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            p2.alignment = PP_ALIGN.LEFT
        return slide

    def add_image_slide(title, img_path, subtitle=''):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2),
                                         Inches(12), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(14)
            p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            p2.alignment = PP_ALIGN.LEFT

        # Image centered
        if Path(img_path).exists():
            slide.shapes.add_picture(str(img_path),
                                     Inches(0.8), Inches(1.3),
                                     Inches(11.7), Inches(5.8))
        return slide

    # Slide 1: Motivation
    s = add_title_slide(
        'Stage 2D: Architecture-Aware Copolymer Modeling',
        'Does copolymer architecture affect property predictions?\n'
        'If so, does chemistry-conditioned architecture modeling help?'
    )
    txBox = s.shapes.add_textbox(Inches(0.8), Inches(4.0),
                                 Inches(11.7), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    bullets = [
        'Stage 2A (Frac): h_poly = f_A·h_A + f_B·h_B  — architecture-blind',
        'Stage 2D0: h_poly = h_mix + α·e_arch  — global architecture offset',
        'Stage 2D1: h_poly = h_mix + α·r_arch(h_A, h_B, ...)  — chemistry-conditioned',
        '7 model variants tested across 5-fold a_held_out cross-validation',
    ]
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f'• {b}'
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Slide 2: Model progression
    add_image_slide('Model Progression: Frac → 2D0 → 2D1',
                    OUT / 'fig1_model_progression.png',
                    'Increasing architecture awareness')

    # Slide 3: Overall performance
    add_image_slide('Overall Prediction Performance',
                    OUT / 'fig2_overall_performance.png',
                    'All architecture-aware variants significantly improve over Frac (p < 0.05)')

    # Slide 4: Arch-deviation performance
    add_image_slide('Architecture-Deviation Prediction',
                    OUT / 'fig3_architecture_deviation.png',
                    '2D1 variants better capture architecture-specific property differences')

    # Slide 5: Head-to-head
    add_image_slide('Best 2D0 vs Best 2D1',
                    OUT / 'fig4_head_to_head.png',
                    'Overall R² gains are small (n.s.); architecture-deviation gains are significant')

    # Slide 6: Final recommendation
    s = add_image_slide('Summary & Recommendation',
                        OUT / 'fig5_stage2d_summary.png')
    txBox = s.shapes.add_textbox(Inches(0.5), Inches(6.5),
                                 Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = ('Recommendation: 2D1-arch (best overall R²; '
              'significantly better architecture-deviation R²)')
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xC0, 0x39, 0x2B)
    p.alignment = PP_ALIGN.CENTER

    pptx_path = OUT / 'stage2d_supervisor_summary.pptx'
    prs.save(str(pptx_path))
    print(f"  → stage2d_supervisor_summary.pptx")


# ====================================================================
# SUMMARY PARAGRAPH
# ====================================================================
def make_summary_md():
    text = """# Stage 2D: Supervisor Summary

Stage 2D investigated whether incorporating copolymer architecture information (alternating, random, block) improves property predictions for electron affinity (EA) and ionisation potential (IP). We compared three model tiers: **Frac** (composition-only baseline), **2D0** (global architecture offset via a learned embedding), and **2D1** (chemistry-conditioned architecture modeling via an interaction MLP that conditions on both monomer identities and architecture). Seven variants were evaluated across 5-fold architecture-held-out cross-validation.

**Architecture matters.** All six architecture-aware variants significantly outperform Frac (p < 0.05 on every variant), improving overall R² by +0.006–0.008 for EA and +0.013–0.016 for IP.

**2D1 provides a modest but meaningful improvement over 2D0.** Overall R² differences are small and not statistically significant (EA: Δ = +0.0010, p = 0.51; IP: Δ = +0.0014, p = 0.12). However, on the more targeted architecture-deviation metric — which measures how well a model captures property differences *between* architectures sharing the same monomers — 2D1 significantly outperforms 2D0 for both targets (EA: Δ = +0.019, p = 0.027; IP: Δ = +0.015, p = 0.007).

**Recommended model: 2D1-arch** (per-architecture alpha × interaction MLP). It achieves the highest mean R² (0.9805) and significantly better architecture-deviation predictions. The additional complexity is justified when architecture-specific prediction accuracy matters, which is the core scientific question of this stage.
"""
    (OUT / 'stage2d_supervisor_summary.md').write_text(text.strip() + '\n')
    print(f"  → stage2d_supervisor_summary.md")


# ====================================================================
# MAIN
# ====================================================================
if __name__ == '__main__':
    print("Generating Stage 2D Supervisor Presentation")
    print("=" * 60)

    print("\nFigure 1: Model Progression")
    make_fig1()

    print("\nFigure 2: Overall Performance")
    make_fig2()

    print("\nFigure 3: Architecture-Deviation Performance")
    make_fig3()

    print("\nFigure 4: Head-to-Head Comparison")
    make_fig4()

    print("\nFigure 5: Summary Figure")
    make_fig5()

    print("\nTable 1: Summary Table")
    make_table1()

    print("\nSlide Deck")
    make_pptx()

    print("\nSupervisor Summary")
    make_summary_md()

    print("\n" + "=" * 60)
    print(f"All outputs saved to: {OUT}")
    print("=" * 60)
